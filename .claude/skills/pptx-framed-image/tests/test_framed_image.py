"""Tests for the pptx-framed-image skill.

Run standalone:  python tests/test_framed_image.py
(also discoverable by pytest). No network, no template needed — a blank
Presentation is built in-memory.
"""
import os
import sys
import tempfile

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "scripts"))

from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn

import framed_image as FI
import nature_images as NI

FRAME = (6270019, 304800, 2593200, 3705000)  # a real OCTO frame's bounds (EMU)
ADJ1, ADJ2 = 50000, 0


def _tmp_img(w=180, h=257, kind="sunset"):
    p = os.path.join(tempfile.gettempdir(), f"_nat_{kind}_{w}x{h}.png")
    NI.generate_to(p, kind, w, h, seed=1)
    return p


def test_round2diag_geom_values():
    g = FI.round2diag_geom(ADJ1, ADJ2)
    assert g.get("prst") == "round2DiagRect"
    gds = g.findall(qn("a:avLst") + "/" + qn("a:gd"))
    vals = {gd.get("name"): gd.get("fmla") for gd in g.iter(qn("a:gd"))}
    assert vals == {"adj1": f"val {ADJ1}", "adj2": f"val {ADJ2}"}, vals
    print("ok  round2diag_geom builds correct preset + adjustments")


def test_picture_gets_frame_geometry_and_bounds():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    img = _tmp_img()
    pic = FI.place_image_in_frame(slide, img, *FRAME, geom=FI.round2diag_geom(ADJ1, ADJ2))

    # geometry cloned onto the picture (not a plain rect, not a custom path)
    spPr = pic._element.spPr
    prst = spPr.find(qn("a:prstGeom"))
    assert prst is not None and prst.get("prst") == "round2DiagRect"
    assert spPr.find(qn("a:custGeom")) is None
    # exactly one geometry child
    assert len(spPr.findall(qn("a:prstGeom"))) == 1

    # placed exactly at the frame bounds
    assert (pic.left, pic.top, pic.width, pic.height) == tuple(Emu(v) for v in FRAME)

    # geometry sits right after xfrm (schema order) so PowerPoint accepts it
    kids = [c.tag for c in spPr]
    assert kids.index(qn("a:xfrm")) < kids.index(qn("a:prstGeom"))
    print("ok  picture is clipped to the frame preset, at the frame bounds, schema-ordered")


def test_picture_within_slide_bounds():
    prs = Presentation()
    sw, sh = prs.slide_width, prs.slide_height
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    pic = FI.place_image_in_frame(slide, _tmp_img(), *FRAME, geom=FI.round2diag_geom())
    assert pic.left >= 0 and pic.top >= 0
    assert pic.left + pic.width <= sw, "picture runs off the right edge"
    assert pic.top + pic.height <= sh, "picture runs off the bottom edge"
    print("ok  picture stays inside the slide")


def test_none_geom_leaves_plain_rect():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    pic = FI.place_image_in_frame(slide, _tmp_img(), *FRAME, geom=None)
    g = pic._element.spPr.find(qn("a:prstGeom"))
    # default is a plain rectangle (or no explicit geometry) — never the frame preset
    assert g is None or g.get("prst") == "rect"
    print("ok  geom=None leaves the picture a plain rectangle")


def test_nature_scenes_size_and_mode():
    for kind in NI.SCENES:
        im = NI.generate(kind, 240, 343, seed=3)
        assert im.size == (240, 343), (kind, im.size)
        assert im.mode == "RGB", (kind, im.mode)
    print(f"ok  {len(NI.SCENES)} nature scenes render at the requested size")


def test_saved_image_reopens_at_frame_aspect():
    # aspect the caller must feed (frame w/h) — generator honours exact px size
    w, h = 700, int(700 / (FRAME[2] / FRAME[3]))
    p = _tmp_img(w, h, "mountains")
    from PIL import Image
    assert Image.open(p).size == (w, h)
    print("ok  generated image matches the requested frame-aspect size")


def test_cover_crop_matches_frame_aspect_without_stretch():
    from PIL import Image
    frame_aspect = FRAME[2] / FRAME[3]  # ~0.700
    src = _tmp_img(500, 683, "forest")  # ~0.732, wrong for the frame
    assert abs(Image.open(src).size[0] / Image.open(src).size[1] - frame_aspect) > 0.02
    dst = os.path.join(tempfile.gettempdir(), "_cropped_forest.png")
    FI.cover_crop_to_aspect(src, dst, frame_aspect)
    w, h = Image.open(dst).size
    assert abs(w / h - frame_aspect) < 0.005, (w, h, w / h, frame_aspect)
    # cover-crop trims one axis only; height is preserved when the source is too wide
    assert h == 683 and w < 500
    print("ok  cover_crop_to_aspect matches the frame aspect (no stretch)")


def _inject_layout_shape(slide, xml):
    """Append a raw shape to the slide's LAYOUT spTree (for the audit tests)."""
    from pptx.oxml import parse_xml
    spTree = slide.slide_layout._element.find(
        qn("p:cSld") + "/" + qn("p:spTree"))
    spTree.append(parse_xml(xml))


def test_frame_obstructions_flags_edge_line_and_clean_is_empty():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    fl, ft, fw, fh = FRAME
    # clean to start
    assert FI.frame_obstructions(slide, *FRAME) == []
    # a stray dark vertical line sitting exactly on the frame's right edge
    x_edge = fl + fw
    line = (
        f'<p:cxnSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        f'xmlns:a="{FI._A}"><p:nvCxnSpPr>'
        f'<p:cNvPr id="909" name="StrayLine"/><p:cNvCxnSpPr/><p:nvPr/></p:nvCxnSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{x_edge}" y="{ft + 100000}"/>'
        f'<a:ext cx="0" cy="1500000"/></a:xfrm>'
        f'<a:prstGeom prst="straightConnector1"><a:avLst/></a:prstGeom>'
        f'<a:ln w="9525"><a:solidFill><a:srgbClr val="0E2356"/></a:solidFill></a:ln>'
        f'</p:spPr></p:cxnSp>'
    )
    _inject_layout_shape(slide, line)
    found = FI.frame_obstructions(slide, *FRAME)
    ids = {f["id"] for f in found}
    assert "909" in ids, found
    assert any(f["stroke"] is True for f in found)
    print("ok  frame_obstructions flags a stray edge line and is empty when clean")


def test_frame_obstructions_ignores_covered_inner_shape():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fl, ft, fw, fh = FRAME
    # a stroked rect well inside the frame -> fully covered by the picture -> ignore
    inner = (
        f'<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        f'xmlns:a="{FI._A}"><p:nvSpPr>'
        f'<p:cNvPr id="910" name="InnerCovered"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr><a:xfrm><a:off x="{fl + fw // 3}" y="{ft + fh // 3}"/>'
        f'<a:ext cx="{fw // 4}" cy="{fh // 4}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f'<a:ln w="9525"><a:solidFill><a:srgbClr val="0E2356"/></a:solidFill></a:ln>'
        f'</p:spPr></p:sp>'
    )
    _inject_layout_shape(slide, inner)
    ids = {f["id"] for f in FI.frame_obstructions(slide, *FRAME)}
    assert "910" not in ids, "inner covered shape should not be flagged"
    print("ok  frame_obstructions ignores a shape fully covered by the frame")


def main():
    fns = [v for k, v in sorted(globals().items()) if k.startswith("test_")]
    for fn in fns:
        fn()
    print(f"\nALL {len(fns)} TESTS PASSED")


if __name__ == "__main__":
    main()
