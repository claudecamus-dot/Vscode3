"""Insert an image into a template *frame* so it takes the frame's EXACT shape.

Template placeholders in the OCTO deck are not plain rectangles: the "cadre
blanc" visual frame is a ``round2DiagRect`` (one diagonal pair of corners
rounded, the other square). Rounding the PNG with a uniform radius in PIL never
matches that — the corners poke out or leave white gaps.

The robust fix here is geometry-level, not pixel-level: drop the picture at the
frame's rendered bounds, then **clone the frame's own ``<a:prstGeom>`` onto the
picture's ``<p:spPr>``**. PowerPoint then clips the image to the identical
preset shape (same preset, same adjustment values), so the image edges follow
the frame's rounded/diagonal corners exactly, for any preset — not just this
one.

Public API
----------
- ``frame_geometry(group_shape, inner_name)`` -> (left, top, width, height, prstGeom_el)
      Read a layout frame's rendered bounds + geometry element to clone.
- ``round2diag_geom(adj1=50000, adj2=0)`` -> prstGeom element (build one by hand).
- ``place_image_in_frame(slide, image_path, left, top, width, height, geom)``
      Insert the picture and clip it to ``geom`` (a prstGeom element).
- ``cover_crop_to_aspect(src, dst, aspect)`` -> dst
      Center-crop a source image to the frame's width/height ratio so the preset
      clip lands with no distortion (the "crop first" step, done for you).
- ``frame_obstructions(slide, left, top, width, height, margin=30000)`` -> list[dict]
      Audit the slide's *layout and master* for stroked/filled shapes that sit on
      or poke past the frame perimeter — the real source of partial borders,
      slivers and stray lines around a filled frame. Run it before rendering.

Why an obstruction audit
------------------------
The visible border/fill of an OCTO framed image does NOT live on the picture:
it comes from the frame shape (and its neighbours) in the slide **layout/master**.
Enlarge or refill a frame and you expose whatever the old picture used to cover —
the frame's own ``<a:ln>`` border, a stray ``straightConnector1`` at the edge, a
teardrop/donut poking one pixel past the new bounds. Geometry checks and a quick
glance at the render (dark photo content hides thin lines) both miss these.
``frame_obstructions`` reports them by ID so you can neutralise the right shape.
"""
import copy

from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.util import Emu

_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def round2diag_geom(adj1: int = 50000, adj2: int = 0):
    """Build a ``round2DiagRect`` preset-geometry element (values in 1/1000 %)."""
    xml = (
        f'<a:prstGeom xmlns:a="{_A}" prst="round2DiagRect"><a:avLst>'
        f'<a:gd name="adj1" fmla="val {int(adj1)}"/>'
        f'<a:gd name="adj2" fmla="val {int(adj2)}"/>'
        f"</a:avLst></a:prstGeom>"
    )
    return parse_xml(xml)


def _iter_leaf_shapes(shapes):
    for sh in shapes:
        if sh.shape_type == 6:  # GROUP
            yield from _iter_leaf_shapes(sh.shapes)
        else:
            yield sh


def frame_geometry(group_shape, inner_name):
    """Return ``(left, top, width, height, prstGeom_el)`` for a framed placeholder.

    ``group_shape`` is the layout group that positions the frame; ``inner_name``
    is the name of the shaped sub-shape carrying the ``prstGeom`` to clone. The
    *rendered* bounds come from the group (its child coordinate space may be
    translated/negative), the geometry comes from the inner shape.
    """
    left, top = group_shape.left, group_shape.top
    width, height = group_shape.width, group_shape.height
    geom = None
    for sub in _iter_leaf_shapes([group_shape]):
        if sub.name == inner_name:
            g = sub._element.spPr.find(qn("a:prstGeom"))
            if g is None:
                g = sub._element.spPr.find(qn("a:custGeom"))
            geom = copy.deepcopy(g) if g is not None else None
            break
    return left, top, width, height, geom


def _set_geometry(pic, geom_el):
    """Replace the picture's geometry child with a clone of ``geom_el``.

    Order in ``<p:spPr>`` matters (schema): ``xfrm`` then the geometry then
    fill/line — so we insert right after ``xfrm`` rather than appending.
    """
    spPr = pic._element.spPr
    for tag in ("a:prstGeom", "a:custGeom"):
        for e in spPr.findall(qn(tag)):
            spPr.remove(e)
    new = copy.deepcopy(geom_el)
    xfrm = spPr.find(qn("a:xfrm"))
    if xfrm is not None:
        xfrm.addnext(new)
    else:
        spPr.insert(0, new)
    return pic


def place_image_in_frame(slide, image_path, left, top, width, height, geom=None):
    """Add ``image_path`` at the given EMU bounds and clip it to ``geom``.

    ``geom`` is a ``prstGeom``/``custGeom`` element (e.g. from ``frame_geometry``
    or ``round2diag_geom``). If ``None``, the picture stays a plain rectangle.
    The image is expected to already match the frame's aspect ratio (crop it
    beforehand, e.g. with :func:`cover_crop_to_aspect`) so the preset clip lands
    cleanly without distortion.
    """
    pic = slide.shapes.add_picture(
        image_path, Emu(int(left)), Emu(int(top)), Emu(int(width)), Emu(int(height))
    )
    if geom is not None:
        _set_geometry(pic, geom)
    return pic


def cover_crop_to_aspect(src, dst, aspect):
    """Center-crop ``src`` to ``aspect`` (= width/height) and write it to ``dst``.

    "Cover" fit: keeps the whole shorter dimension and trims the overflow off the
    long one, preserving the source's proportions (no letterbox, no stretch).
    Feed the result to :func:`place_image_in_frame` so the preset clip is
    distortion-free. Returns ``dst``.
    """
    from PIL import Image

    im = Image.open(src)
    w, h = im.size
    cur = w / h
    if abs(cur - aspect) > 1e-4:
        if cur > aspect:  # too wide -> trim left/right
            new_w = int(round(h * aspect))
            x0 = (w - new_w) // 2
            im = im.crop((x0, 0, x0 + new_w, h))
        else:  # too tall -> trim top/bottom
            new_h = int(round(w / aspect))
            y0 = (h - new_h) // 2
            im = im.crop((0, y0, w, y0 + new_h))
    if im.mode not in ("RGB", "L"):
        im = im.convert("RGB")
    im.save(dst)
    return dst


def _stroke_visible(spPr):
    """True/False if the shape draws a line; None if it inherits (unknown)."""
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        return None
    if ln.find(qn("a:noFill")) is not None:
        return False
    for f in ("a:solidFill", "a:gradFill", "a:pattFill"):
        if ln.find(qn(f)) is not None:
            return True
    return None


def _fill_visible(spPr):
    for f in ("a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill"):
        if spPr.find(qn(f)) is not None:
            return True
    if spPr.find(qn("a:noFill")) is not None:
        return False
    return None


def _bounds(shape):
    try:
        if None in (shape.left, shape.top, shape.width, shape.height):
            return None
        return (shape.left, shape.top, shape.left + shape.width, shape.top + shape.height)
    except (TypeError, ValueError):
        return None


def frame_obstructions(slide, left, top, width, height, margin=30000):
    """Report layout/master shapes that would draw an edge on/around the frame.

    Scans the top-level shapes of the slide's layout and master for anything with
    a visible **stroke or fill** whose bounds touch the frame perimeter band
    (frame rect grown by ``margin`` EMU) while extending past the frame edge — i.e.
    a shape the filling picture does *not* fully cover. That is exactly what shows
    up as a partial border, a sliver, or a stray line beside a framed image.

    Returns a list of dicts: ``{source, id, name, prst, bounds, stroke, fill,
    reason}``. Empty list == the frame region is clean. Does not descend into
    groups (bounds inside a group are in the group's own coordinate space); the
    real-world offenders (frame border, edge connectors) are top-level.
    """
    fl, ft, fr, fb = left, top, left + width, top + height
    ex_l, ex_t, ex_r, ex_b = fl - margin, ft - margin, fr + margin, fb + margin
    eps = margin
    findings = []
    parts = [("layout", slide.slide_layout)]
    try:
        parts.append(("master", slide.slide_layout.slide_master))
    except Exception:
        pass
    for source, part in parts:
        for sh in part.shapes:
            spPr = getattr(sh._element, "spPr", None)
            if spPr is None:
                continue
            b = _bounds(sh)
            if b is None:
                continue
            sl, st, sr, sb = b
            # intersects the grown frame rect?
            if sr < ex_l or sl > ex_r or sb < ex_t or st > ex_b:
                continue
            # reaches the frame perimeter band (a stroke here draws on the edge)
            touches = (
                sl < fl + margin or st < ft + margin
                or sr > fr - margin or sb > fb - margin
            )
            # extends past the *actual* frame edge (fill there is not covered)
            pokes = sl < fl - eps or st < ft - eps or sr > fr + eps or sb > fb + eps
            stroke = _stroke_visible(spPr)
            fill = _fill_visible(spPr)
            if stroke is True and touches:
                reason = "stroked shape on the frame perimeter (draws a border/line)"
            elif fill is True and pokes:
                reason = "filled shape pokes past the frame edge"
            else:
                continue
            g = spPr.find(qn("a:prstGeom"))
            cnv = sh._element.find(qn("p:nvSpPr") + "/" + qn("p:cNvPr"))
            if cnv is None:  # pics/cxn use different nv element
                cnv = sh._element.find(".//" + qn("p:cNvPr"))
            findings.append({
                "source": source,
                "id": cnv.get("id") if cnv is not None else None,
                "name": cnv.get("name") if cnv is not None else None,
                "prst": g.get("prst") if g is not None else None,
                "bounds": b,
                "stroke": stroke,
                "fill": fill,
                "reason": reason,
            })
    return findings
