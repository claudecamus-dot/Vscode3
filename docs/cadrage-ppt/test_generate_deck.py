"""Tests fonctionnels du générateur de synthèse PPT (cadrage BMAD IAP).

Complète la vérification manuelle (rendu PowerPoint COM + relecture, cf.
`pptx-verify`) par des assertions rejouables sur le `.pptx` réellement généré
— même esprit que `test-export-ppt.py`/`test-ppt-charte.py` (VSCode1) et
`test_export_pptx_renders_cleanly_in_a_real_engine` (VSCode2), adapté à un
générateur script (pas une app web) :

  - structure : nombre de slides, géométrie (`D.verifier_geometrie`, déjà
    appelé par `build()`) ;
  - **cadres photo bien calés** : pour chaque chapitre + la slide vision,
    l'image posée doit avoir EXACTEMENT les bornes du cadre du template et
    porter le bon `prstGeom` cloné — pas juste "une image est présente
    quelque part". Trouvé la raison d'être de ce test : une image sur la
    slide 8 avait été jugée "pas bien calée" à l'œil (en fait une photo trop
    pâle qui se fondait dans le fond, cf. mémoire de session) ; ce test ne
    remplace pas l'œil (il ne juge pas la qualité de la photo) mais aurait
    immédiatement confirmé que le cadrage géométrique, lui, était correct —
    et détecterait un vrai décalage si l'un survenait.
  - aucun cadre laissé vide (« ici mettre une Photo » résiduel) ;
  - régression : l'encart numéro de chapitre ne doit jamais réhériter le
    retrait de puce du master (cf. `_sans_puce`, bug trouvé et corrigé) ;
  - obstructions de cadre limitées aux deux formes décoratives connues
    (badge logo/numéro) — une nouvelle obstruction serait un vrai défaut ;
  - aucune police générique (Arial/Calibri/...) explicitement posée ;
  - rendu réel via LibreOffice : le fichier s'ouvre et produit bien une page
    par slide (pas un fichier que python-pptx parse mais qu'aucun moteur
    n'ouvre proprement).

Usage : python test_generate_deck.py
"""
import os
import re
import subprocess
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

import generate_deck as gen

echecs = 0


def check(cond, msg):
    global echecs
    if cond:
        print(f"  ok   {msg}")
    else:
        echecs += 1
        print(f"  FAIL {msg}")


def _images(slide):
    return [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]


def _soffice_path():
    defaut = r"C:\Program Files\LibreOffice\program\soffice.exe"
    return defaut if os.path.exists(defaut) else "soffice"


def _verifier_rendu_reel(pptx_path, n_slides_attendu, tmp_dir):
    """Convertit en PDF via LibreOffice et compte les pages — même principe
    que le test VSCode2 éponyme : un .pptx qui parse avec python-pptx peut
    quand même être refusé/tronqué par un vrai moteur de rendu."""
    try:
        result = subprocess.run(
            [_soffice_path(), "--headless", "--convert-to", "pdf", "--outdir", tmp_dir, pptx_path],
            capture_output=True, timeout=120,
        )
    except (FileNotFoundError, subprocess.TimeoutExpired) as e:
        return None, f"LibreOffice indisponible ({e}) — vérification réelle non faite"
    if result.returncode != 0:
        return False, f"LibreOffice a échoué : {result.stderr.decode(errors='replace')[:300]}"
    pdf_path = os.path.join(tmp_dir, os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf")
    if not os.path.exists(pdf_path):
        return False, "LibreOffice n'a produit aucun PDF"
    pdf_bytes = open(pdf_path, "rb").read()
    if len(pdf_bytes) < 2000:
        return False, "PDF quasi vide — rendu suspect"
    page_count = len(re.findall(rb"/Type\s*/Page[^s]", pdf_bytes))
    if page_count != n_slides_attendu:
        return False, f"{page_count} page(s) rendue(s) pour {n_slides_attendu} slide(s) exportée(s)"
    return True, f"{page_count} page(s) rendues, conforme aux {n_slides_attendu} slides"


def main():
    problemes = gen.build()
    out = os.path.join(gen.HERE, "bmad-iap-cadrage-synthese.pptx")
    prs = Presentation(out)

    print("Structure :")
    check(len(prs.slides) == 29, f"29 slides — reçu {len(prs.slides)}")
    check(not problemes, f"géométrie propre (verifier_geometrie) — {len(problemes or [])} problème(s)")
    check(os.path.exists(out) and os.path.getsize(out) > 500_000,
          f"fichier .pptx écrit, taille plausible ({os.path.getsize(out) if os.path.exists(out) else 0} octets)")

    print("Cadres photo bien calés (chapitres — layout '50 - Chapitre', teardrop) :")
    chapitres = [4, 11, 18]
    for idx in chapitres:
        slide = prs.slides[idx - 1]
        cadre = gen._find_frame_by_geom(slide.slide_layout.shapes, "teardrop")
        images = _images(slide)
        check(len(images) == 1, f"slide {idx} : exactement 1 image posée (reçu {len(images)})")
        check(cadre is not None, f"slide {idx} : cadre teardrop trouvé sur le layout")
        if images and cadre:
            pic = images[0]
            l, t, w, h, _ = cadre
            check((pic.left, pic.top, pic.width, pic.height) == (l, t, w, h),
                  f"slide {idx} : image alignée exactement sur le cadre "
                  f"(image=({pic.left},{pic.top},{pic.width},{pic.height}) vs cadre=({l},{t},{w},{h}))")
            g = pic._element.spPr.find(qn("a:prstGeom"))
            check(g is not None and g.get("prst") == "teardrop",
                  f"slide {idx} : image clippée au bon preset (teardrop)")

    print("Cadre photo bien calé (slide vision — layout 'cadre blanc', round2DiagRect) :")
    slide_vision = prs.slides[2]
    cadre_vision = gen._find_frame_in_group(
        slide_vision.slide_layout.shapes, "Google Shape;212;p17", "Google Shape;213;p17")
    images_vision = _images(slide_vision)
    check(len(images_vision) == 1, f"slide 3 (vision) : exactement 1 image posée (reçu {len(images_vision)})")
    check(cadre_vision is not None, "slide 3 (vision) : cadre 'cadre blanc' trouvé sur le layout")
    if images_vision and cadre_vision:
        pic = images_vision[0]
        l, t, w, h, _ = cadre_vision
        check((pic.left, pic.top, pic.width, pic.height) == (l, t, w, h),
              f"slide 3 (vision) : image alignée exactement sur le cadre "
              f"(image=({pic.left},{pic.top},{pic.width},{pic.height}) vs cadre=({l},{t},{w},{h}))")
        g = pic._element.spPr.find(qn("a:prstGeom"))
        check(g is not None and g.get("prst") == "round2DiagRect",
              "slide 3 (vision) : image clippée au bon preset (round2DiagRect)")

    print("Aucun cadre laissé vide (texte gabarit « ici mettre une Photo » résiduel) :")
    texte_complet = "\n".join(
        shp.text_frame.text for slide in prs.slides for shp in slide.shapes if shp.has_text_frame)
    check("ici mettre une Photo" not in texte_complet, "aucun texte gabarit de cadre photo résiduel")

    print("Régression — encart numéro de chapitre sans retrait de puce hérité (bug trouvé/corrigé) :")
    for idx in chapitres:
        slide = prs.slides[idx - 1]
        ph = next((p for p in slide.placeholders if p.placeholder_format.idx == 1), None)
        check(ph is not None, f"slide {idx} : placeholder numéro (idx=1) présent")
        if ph is not None:
            p_el = ph.text_frame.paragraphs[0]._p
            pPr = p_el.find(qn("a:pPr"))
            marL = pPr.get("marL") if pPr is not None else None
            indent = pPr.get("indent") if pPr is not None else None
            bu_none = pPr.find(qn("a:buNone")) is not None if pPr is not None else False
            check(marL == "0" and indent == "0" and bu_none,
                  f"slide {idx} : numéro sans retrait de puce hérité (marL={marL}, indent={indent}, buNone={bu_none})")

    print("Obstructions de cadre — limitées aux 2 formes décoratives connues (badge logo/numéro) :")
    attendus = {"Google Shape;37;p4", "Google Shape;54;p4"}
    for idx in chapitres:
        slide = prs.slides[idx - 1]
        cadre = gen._find_frame_by_geom(slide.slide_layout.shapes, "teardrop")
        if cadre:
            obstructions = gen.frame_obstructions(slide, *cadre[:4])
            noms = {o["name"] for o in obstructions}
            check(noms <= attendus, f"slide {idx} : pas d'obstruction inattendue (trouvé {noms - attendus or 'aucune'})")

    print("Police — aucune police générique explicitement posée (Arial/Calibri/...) :")
    generiques = {"arial", "calibri", "times new roman", "segoe ui"}
    trouve = set()
    for slide in prs.slides:
        for shp in slide.shapes:
            if shp.has_text_frame:
                for p in shp.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.name and r.font.name.lower() in generiques:
                            trouve.add(r.font.name)
    check(not trouve, f"aucune police générique posée{' (trouvé ' + str(sorted(trouve)) + ')' if trouve else ''}")

    print("Rendu réel (LibreOffice — conversion PDF, comptage de pages) :")
    import tempfile
    with tempfile.TemporaryDirectory(prefix="test-ppt-render-") as tmp:
        ok, detail = _verifier_rendu_reel(out, len(prs.slides), tmp)
        if ok is None:
            print(f"  --   {detail}")
        else:
            check(ok, detail)

    print("\nTOUS LES TESTS PASSENT" if echecs == 0 else f"\n{echecs} TEST(S) EN ECHEC")
    sys.exit(0 if echecs == 0 else 1)


if __name__ == "__main__":
    main()
