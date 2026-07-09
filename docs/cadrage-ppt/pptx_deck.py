"""pptx_deck — petite bibliotheque d'aide pour construire des slides python-pptx
"de qualite" : echelle typographique coherente, formes (barres, jauge, cartes),
couleurs, et surtout un controle geometrique automatique (`verifier_geometrie`)
qui detecte toute forme qui sort de la slide — le defaut classique des decks
generes a la main.

Reutilisable hors de ce projet : aucune dependance au domaine metier ici.
Les coordonnees des helpers sont exprimees en POUCES (float) pour la lisibilite.
"""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# --- Echelle typographique (pt) — une seule source de verite ---
TYPE = {
    "title": 26, "h2": 18, "h3": 14, "body": 12, "small": 10.5, "tiny": 9,
    "kpi": 44, "kpi_unit": 16,
}

# Palette des piliers : IDENTIQUE au radar web (radar-svg.js) pour que les
# barres et le radar parlent le meme langage couleur.
PALETTE = ["#2c5cc5", "#1e6b34", "#b3261e", "#b8860b", "#6a3d9a", "#138086"]

INK = "#1c2330"
MUTED = "#6b7280"
LINE = "#e6e8ee"
TRACK = "#eef1f7"
OK = "#1e6b34"
WARN = "#b3261e"
GOLD = "#b8860b"


def rgb(hexa):
    return RGBColor.from_string(hexa.lstrip("#"))


def couleur_pilier(i):
    return PALETTE[i % len(PALETTE)]


def _no_shadow(shape):
    # Les autoshapes heritent parfois d'une ombre du theme : on la coupe.
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def add_text(slide, l, t, w, h, lignes, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT,
             wrap=True):
    """Ajoute une zone de texte. `lignes` = liste de (texte, opts) ; chaque
    element devient un paragraphe. opts: size,bold,italic,color,align,
    space_before,space_after,line_spacing."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    for i, (texte, opts) in enumerate(lignes):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get("align", align)
        if "space_before" in opts:
            p.space_before = Pt(opts["space_before"])
        if "space_after" in opts:
            p.space_after = Pt(opts["space_after"])
        if "line_spacing" in opts:
            p.line_spacing = opts["line_spacing"]
        run = p.add_run()
        run.text = texte
        f = run.font
        f.size = Pt(opts.get("size", TYPE["body"]))
        f.bold = opts.get("bold", False)
        f.italic = opts.get("italic", False)
        f.color.rgb = rgb(opts.get("color", INK))
    return box


def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=1.0, rounded=False,
             radius=0.12):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    _no_shadow(shp)
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = rgb(line)
        shp.line.width = Pt(line_w)
    shp.text_frame.paragraphs[0].text = ""
    return shp


def add_hbar(slide, l, t, w, h, frac, fill, track=TRACK):
    """Barre de progression horizontale (piste + remplissage), coins arrondis."""
    frac = max(0.0, min(1.0, float(frac)))
    add_rect(slide, l, t, w, h, fill=track, rounded=True, radius=0.5)
    if frac > 0:
        wv = max(h, w * frac)  # largeur minimale visible = hauteur (pastille)
        add_rect(slide, l, t, wv, h, fill=fill, rounded=True, radius=0.5)


def add_gauge(slide, l, t, size, frac, fill, track=TRACK, hole=62):
    """Jauge circulaire (anneau) via un graphique doughnut a 2 segments.
    Renvoie le GraphicFrame. Le libelle central est a poser separement."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    frac = max(0.0, min(1.0, float(frac)))
    data = CategoryChartData()
    data.categories = ["v", "r"]
    data.add_series("g", (frac, 1 - frac))
    gf = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(l), Inches(t),
                                Inches(size), Inches(size), data)
    chart = gf.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    plot.has_data_labels = False
    # Taille du trou
    dough = plot._element  # <c:doughnutChart>
    hs = dough.find(qn("c:holeSize"))
    if hs is None:
        hs = dough.makeelement(qn("c:holeSize"), {"val": str(hole)})
        dough.append(hs)
    else:
        hs.set("val", str(hole))
    # Couleurs des 2 segments
    pts = plot.series[0].points
    for pt_, col in ((pts[0], fill), (pts[1], track)):
        pt_.format.fill.solid()
        pt_.format.fill.fore_color.rgb = rgb(col)
        pt_.format.line.color.rgb = rgb("#ffffff")
        pt_.format.line.width = Pt(1)
    return gf


def add_card(slide, l, t, w, h, accent):
    """Carte blanche a coins arrondis + liseré couleur a gauche (style infographie)."""
    add_rect(slide, l, t, w, h, fill="#ffffff", line=LINE, line_w=0.75,
             rounded=True, radius=0.06)
    add_rect(slide, l, t, 0.07, h, fill=accent, rounded=True, radius=0.5)


def add_dot(slide, x, y, d, color):
    """Petite pastille ronde pleine (puce de legende / marqueur de chip)."""
    return add_rect(slide, x, y, d, d, fill=color, rounded=True, radius=0.5)


def add_range_bar(slide, l, t, w, h, mn, mx, scale_max, fill, marker=None,
                  track=TRACK):
    """Barre d'amplitude min..max sur une echelle 0..scale_max (piste complete +
    segment colore couvrant la plage). `marker` (ex. moyenne) pose un repere
    vertical. Sert a montrer une dispersion sur l'echelle reelle, pas en relatif."""
    add_rect(slide, l, t, w, h, fill=track, rounded=True, radius=0.5)
    fa = max(0.0, min(1.0, mn / scale_max))
    fb = max(0.0, min(1.0, mx / scale_max))
    seg_w = max(h, w * (fb - fa))  # largeur minimale = hauteur (pastille)
    add_rect(slide, l + w * fa, t, seg_w, h, fill=fill, rounded=True, radius=0.5)
    if marker is not None:
        fm = max(0.0, min(1.0, marker / scale_max))
        mx_x = l + w * fm - 0.015
        add_rect(slide, mx_x, t - 0.05, 0.03, h + 0.10, fill=INK, rounded=True,
                 radius=0.5)


def _compte_lignes(texte, cpl):
    """Nombre de lignes apres un repli mot-a-mot pour une largeur de `cpl`
    caracteres. Estimateur volontairement simple (pas de mesure de police reelle) :
    calibre empiriquement via `cpl`, pas cense etre pixel-parfait."""
    total = 0
    for para in str(texte).split("\n"):
        cur = 0
        n = 1
        for mot in para.split():
            ajout = (1 if cur else 0) + len(mot)
            if cur + ajout > cpl and cur:
                n += 1
                cur = len(mot)
            else:
                cur += ajout
        total += n
    return total


def estimer_lignes(texte, largeur_in, taille_pt, cpi_ref=11.0, taille_ref=10.5):
    """Estime le nombre de lignes qu'occupera `texte` une fois reparti mot-a-mot
    sur `largeur_in` pouces a la taille de police `taille_pt`. Les caracteres par
    pouce sont derives de `cpi_ref` (calibre a `taille_ref` pt) par une regle de
    trois : une police 2x plus petite loge environ 2x plus de caracteres/pouce."""
    if not texte:
        return 1
    cpi = cpi_ref * (taille_ref / taille_pt)
    cpl = max(6, int(largeur_in * cpi))
    return _compte_lignes(texte, cpl)


def ajuster_police(textes, largeur_in, taille_max, taille_min, budget_ok, pas=0.5,
                   cpi_ref=11.0, taille_ref=10.5):
    """Adapte la taille de police a la longueur des phrases a restituer : cherche,
    par pas de `pas` pt entre `taille_max` et `taille_min`, la plus GRANDE taille
    telle que `budget_ok(taille, lignes_max)` soit vrai — ou `lignes_max` est le
    nombre de lignes necessaires au plus long de `textes` une fois reparti sur
    `largeur_in` pouces a cette taille (voir `estimer_lignes`).

    `budget_ok` encapsule la contrainte geometrique propre a l'appelant (ex. :
    n cartes empilees doivent tenir dans la bande disponible) — cette fonction
    reste agnostique du domaine. Objectif : ne JAMAIS tronquer/faire deborder une
    phrase — si aucune taille ne satisfait `budget_ok`, on degrade sur
    `taille_min` (texte tres dense) plutot que de laisser un texte coupe.

    Renvoie (taille, lignes_max)."""
    taille = taille_max
    while True:
        lignes_max = max((estimer_lignes(t, largeur_in, taille, cpi_ref, taille_ref)
                          for t in textes), default=1)
        if budget_ok(taille, lignes_max) or taille <= taille_min:
            return (max(taille, taille_min), lignes_max)
        taille = max(taille_min, round(taille - pas, 2))


def tronquer_a_lignes(texte, largeur_in, taille_pt, max_lignes, cpi_ref=11.0,
                      taille_ref=10.5):
    """Tronque `texte` (avec une ellipse finale) pour qu'il tienne dans
    `max_lignes` lignes une fois reparti sur `largeur_in` pouces a `taille_pt`.
    Dernier recours quand meme `ajuster_police` a sa taille plancher ne suffit
    plus a eviter un debordement geometrique — mieux vaut un texte coupe
    proprement qu'une forme qui deborde de la slide. Ne fait rien si le texte
    tient deja dans `max_lignes`."""
    if estimer_lignes(texte, largeur_in, taille_pt, cpi_ref, taille_ref) <= max_lignes:
        return texte
    cpi = cpi_ref * (taille_ref / taille_pt)
    cpl = max(6, int(largeur_in * cpi))
    limite = max(1, cpl * max_lignes - 1)
    tronque = str(texte)[:limite].rstrip()
    dernier_espace = tronque.rfind(" ")
    if dernier_espace > limite * 0.6:
        tronque = tronque[:dernier_espace]
    return tronque.rstrip(" ,;:.") + "…"


def verifier_geometrie(prs, marge_in=0.02):
    """Retourne la liste des problemes : toute forme dont les bords depassent la
    slide (au-dela d'une petite marge de tolerance). Liste vide = OK."""
    W, H = prs.slide_width, prs.slide_height
    tol = Inches(marge_in)
    problemes = []
    for si, slide in enumerate(prs.slides, start=1):
        for shp in slide.shapes:
            try:
                l, t, w, h = shp.left, shp.top, shp.width, shp.height
            except Exception:
                continue
            if None in (l, t, w, h):
                continue
            nom = shp.name or "shape"
            if l < -tol or t < -tol or (l + w) > W + tol or (t + h) > H + tol:
                problemes.append(
                    f"slide {si}: '{nom}' hors cadre "
                    f"(l={Emu(l).inches:.2f} t={Emu(t).inches:.2f} "
                    f"r={Emu(l + w).inches:.2f} b={Emu(t + h).inches:.2f} ; "
                    f"slide {Emu(W).inches:.2f}x{Emu(H).inches:.2f})")
    return problemes


def theme_colors(prs):
    """Lit le nuancier du theme (dk1/lt1/dk2/lt2/accent1..6) du 1er master et le
    renvoie en dict {nom: '#RRGGBB'}. Sert a adapter le deck a la charte du
    template fourni (couleur d'accent de marque) sans rien coder en dur.
    Renvoie {} si le theme est introuvable (l'appelant prevoit un repli)."""
    import re
    try:
        part = prs.slide_masters[0].part
        theme_part = next((r.target_part for r in part.rels.values()
                           if "theme" in r.reltype), None)
        if theme_part is None:
            return {}
        xml = theme_part.blob.decode("utf-8", "ignore")
    except Exception:
        return {}
    m = re.search(r"<a:clrScheme.*?</a:clrScheme>", xml, re.S)
    if not m:
        return {}
    seg = m.group(0)
    out = {}
    for name in ("dk1", "lt1", "dk2", "lt2", "accent1", "accent2",
                 "accent3", "accent4", "accent5", "accent6"):
        mm = re.search(
            r"<a:" + name + r">.*?(?:srgbClr val=\"([0-9A-Fa-f]{6})\"|"
            r"sysClr[^>]*lastClr=\"([0-9A-Fa-f]{6})\")", seg, re.S)
        if mm:
            out[name] = "#" + (mm.group(1) or mm.group(2)).upper()
    return out
