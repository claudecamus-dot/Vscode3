"""Génère une synthèse PPT (46 slides) des RÉSULTATS du cadrage BMAD IAP
(docs/bmad-iap-cadrage.md) à partir des helpers pptx_deck, dessinée
PAR-DESSUS le vrai template de marque OCTO (template-octo.pptx) —
masters/layouts/thème conservés, pas un deck sur canevas vierge.

Structure v2.9 en 9 chapitres, sur le fil rouge narratif des decks SCALE
(docs/Import/notes-extraction-scale.md : POURQUOI → QUI → QUOI → COMMENT →
RÉSULTAT) : 01 Exec summary (le pitch de l'offre en 3 faces, puis la démarche
avec ou sans agentic — en ouverture)
· 02 Contexte (pourquoi) · 03 Personas (qui) · 04 Besoins & douleurs
(le pourquoi, mesuré) · 05 Proposition (quoi — thèse `why_iap` en ouverture,
méthode scorée, cible d'organisation)
· 06 IA (le quoi, côté IA — sous gate) · 07 Démarche (comment — trajectoire
fusionnée avec la vue bout-en-bout, fil humain, activités humaines avec/sans
l'outil (v2.6), parcours de mission (v2.9), schéma de fonctionnement,
inventaire des agents, livrables)
· 08 Outillage IAP (avec quoi — ouvre sur le schéma d'architecture en contexte
client (v2.6), la slide « ce que le module met dans les mains du consultant »
annoncée par le plan v2.5, puis ambition A/B/C et lien SI) · 09 KPI (la preuve
— 3 familles → mise en place → grille de maturité
→ cas chiffré, clôture). L'IA reste tirée APRÈS la proposition (doctrine :
« l'IA amplifie l'organisation, elle n'est jamais la réponse à un problème
d'abord organisationnel ») ; la Démarche vient APRÈS l'IA pour que le
« comment » enchaîne directement sur l'outillage puis la preuve.
L'executive summary (slide 2) reprend le même fil en 5 blocs
OFFRE/POURQUOI/QUOI/COMMENT/RÉSULTAT avec renvoi aux chapitres (v2.10 : le
bloc OFFRE, qui référence le chapitre 01, manquait — le sommaire ne
mentionnait aucun chapitre avant "02–04").

v2.6 : le sous-chapitre « Exemples » de la Proposition (séparateur + 3 slides
illustratives) est SUPPRIMÉ à la demande — git garde l'historique (v2.5) ; en
échange, 2 slides nouvelles (activités humaines de la démarche ; architecture
IAP en contexte client) et un badge de série « déploiement agentic chez le
client » sur les 4 slides de proposition agentic du chapitre IA (3 agents
candidats + export markdown), renvoyant au schéma du chapitre Outillage IAP.

v2.7 (2026-09-01) : 2 slides de plus (40 -> 42) sur ce que le cadrage
déclarait faisant foi pour le deck sans y être redescendu — « qui achète,
contre quoi » (les 4 achats alternatifs, chapitre Contexte, après les
déclencheurs) et « conditions de réussite et non-engagement » (ce que la
mission exige du client et ce que son absence déclenche, chapitre Démarche,
après le fil humain).

v2.8 (2026-09-02) : nouveau chapitre 01 « Exec summary », EN OUVERTURE du
deck (avant slide_vision, juste après le sommaire) — 3 slides de plus
(42 -> 45), tous les chapitres suivants glissent de +1. Reprend la slide 2 du
pitch source (l'offre : chapô + citation-thèse, verbatim) et son schéma du
parcours de mission — redessiné EN NATIF (pas une insertion d'image), sur 3
registres (mouvements du socle toujours présents, variantes conditionnées au
contexte, mécanismes additifs) — puis une synthèse en une page de l'offre qui
résume aussi le reste du deck (COMPRENDRE/DÉFINIR/FAIRE ADOPTER & PROUVER,
avec renvoi aux chapitres). Le nœud « Discovery gaspillages » du schéma garde
« 6 catégories » tel quel dans le document source, alors que le chapitre
Besoins & douleurs du deck compte 8 familles de gaspillage — divergence DU
DOCUMENT SOURCE (marqué WIP), remontée pour arbitrage, non corrigée ici.

v2.9 (2026-09-02) : le chapitre 01 « Exec summary » est REFONDU pour parler au
prospect plutôt que de résumer le deck deux fois (45 -> 46 slides ; 9 chapitres
inchangés). Arbitrages : `slide_offre_synthese` est SUPPRIMÉE — le sommaire du
deck reste `slide_executive_summary`, qui la précédait et ne bouge pas ;
`slide_offre_iap` (le grand schéma du parcours de mission) DÉMÉNAGE au chapitre
07 · Démarche, en tête du bloc des schémas, avec le kicker et l'or du chapitre
d'accueil. À leur place, deux slides neuves :
  - `slide_pitch_iap` — trois cartes ÉGALES à ne pas confondre (les douleurs de
    ces organisations ; notre démarche outillée par un module agentic, côté
    CONSULTANT ; déployer de l'agentic chez le CLIENT, en option sous gate IA),
    deux teintes et deux silhouettes distinctes pour les deux faces de l'agentic
    — un sponsor qui lit deux fois « agentic » sans étiquette croit qu'on lui
    vend la même chose deux fois, ce qui rendrait la slide fausse. Sous la
    rangée, une bande grise pleine largeur (poids visuel moindre, pas une 4e
    carte) sur le MATÉRIAU de cadrage : trois contextes sectoriels réels, sans
    jamais affirmer une mission — aucun REX n'existe dans les sources
    (docs/bmad-iap-cadrage.md:115 décrit une note de rédaction, `rex-library.md`
    est planifié et non peuplé).
  - `slide_demarche_avec_sans_agentic` — le niveau ZOOMÉ-ARRIÈRE de
    l'« avec ou sans agentic » : une ligne horizontale de trois temps
    (COMPRENDRE / DÉFINIR / FAIRE ADOPTER & PROUVER) et, sous chaque temps, des
    pastilles empilées (socle gris-navy, module teal, agentic client violet).
    Ne redessine ni `slide_activites_humaines` (phase par phase, ch. 07), ni
    `slide_iap_contexte_client` (topologie, ch. 08), ni `slide_export_markdown`
    (bifurcation, ch. 06) — renvois littéraux « décliné chapitre 07 » et
    « déployé chapitre 08 » portés sur la slide.

v2.10 (2026-09-03) : le chapitre 01 · Exec summary (slide_pitch_iap +
slide_demarche_avec_sans_agentic) était déjà groupé par son propre intercalaire
(`slide_chapitre(prs, "01", ...)`), mais restait invisible du sommaire
(`slide_executive_summary`) — sa rangée POURQUOI/QUOI/COMMENT/RÉSULTAT ne
renvoyait qu'aux chapitres 02 à 09. Un bloc OFFRE, accroché en tête de rangée
(même style « accent » — fond plein — que RÉSULTAT, pour border la rangée :
l'offre ouvre, la preuve ferme), renvoie maintenant au chapitre 01 et reprend
au mot près les deux sous-titres des slides qu'il annonce.

v2.11 (2026-09-03) : deux défauts graphiques signalés par relecture réelle
(vraie taille de deck, PAS le self-check géométrique) — départagés en
comparant le rendu LibreOffice ET un rendu PowerPoint réel (COM) sur les
mêmes slides, pour ne corriger que ce qui existe dans l'artefact que
l'utilisateur ouvre :
  - Slide 1 (couverture) : le trait qui semble mal rejoindre le coin arrondi
    du bandeau version est un ARTEFACT DE RENDU LIBREOFFICE (groupe pivoté à
    180° du template `template-octo.pptx`, composé différemment par les deux
    moteurs) — absent du rendu PowerPoint réel. Rien à corriger côté
    générateur ; toucher le template partagé (masters/layouts hors périmètre
    de ce script) aurait été le mauvais geste pour un défaut qui n'existe pas
    dans le livrable réel.
  - Slide 4 (`slide_pitch_iap`) : DÉFAUT RÉEL, confirmé dans les deux rendus.
    La carte « CE QU'ILS VIVENT » a 4 puces (4 personas) quand les 2 autres
    cartes n'en ont que 3, mais le budget vertical `dispo` était partagé sans
    tenir compte du nombre d'items — la 4e puce débordait sous le chip de
    pied, dont le fond plein la masquait entièrement (seul son disque de
    puce dépassait, visible comme une virgule rouge au-dessus du bouton).
    Corrigé par une taille/interligne de puce ADAPTATIVE (calculée par carte
    selon ce qui tient réellement dans `dispo`, avec marge de sécurité),
    plutôt que par un chiffre choisi à la main — corrige la classe de bug,
    pas seulement cette instance. Le mécanisme d'anomalies de build
    (`_ANOMALIES_BUILD`, jusque-là réservé aux photos manquantes) est
    généralisé à ce type de débordement : si aucune taille ne suffit même au
    plancher, le build le signale désormais comme un vrai défaut au lieu
    d'un print perdu — c'est l'évolution du check graphique demandée.

v2.12 (2026-09-03) : deuxième passe sur les mêmes 4 slides signalées, cette
fois avec un rendu PowerPoint réel de CHAQUE slide (pas un échantillon) —
deux défauts réels supplémentaires trouvés, tous deux présents dans les DEUX
moteurs de rendu (donc jamais des artefacts LibreOffice comme la slide 1) :
  - `slide_executive_summary` (slide 2) était le SEUL appel `content_slide()`
    de tout le générateur (~30 autres) sans `color=` explicite — son kicker
    retombait sur l'accent cyan générique au lieu du NAVY du chapitre 01
    qu'elle ouvre, déjà porté par ses deux slides suivantes. Corrigé en
    passant `color=NAVY`.
  - `slide_cover` (slide 1) affichait une version gelée sur "v2.8 ·
    2026-09-02" depuis 4 bumps de version consécutifs (v2.9 à v2.11) — la
    chaîne était écrite en dur dans la fonction plutôt que dérivée d'une
    source unique. Introduit `VERSION_DECK`/`DATE_VERSION_DECK` (constantes
    de module) + un test de régression (`test_generate_deck.py`) qui compare
    `VERSION_DECK` à la dernière entrée "vX.Y" du docstring de ce module et
    au texte réellement posé sur le placeholder de couverture — ce test a
    lui-même détecté l'écart en cours d'écriture de cette entrée.

v2.13 (2026-09-03) : diagnostic étage 2 rafraîchi (2 jours périmé) a trouvé
`content_slide(prs, kicker, title, color=None)` — sur 34 vrais appels
(37 occurrences moins 3 en commentaire), TOUS passent déjà `color=`
explicitement depuis ce module. Le repli `color or ACCENT` ne protégeait
donc plus personne : sa seule fonction résiduelle était de laisser un futur
appel oublié retomber en cyan silencieux, exactement le défaut de v2.12.
`color` devient un paramètre obligatoire (zéro régression mesurée, 34/34
appels déjà conformes) — un appel qui l'omettrait échoue maintenant au
build (`TypeError`), pas au rendu.

`slide_executive_summary` DÉMÉNAGE (arbitrage utilisateur 2026-09-03) de
juste après la couverture à juste après l'intercalaire du chapitre 01 —
modifié manuellement par l'utilisateur sur l'export, reporté dans `build()`
pour que toute régénération le conserve.

Brouillon `slide_synthese_pourquoi_quoi_comment` ajouté (demande
utilisateur) mais **NON câblé dans `build()`** — généré et vérifié en
isolation pour validation avant intégration, sur consigne explicite
("pour l'instant ne génère que cette slide").

Séparateurs : chapitres = intercalaire teardrop (photo + numéro, layout dédié) ;
sous-chapitres = `slide_sous_chapitre` (bloc-titre léger, sans photo ni numéro —
sans appelant depuis la v2.6, machinerie conservée).

Centré sur les résultats du cadrage (mission, doctrine, méthode, maturité,
ambition, KPIs, schéma de fonctionnement) plutôt que sur tout le détail de
mise en œuvre. Le commit 4f0c9b7 avait retiré ce détail d'implémentation ;
l'arbitrage utilisateur du 2026-07-21 rouvre ce périmètre sur DEUX points
précis seulement, désormais dans le deck :
  - l'architecture des 11 agents-workflows (slide_architecture_agents,
    inventaire des composants — complémentaire du schéma de flux, pas un
    doublon) ;
  - l'étude des personas / product discovery (slide_personas, réouverture de
    la discovery fusionnée en MVP1, §Décision de cadrage ligne 236).
Restent hors périmètre (toujours du support interne, pas une synthèse
exécutive) : le schéma des workflows détaillé, la roadmap MVP et les points
ouverts — les trois autres slides retirées au même commit ne sont PAS
réintroduites.

Usage : python generate_deck.py
Sortie : bmad-iap-cadrage-synthese.pptx (à côté de ce script).
"""
import sys, os
sys.path.insert(0, os.path.dirname(__file__))
import pptx_deck as D
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn

# Source unique du numero de version affiche sur la couverture (slide_cover) —
# jusqu'a v2.12 c'etait une chaine gelee dans slide_cover, jamais mise a jour a
# 4 bumps de version consecutifs (v2.9 a v2.11 ont toutes laisse "v2.8 · date
# perimee" sur la SLIDE LA PLUS VISIBLE du deck). Un seul endroit a changer
# desormais.
VERSION_DECK = "v2.13"
DATE_VERSION_DECK = "2026-09-03"

HERE = os.path.dirname(__file__)
TEMPLATE = os.path.join(HERE, "template-octo.pptx")
REPO_ROOT = os.path.abspath(os.path.join(HERE, "..", ".."))
sys.path.insert(0, os.path.join(REPO_ROOT, ".claude", "skills", "pptx-framed-image", "scripts"))
from framed_image import place_image_in_frame, frame_obstructions  # noqa: E402
import nature_images  # noqa: E402
import stock_images  # noqa: E402

IMG_DIR = os.path.join(HERE, "_img")
os.makedirs(IMG_DIR, exist_ok=True)
IMG_MANIFEST = os.path.join(HERE, "images-manifest.json")

LAYOUT_COUVERTURE = 8   # "40 - Couverture [1]" — idx0 titre, idx1 sous-titre, idx2/idx3 crédit+date
LAYOUT_TITRE_SEUL = 5   # "04 - Titre seul" — idx0 titre, garde logo/pied de page/n° de slide
LAYOUT_VIDE = 0         # "06 - Slide vide" — pas de placeholder, juste logo + badge de pagination
LAYOUT_CHAPITRE = 2     # "50 - Chapitre [1]" — idx0 titre (grand), idx1 numéro ; cadre photo teardrop
LAYOUT_VISUEL_DROITE = 15  # "63 - Titre, contenu et visuel à droite - cadre blanc"

# --- Géométrie du template OCTO réel (10 x 5.625 in, 16:9) — cf.
# docs/vscode1-export/template-octo.md §4-5, vérifiée localement contre
# template-octo.pptx (mêmes dims/layouts/thème). Contenu dessiné dans la
# zone de contenu du layout « Titre seul » (sous le titre, au-dessus du
# pied de page), marge gauche alignée sur le placeholder titre (0.615 in),
# marge droite plafonnée avant le badge de pagination bas-droit.
SLIDE_W, SLIDE_H = 10.0, 5.625
MARGIN = 0.615
BORD_DROIT = 9.15
CONTENT_TOP = 1.15
CONTENT_BOTTOM = 5.45
CONTENT_W = BORD_DROIT - MARGIN
CONTENT_H = CONTENT_BOTTOM - CONTENT_TOP
GAP = 0.2

def _exiger_template():
    """Garde à l'import (finding robustesse, audit 2026-07-23) : sans elle, un
    template absent remontait en FileNotFoundError brute depuis python-pptx. Le
    générateur est lancé à la main — l'échec doit nommer le fichier attendu et
    son emplacement, sans exiger de lire la stack."""
    if not os.path.isfile(TEMPLATE):
        raise SystemExit(
            "generate_deck : template introuvable — placer template-octo.pptx "
            f"à côté du générateur (attendu : {TEMPLATE})"
        )


_exiger_template()
TH = D.theme_colors(Presentation(TEMPLATE))
NAVY = TH.get("dk1", D.INK)          # #0E2356 — texte principal, titres
WHITE = TH.get("lt1", "#FFFFFF")
ACCENT = TH.get("accent3", D.PALETTE[0])   # #00D2DD — cyan OCTO, identité du deck
MUTED = TH.get("lt2", D.MUTED)       # #586586 — slate 600, texte secondaire
LINE = TH.get("accent5", D.LINE)     # #CFD3DD — slate 200, bordures de cards
TRACK = TH.get("accent6", D.TRACK)   # #E7E9EE — slate 100, fonds d'encarts

SEVERITE = ["#1e6b34", "#5b8a3c", "#b8860b", "#c1651e", "#b3261e"]  # D0..D4, vert -> rouge


def _rgb(hexcolor):
    return RGBColor.from_string(hexcolor.lstrip("#").upper())


def new_prs():
    _exiger_template()
    prs = Presentation(TEMPLATE)
    # Retire les 9 slides d'exemple du template — masters/layouts/thème conservés.
    # Il faut aussi supprimer la relation (drop_rel), sinon les parties
    # ppt/slides/slideN.xml orphelines entrent en collision de nom avec les
    # nouvelles slides ajoutées ensuite (même numérotation réutilisée).
    xml_slides = prs.slides._sldIdLst
    for sld in list(xml_slides):
        rId = sld.get(D.qn("r:id"))
        prs.part.drop_rel(rId)
        xml_slides.remove(sld)
    return prs


# Couleur de chapitre par groupe de slides — signal de navigation plus fort
# que le seul kicker textuel (piste retenue dans analyse-template-alternatif.md).
# Un code couleur par chapitre, passé explicitement à chaque appel de
# content_slide() ET repris sur l'intercalaire du chapitre (v2.9, 9 chapitres) :
#   01 Exec summary        = NAVY           (le pitch de l'offre, en ouverture)
#   02 Contexte            = D.PALETTE[0]  (bleu)
#   03 Personas            = D.PALETTE[5]  (teal)
#   04 Besoins & douleurs  = D.PALETTE[2]  (rouge)
#   05 Proposition         = D.PALETTE[1]  (vert)
#   06 IA                  = D.PALETTE[4]  (violet)
#   07 Démarche            = D.PALETTE[3]  (or)
#   08 Outillage IAP       = D.PALETTE[5]  (teal — réutilisé, comme KPI réutilise le bleu)
#   09 KPI                 = D.PALETTE[0]  (bleu)


def content_slide(prs, kicker, title, color):
    # v2.13 : `color` etait optionnel (repli sur ACCENT, cyan generique) —
    # diagnostic du 2026-09-03 : 34/34 appels reels de ce module passent deja
    # `color=` explicitement, le repli ne protegeait plus rien, il masquait
    # silencieusement un oubli (cf. le defaut kicker de slide_executive_summary,
    # v2.12). Obligatoire desormais : un oubli echoue au build, pas au rendu.
    layout = prs.slide_masters[0].slide_layouts[LAYOUT_TITRE_SEUL]
    s = prs.slides.add_slide(layout)
    ph = s.shapes.placeholders[0]
    box_w = Emu(ph.width).inches
    kicker_color = color
    texte_complet = f"{kicker.upper()}   ·   {title}"
    taille, _ = D.ajuster_police([texte_complet], box_w, 17, 12,
                                  lambda t, lignes_max: lignes_max <= 1)
    tf = ph.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = kicker.upper() + "   ·   "
    r1.font.bold = True
    r1.font.size = Pt(taille)
    r1.font.color.rgb = _rgb(kicker_color)
    r2 = p.add_run()
    r2.text = title
    r2.font.bold = True
    r2.font.size = Pt(taille)
    r2.font.color.rgb = _rgb(NAVY)
    return s


def slide_sous_chapitre(prs, kicker, titre, sous_titre, color):
    """Séparateur de SOUS-chapitre (léger) : PAS l'intercalaire teardrop des
    chapitres (layout dédié + photo), juste un bloc-titre de section sur le layout
    « titre seul ». Introduit un groupe logique DANS un chapitre. Reste plus léger
    qu'un chapitre : pas de numéro, pas de photo, garde le pied de page du master.
    SANS APPELANT depuis la v2.6 : le seul groupe qui l'utilisait — « Exemples »
    dans la Proposition — a été supprimé à la demande (l'arbitrage 2026-07-22
    « un séparateur pour les Exemples » est caduc, les 3 slides d'exemple vivent
    dans git, v2.5). Machinerie des séparateurs à deux niveaux conservée (cf.
    CLAUDE.md §docs/cadrage-ppt) pour un futur groupe logique."""
    layout = prs.slide_masters[0].slide_layouts[LAYOUT_TITRE_SEUL]
    s = prs.slides.add_slide(layout)
    # Vider le placeholder titre (sinon prompt résiduel) — on pose notre propre bloc.
    s.shapes.placeholders[0].text_frame.text = ""
    bar_top, bar_h = 2.05, 1.55
    D.add_rect(s, MARGIN, bar_top, 0.14, bar_h, fill=color, rounded=True, radius=0.5)
    tx = MARGIN + 0.45
    tw = CONTENT_W - 0.45
    D.add_text(s, tx, bar_top, tw, bar_h, [
        (kicker.upper() + "  ·  SOUS-CHAPITRE",
         dict(size=D.TYPE["tiny"], bold=True, color=color, line_spacing=1.0)),
        (titre, dict(size=34, bold=True, color=NAVY, space_before=8, line_spacing=1.0)),
        (sous_titre, dict(size=D.TYPE["small"], color=MUTED, italic=True, space_before=12, line_spacing=1.2)),
    ], anchor=MSO_ANCHOR.MIDDLE)
    return s


def _sans_puce(paragraph):
    """Retire l'indentation de puce héritée (marL/indent) et désactive le
    caractère de puce. Cause réelle du bug "01" qui passe à la ligne dans le
    petit encart numéro du layout Chapitre : son style de liste hérité pose
    marL=0.5in (indentation de puce) dans un encart large de 0.546in — il ne
    reste presque plus de largeur utile, donc chaque caractère wrap. Le REX
    V3 (VSCode1) corrige exactement ça avec un pPr marL=0/indent=0/buNone
    explicite ; python-pptx n'expose pas ces attributs, d'où la manipulation
    XML directe."""
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("marL", "0")
    pPr.set("indent", "0")
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def _find_frame_by_geom(shapes, prst):
    """Cadre non groupé (top-level) portant un prstGeom donné — variante de
    pptx-framed-image.frame_geometry pour le cas où le cadre n'est pas niché
    dans un groupe (le layout Chapitre du template, à la différence des
    layouts « cadre blanc », place son cadre teardrop directement)."""
    for sh in shapes:
        spPr = getattr(sh._element, "spPr", None)
        if spPr is None:
            continue
        g = spPr.find(qn("a:prstGeom"))
        if g is not None and g.get("prst") == prst:
            return sh.left, sh.top, sh.width, sh.height, g
    return None


def _find_frame_in_group(shapes, group_name, inner_name):
    from framed_image import frame_geometry
    for sh in shapes:
        if sh.name == group_name:
            return frame_geometry(sh, inner_name)
    return None


# scene -> requête Openverse (photo réelle) ; la génération procédurale
# (nature_images) reste le nom de "scene" utilisé comme repli hors-ligne.
_REQUETES_PHOTO = {
    "mountains": "mountains landscape",
    "forest": "green forest sunlight",
    # Littoral rocheux turquoise (chapitre Besoins & douleurs) : « ocean waves aerial » (horizon
    # brumeux délavé en blanc) puis « turquoise sea water aerial » (0 résultat
    # Openverse -> repli procédural à ciel pâle) échouaient tous deux à ancrer le
    # haut du cadre teardrop sur le fond blanc de la slide. « turquoise water »
    # (seed 0) renvoie une vue plongeante roche+eau+écume, texturée et contrastée
    # sur les quatre bords — VÉRIFIÉE au rendu réel le 2026-07-21.
    "ocean": "turquoise water",
    "sunset": "sunset sky",
    # Chapitres à photo (restructurations 7 puis 8 puis 9 chapitres) : scènes réelles
    # distinctes, VÉRIFIÉES au rendu réel — une requête mot-clé n'a aucun jugement (cf. « plage
    # bondée », « desert dune » seed 0 → fossile de musée, « winding river » →
    # cloître de monastère), donc chaque photo est validée à l'œil (fetch du _brut
    # puis lecture image avant câblage). Le repli nature_images (procédural) ne se
    # déclenche que si Openverse est indisponible (SSL/0-résultat) ET que le nom de
    # scène est connu du fallback (forest/meadow/mountains/ocean/sunset/tropical) —
    # sinon le générateur PLANTE (ValueError unknown scene). Préférer une vraie photo
    # à du procédural ; cf. mémoire reference-deck-image-fetcher.
    #   dunes  (Proposition) = vue aérienne NASA ; nightsky (IA) = astrophoto ;
    #   canyon (Démarche)    = strates de roche (nom NEUF → repli qui PLANTE, comme
    #                          dunes/nightsky : dépend d'un vrai fetch) ;
    #   meadow (KPI, seed 1) = asters/verges d'or (nom CONNU du fallback → sûr).
    "dunes": "sand dunes",
    "nightsky": "starry night sky",
    "canyon": "canyon landscape",
    "meadow": "meadow wildflowers",
    # tropical (Outillage IAP, chapitre 08 — nommé chapitre 07 en v2.5) : nom CONNU
    # du fallback procédural (forest/meadow/mountains/ocean/sunset/tropical), donc
    # sûr même hors ligne. Photo à VÉRIFIER au rendu réel comme les autres.
    "tropical": "tropical palm leaves",
    # wheatfield (Exec summary, chapitre 01, nouveau v2.8) : épis de blé doré, gros
    # plan texturé — la récolte/le résultat, en écho au thème du chapitre (l'offre
    # ET sa synthèse, « ce que la mission produit »). « golden wheat field sunset »/
    # « wheat field golden hour » (0 résultat Openverse en aspect carré, le cadre
    # teardrop de ce layout est carré — pas « tall » comme les autres chapitres)
    # échouaient ; « wheat field » simple RENVOIE un résultat, gros plan contrasté
    # sur les 4 bords — VÉRIFIÉE au rendu réel le 2026-09-02. Nom NEUF → repli
    # procédural qui PLANTE sauf mapping _SCENE_REPLI (ci-dessous).
    "wheatfield": "wheat field",
}


# Repli procedural : nature_images ne connait que 6 scenes
# (forest/meadow/mountains/ocean/sunset/tropical). Les noms NEUFS choisis pour
# les chapitres (dunes, nightsky, canyon) n'y sont pas — hors reseau ou sur
# 0-resultat Openverse, generate_to levait ValueError HORS du try, ce qui tuait
# build() en entier : 0 slide produite alors que 37 des 40 n'ont pas de photo
# (mesure du 2026-09-01). On mappe donc chaque nom neuf sur la scene connue la
# plus proche visuellement. Ce n'est PAS la meme image — c'est un repli assume,
# dont le but est que le deck sorte, pas qu'il soit identique.
_SCENE_REPLI = {
    "dunes": "sunset",       # tons chauds sable/orange
    "nightsky": "sunset",    # composition de ciel (clair au lieu de sombre)
    "canyon": "mountains",   # relief rocheux
    "wheatfield": "meadow",  # champ ouvert, tons chauds proches
}

# Anomalies relevees pendant le build (pas seulement d'image, malgre le nom
# historique), fusionnees dans `problemes` par build(). Sans cela, un defaut
# ne sortait qu'en print : le build annoncait « GEOMETRIE: OK » avec des
# photos manquantes (cadre introuvable/repli impossible) OU, depuis v2.11,
# avec un contenu qui deborde silencieusement sous un element de pied de
# carte (cf. slide_pitch_iap : le self-check geometrique de pptx_deck.py ne
# mesure QUE les formes que NOUS dessinons hors-cadre — pas un debordement de
# texte dans son propre panneau).
_ANOMALIES_BUILD = []


def _remplir_cadre(slide, cadre, scene, seed=0):
    """Pose une vraie photo libre de droit (Openverse, CC0) à l'aspect exact
    du cadre, repli sur la génération procédurale (nature_images) si le
    réseau/l'API n'est pas disponible — cf. pptx-framed-image, greffé depuis
    VSCode1. Une photo réelle lit mieux qu'un aplat vectoriel généré, constat
    fait en comparant au REX "⛱️ L'Été de l'IA" (VSCode1) qui utilise de
    vraies photos sur ces mêmes cadres."""
    if cadre is None:
        msg = f"cadre introuvable pour la scène '{scene}' — image non posée"
        print(f"  {msg}")
        _ANOMALIES_BUILD.append(msg)
        return
    left, top, width, height, geom = cadre
    aspect = Emu(width).inches / Emu(height).inches
    px_w = 960
    px_h = int(round(px_w / aspect))
    path = os.path.join(IMG_DIR, f"{scene}_{seed}_{px_w}x{px_h}.png")
    if not os.path.exists(path):
        requete = _REQUETES_PHOTO.get(scene, scene)
        aspect_ratio = "wide" if aspect > 1.15 else "tall" if aspect < 0.85 else "square"
        try:
            brut = os.path.join(IMG_DIR, f"_brut_{scene}_{seed}.jpg")
            stock_images.fetch_to(brut, requete, seed=seed, aspect_ratio=aspect_ratio,
                                   manifest_path=IMG_MANIFEST)
            from framed_image import cover_crop_to_aspect
            cover_crop_to_aspect(brut, path, aspect)
            print(f"  photo réelle posée pour '{scene}' ({requete!r}, via Openverse CC0)")
        except Exception as e:
            repli = _SCENE_REPLI.get(scene, scene)
            note = f" (scène '{scene}' inconnue du repli -> '{repli}')" if repli != scene else ""
            print(f"  Openverse indisponible pour '{scene}' ({e}) — repli sur nature_images{note}")
            try:
                nature_images.generate_to(path, repli, px_w, px_h, seed=seed)
            except Exception as e2:
                # Degrader, jamais planter : la slide sort sans photo et le
                # defaut remonte dans `problemes`, il ne disparait pas.
                msg = f"aucune image pour '{scene}' : Openverse KO ({e}) et repli KO ({e2})"
                print(f"  {msg}")
                _ANOMALIES_BUILD.append(msg)
                return
    place_image_in_frame(slide, path, left, top, width, height, geom=geom)


def slide_chapitre(prs, numero, titre, couverture, color, scene, seed=0):
    """Slide d'intercalaire de chapitre — vrai layout dédié du template
    (« 50 - Chapitre [1] »), repris tel qu'utilisé dans le REX
    "⛱️ L'Été de l'IA" (VSCode1) : cadre photo teardrop rempli (pas laissé
    vide), numéro à 17pt (pas la taille par défaut d'un texte de titre — un
    premier essai à 28pt débordait du petit encart sur le badge logo voisin,
    trouvé au rendu, cf. mémoire de session). Couleur de chapitre appliquée
    au numéro et au titre — le REX source ne le faisait pas, ajouté ici pour
    rester cohérent avec le code couleur déjà en place sur tout le deck."""
    layout = prs.slide_masters[0].slide_layouts[LAYOUT_CHAPITRE]
    s = prs.slides.add_slide(layout)
    phs = {ph.placeholder_format.idx: ph for ph in s.placeholders}

    phs[0].text_frame.text = titre
    p2 = phs[0].text_frame.add_paragraph()
    p2.text = couverture
    for r in p2.runs:
        r.font.size = Pt(D.TYPE["small"])
        r.font.italic = True
        r.font.color.rgb = _rgb(MUTED)
    for p in phs[0].text_frame.paragraphs[:1]:
        for r in p.runs:
            r.font.color.rgb = _rgb(color)

    tf1 = phs[1].text_frame
    tf1.text = numero
    # Marges par défaut (~0.1in/côté) mangent la largeur du minuscule encart
    # (0.55in) et forcent "01" à passer à la ligne — invisible tant qu'on ne
    # zéroute pas les marges comme le fait l'exemple qui fonctionne (REX V3).
    tf1.margin_left = tf1.margin_right = tf1.margin_top = tf1.margin_bottom = 0
    tf1.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in tf1.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        _sans_puce(p)
        for r in p.runs:
            r.font.size = Pt(17)
            r.font.color.rgb = _rgb(color)

    cadre = _find_frame_by_geom(s.slide_layout.shapes, "teardrop")
    for pb in frame_obstructions(s, *cadre[:4]) if cadre else []:
        print(f"  [obstruction] chapitre {numero}:", pb["source"], pb["name"], pb["reason"])
    _remplir_cadre(s, cadre, scene, seed)
    return s


def dot_scale(slide, x, y, n, score, color, d=0.14, gap=0.06, empty_color=None):
    """Jauge à points 0..n (score plein en `color`, reste en `empty_color`) —
    pattern repris de la carte de recommandation valeur/complexité observée
    dans l'autre template analysé (analyse-template-alternatif.md §4)."""
    empty_color = empty_color or TRACK
    for i in range(n):
        fill = color if i < score else empty_color
        D.add_dot(slide, x + i * (d + gap), y, d, fill)


def col_x(i, n, w=CONTENT_W, x0=MARGIN, gap=GAP):
    col_w = (w - (n - 1) * gap) / n
    return x0 + i * (col_w + gap), col_w


def chip(slide, x, y, w, h, label, color, text_color="#ffffff", size=D.TYPE["tiny"]):
    D.add_rect(slide, x, y, w, h, fill=color, rounded=True, radius=0.5)
    D.add_text(slide, x, y, w, h, [(label, dict(size=size, bold=True, color=text_color,
                align=PP_ALIGN.CENTER))], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)


# --- Helpers du schéma « parcours de mission » (slide_offre_iap, v2.8) : pas de
# CONNECTOR/oval réutilisable ailleurs dans le générateur avant ce schéma, donc
# petits helpers dédiés plutôt qu'un détour par pptx_deck (déjà surchargé de
# add_rect/add_card génériques — ceux-ci sont spécifiques à ce diagramme).
def _oval(slide, x, y, w, h, fill=None, line=None, line_w=1.0):
    """Ellipse simple (nœuds « entrée/sortie » du schéma de parcours)."""
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.shadow.inherit = False
    except Exception:
        pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(fill)
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = _rgb(line)
        shp.line.width = Pt(line_w)
    shp.text_frame.paragraphs[0].text = ""
    return shp


def _dashed_rect(slide, x, y, w, h, fill, line, line_w=1.0, radius=0.12):
    """Rectangle à bordure pointillée (« mécanisme additif » du schéma de parcours) —
    python-pptx n'expose le style de trait qu'en LineFormat.dash_style, pas via
    D.add_rect (qui ne prend pas ce paramètre)."""
    shp = D.add_rect(slide, x, y, w, h, fill=fill, line=line, line_w=line_w,
                      rounded=True, radius=radius)
    shp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return shp


def _noeud_socle(slide, x, y, w, h, titre, sous_titre=None, oval=False):
    """Nœud « mouvement du socle » (toujours présent) du schéma de parcours —
    fill bleu-gris clair, bordure navy ; ellipse pour les nœuds d'entrée/sortie."""
    fill = "#dce6f5"
    if oval:
        _oval(slide, x, y, w, h, fill=fill, line=NAVY, line_w=1.0)
    else:
        D.add_rect(slide, x, y, w, h, fill=fill, line=NAVY, line_w=1.0, rounded=True, radius=0.14)
    lignes = [(titre, dict(size=7, bold=True, color=NAVY, align=PP_ALIGN.CENTER, line_spacing=1.0))]
    if sous_titre:
        lignes.append((sous_titre, dict(size=5.8, color=MUTED, align=PP_ALIGN.CENTER,
                                         italic=True, space_before=1, line_spacing=1.0)))
    D.add_text(slide, x + 0.04, y, w - 0.08, h, lignes, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)


def _pilule_variante(slide, x, y, w, h, texte, size=6.2):
    """Pilule « variante conditionnée au contexte » (sable/or) du schéma de parcours.
    Si `h` est None, la hauteur est calculée à partir du texte (pilules « si contexte
    politique », plus longues que les pilules courtes « Contexte léger/politique ») —
    retourne toujours la hauteur effectivement utilisée."""
    pad = 0.03
    if h is None:
        lignes = _lignes(texte, w - 2 * pad, size)
        h = 2 * pad + lignes * (size * 1.15 / 72.0)
    D.add_rect(slide, x, y, w, h, fill="#f7ecd2", line=D.PALETTE[3], line_w=1.0,
               rounded=True, radius=0.35)
    D.add_text(slide, x + 0.05, y, w - 0.10, h, [
        (texte, dict(size=size, bold=True, color=D.PALETTE[3], align=PP_ALIGN.CENTER, line_spacing=1.05)),
    ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    return h


def _note_mecanisme(slide, x, y, w, titre, corps, title_size=6.4, body_size=6.0, pad=0.04):
    """Encadré pointillé pâle = « mécanisme additif » (extension, checklist transverse)
    du schéma de parcours — hauteur calculée à partir du corps, jamais fixe (cf. défaut
    « panneau sur-étiré » du dépôt) ; retourne la hauteur effectivement utilisée."""
    lignes = _lignes(corps, w - 2 * pad, body_size)
    h = 2 * pad + (title_size * 1.1 / 72.0) + 0.02 + lignes * (body_size * 1.15 / 72.0)
    _dashed_rect(slide, x, y, w, h, fill="#fdf8ec", line=D.PALETTE[3], line_w=0.9, radius=0.10)
    D.add_text(slide, x + pad, y + pad * 0.6, w - 2 * pad, h - pad * 1.2, [
        (titre, dict(size=title_size, bold=True, color=D.PALETTE[3], line_spacing=1.05)),
        (corps, dict(size=body_size, color=MUTED, italic=True, space_before=2, line_spacing=1.15)),
    ])
    return h


def _fleche_h(slide, x, y, w, h, color=MUTED, size=10):
    """Flèche « → » centrée dans une cellule (vocabulaire de flux du schéma de
    parcours — même simplification texte que slide_iap_contexte_client)."""
    D.add_text(slide, x, y, w, h, [
        ("→", dict(size=size, bold=True, color=color, align=PP_ALIGN.CENTER)),
    ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)


# --- Badge de série (v2.6, point ④) : les 4 slides « proposition de déploiement
# agentic chez le client » du chapitre 06 · IA (3 agents candidats + export
# markdown) portent le MÊME petit badge — signal visuel récurrent et discret qui
# les relie à la zone « déploiement agentic » du schéma d'architecture
# (slide_iap_contexte_client, chapitre 08). Renvoi par CHAPITRE, jamais par
# numéro de page (les numéros bougent). Violet D.PALETTE[4] = couleur du
# chapitre IA, la même que la zone du schéma.
BADGE_AGENTIC_W = 2.3


def badge_deploiement_agentic(slide):
    x = BORD_DROIT - BADGE_AGENTIC_W
    h = 0.42
    D.add_rect(slide, x, CONTENT_TOP, BADGE_AGENTIC_W, h, fill="#ffffff",
               line=D.PALETTE[4], line_w=1.0, rounded=True, radius=0.18)
    D.add_text(slide, x + 0.12, CONTENT_TOP, BADGE_AGENTIC_W - 0.24, h, [
        ("DÉPLOIEMENT AGENTIC CHEZ LE CLIENT",
         dict(size=6, bold=True, color=D.PALETTE[4], line_spacing=1.1)),
        ("cf. schéma d'architecture · chapitre 08",
         dict(size=6, italic=True, color=MUTED, space_before=1)),
    ], anchor=MSO_ANCHOR.MIDDLE)


# Le glyphe "⟲" (U+27F2) n'a pas de variante GRASSE dans la police du template
# (rendu LibreOffice = case vide/tofu dans un run bold) alors que sa variante
# normale s'affiche — même correctif que slide_trajectoire/slide_schema_*
# /slide_livrables_ppt : forcer bold=False pour ce SEUL caractère. Voir
# CLAUDE.md §docs/cadrage-ppt.
_GLYPHES_SANS_GRAS = ("⟲",)


def _header_cell(slide, x, y, w, h, label, size=7, color=MUTED, bold=True,
                 anchor=MSO_ANCHOR.TOP):
    """En-tête de colonne en un seul paragraphe multi-runs : chaque caractère de
    `_GLYPHES_SANS_GRAS` est posé en bold=False même si le libellé est en gras,
    pour éviter le tofu du "⟲" en fonte grasse (cf. _GLYPHES_SANS_GRAS)."""
    import re as _re
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    p = tf.paragraphs[0]
    motif = "(" + "|".join(_re.escape(g) for g in _GLYPHES_SANS_GRAS) + ")"
    for part in _re.split(motif, label):
        if not part:
            continue
        r = p.add_run()
        r.text = part
        r.font.size = Pt(size)
        r.font.bold = bool(bold) and part not in _GLYPHES_SANS_GRAS
        r.font.color.rgb = _rgb(color)
    return box


def _lignes(texte, largeur_in, taille_pt):
    """Nombre de lignes estimé pour `texte` (helper de dimensionnement des
    panneaux à la hauteur de leur contenu — cf. « panneau sur-étiré »)."""
    return max(1, D.estimer_lignes(texte, largeur_in, taille_pt))


# ---------------------------------------------------------------- slide 1
def slide_cover(prs):
    layout = prs.slide_masters[0].slide_layouts[LAYOUT_COUVERTURE]
    s = prs.slides.add_slide(layout)
    phs = {ph.placeholder_format.idx: ph for ph in s.placeholders}
    phs[0].text_frame.text = "BMAD IAP"
    phs[1].text_frame.text = "Infra as a Product Transformation Pack — synthèse de cadrage"
    phs[2].text_frame.text = "OCTO Technology"
    # v2.8 (2026-09-02) : nouveau chapitre 01 « Exec summary » en ouverture du
    # deck (l'offre du pitch + sa synthèse) — tous les chapitres suivants
    # glissent de +1. (v2.6 : sous-chapitre « Exemples » supprimé (séparateur +
    # 3 slides illustratives, à la demande) ; nouvelles slides « activités
    # humaines avec/sans l'outil » (Démarche) et « architecture IAP en contexte
    # client » (ouvre l'Outillage IAP) ; badge de série « déploiement agentic
    # chez le client » sur les 4 slides de proposition agentic (chapitre IA).
    # v2.5 : restructuration 8 chapitres sur le fil rouge SCALE — fusion
    # trajectoire/bout-en-bout, executive summary réancré, chapitre Outillage
    # IAP. v2.4 : fil humain de la trajectoire.)
    phs[3].text_frame.text = f"{VERSION_DECK} · {DATE_VERSION_DECK}"
    # Bandeau de métadonnées (statut/langue/confidentialité/sources) retiré sur
    # demande — la couverture ne garde que titre, sous-titre, entité et version.
    return s


# ---------------------------------------------------------------- slide 2
# Réancré (v2.5, chantier ②) sur le fil rouge narratif SCALE
# (docs/Import/notes-extraction-scale.md) : 4 blocs POURQUOI → QUOI → COMMENT →
# RÉSULTAT, chacun une claim d'une ligne + le renvoi aux chapitres — pattern 7
# du catalogue deck-design-library (« rangée de cartes sur bandeau, une en
# accent ») : le RÉSULTAT (la preuve, ce que le sponsor achète in fine) est la
# seule carte en fill navy plein.
def slide_executive_summary(prs):
    # v2.12 : seul appel content_slide() de tout le deck SANS color= explicite
    # (les ~30 autres en passent tous un) -> kicker retombait sur ACCENT (cyan
    # generique) au lieu du NAVY du chapitre 01 Exec summary, dont cette slide
    # ouvre le fil (slide_pitch_iap et slide_demarche_avec_sans_agentic, juste
    # apres, portent deja color=NAVY) -- visible dans les DEUX rendus (LibreOffice
    # ET PowerPoint), manquee lors d'une premiere relecture concentree sur le
    # contenu (bloc OFFRE) plutot que sur la couleur du kicker.
    s = content_slide(prs, "Executive summary",
                       "Du pourquoi à la preuve : une transformation cadrée de bout en bout",
                       color=NAVY)

    headline_h = 0.62
    D.add_text(s, MARGIN, CONTENT_TOP, CONTENT_W, headline_h, [
        ("Transformer l'infrastructure en plateforme opérée comme un produit, ET traiter "
         "structurellement le gaspillage qui l'en empêche — le deck suit le fil : l'offre, "
         "pourquoi, qui, quoi, comment, avec quoi, la preuve.",
         dict(size=D.TYPE["small"], color=NAVY, italic=True, line_spacing=1.3)),
    ])

    # v2.10 : bloc OFFRE ajouté en tête de rangée (arbitrage utilisateur) — le
    # sommaire ne référençait aucun chapitre avant "Chapitres 02–04" alors que
    # le chapitre 01 · Exec summary (slide_pitch_iap + slide_demarche_avec_sans_agentic,
    # juste après CETTE slide) restait invisible ici. Accent (fond plein) comme
    # RÉSULTAT : les deux bornent la rangée (l'offre ouvre, la preuve ferme).
    items = [
        ("OFFRE", NAVY, "Trois faces à ne pas confondre, une démarche avec ou sans IA.",
         "Leurs douleurs, notre outillage de consultant, l'agentic chez eux en option — "
         "déclinés sur les trois mêmes temps, outillés ou non.",
         "Chapitre 01"),
        ("POURQUOI", D.PALETTE[0], "L'infra subie coûte de plus en plus cher.",
         "Trois déclencheurs, quatre personas interrogés séparément, des douleurs "
         "mesurables plutôt que des plaintes.",
         "Chapitres 02–04"),
        ("QUOI", D.PALETTE[1], "Traiter l'infra comme un produit — et assainir.",
         "Double mission, méthode scorée (impact × faisabilité − prudence IA), IA sous "
         "gate : jamais la réponse à un problème d'abord organisationnel.",
         "Chapitres 05–06"),
        ("COMMENT", D.PALETTE[3], "Trois temps et une boucle, personnes comprises.",
         "Démarche ①②③⟲ avec son fil humain de bout en bout ; l'outillage IAP au "
         "service de la démarche — jamais l'inverse.",
         "Chapitres 07–08"),
        ("RÉSULTAT", NAVY, "Le delta instrumenté T0 → réévaluation.",
         "Trois familles de KPIs, même instrument aux deux instants — la preuve, "
         "pas une opinion.",
         "Chapitre 09"),
    ]
    n = len(items)
    pad = 0.16
    _, cw = col_x(0, n)
    usable = cw - 2 * pad
    desc_size = 8
    line_h = desc_size * 1.3 / 72.0
    claim_h = max(_lignes(c, usable, 9) for _, _, c, _, _ in items) * (9 * 1.2 / 72.0) + 0.06
    desc_h = max(_lignes(d, usable, desc_size) for _, _, _, d, _ in items) * line_h + 0.06
    # étages : label (0.24) + claim + desc + renvoi chapitres (0.26) + respirations
    card_h = 0.14 + 0.24 + claim_h + 0.10 + desc_h + 0.14 + 0.26 + 0.14
    top0 = CONTENT_TOP + headline_h + 0.30
    # bandeau de fond commun (pattern 7) : regroupe les 5 blocs en un seul
    # « bloc de lecture » — le fil se lit d'un trait, flèches dans les inter-colonnes.
    D.add_rect(s, MARGIN - 0.08, top0 - 0.16, CONTENT_W + 0.16, card_h + 0.32,
               fill=TRACK, rounded=True, radius=0.06)
    for i, (etape, color, claim, desc, renvoi) in enumerate(items):
        x, w = col_x(i, n)
        accent = etape in ("OFFRE", "RÉSULTAT")   # bornent la rangée : l'offre ouvre, la preuve ferme
        if accent:
            D.add_rect(s, x, top0, w, card_h, fill=NAVY, rounded=True, radius=0.08)
        else:
            D.add_rect(s, x, top0, w, card_h, fill="#ffffff", line=LINE, line_w=0.75,
                       rounded=True, radius=0.08)
        D.add_text(s, x + pad, top0 + 0.14, w - 2 * pad, 0.24, [
            (etape, dict(size=8, bold=True, color="#8fd6db" if accent else color)),
        ])
        D.add_text(s, x + pad, top0 + 0.38, w - 2 * pad, claim_h, [
            (claim, dict(size=9, bold=True, color="#ffffff" if accent else NAVY,
                         line_spacing=1.2)),
        ])
        D.add_text(s, x + pad, top0 + 0.38 + claim_h + 0.10, w - 2 * pad, desc_h, [
            (desc, dict(size=desc_size, color="#c7cbe0" if accent else MUTED,
                        line_spacing=1.3)),
        ])
        D.add_text(s, x + pad, top0 + card_h - 0.38, w - 2 * pad, 0.26, [
            (renvoi, dict(size=7.5, bold=True, color="#8fd6db" if accent else
                          (D.PALETTE[0] if etape == "RÉSULTAT" else color))),
        ], anchor=MSO_ANCHOR.BOTTOM)
        if i < n - 1:   # le fil : flèche dans l'inter-colonne
            D.add_text(s, x + w - 0.02, top0 + card_h / 2 - 0.12, GAP + 0.04, 0.24, [
                ("→", dict(size=11, bold=True, color=MUTED, align=PP_ALIGN.CENTER)),
            ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    return s


# ---------------------------------------------------------------- chapitre 07 (v2.9)
def slide_offre_iap(prs):
    """DÉPLACÉE en v2.9 (arbitrage utilisateur) du chapitre 01 · Exec summary vers
    le chapitre 07 · Démarche, juste avant slide_schema_fonctionnement : le
    parcours de mission détaillé est un objet de COMMENT, pas d'ouverture — le
    sommaire du deck reste slide_executive_summary. Contenu inchangé ; seuls le
    kicker et la couleur suivent le chapitre d'accueil (D.PALETTE[3], or).

    Nouveau (v2.8) — reprend la slide 2 du pitch source (chapô + citation-thèse,
    VERBATIM — document source, ne pas reformuler) et son schéma du parcours de
    mission, redessiné EN NATIF (pas une insertion de l'image source) sur ses 3
    registres : mouvements du socle (toujours présents, bleu-gris), variantes
    conditionnées au contexte (sable/or), mécanismes additifs (encadrés pointillés
    pâles). Garde « 6 catégories » au nœud Discovery gaspillages tel quel — même si
    le chapitre Besoins & douleurs du deck compte 8 familles de gaspillage, c'est une
    divergence DU DOCUMENT SOURCE (marqué WIP), remontée pour arbitrage, pas
    corrigée ici (cf. docstring de module)."""
    s = content_slide(prs, "Démarche",
                       "Accompagnement Infra as a Product : transformer une fonction infra en produit interne",
                       color=D.PALETTE[3])

    chapo = ("L'offre proposée est une méthodologie d'accompagnement pour transformer « une "
             "fonction infra ou une plateforme interne » en un véritable produit interne : un "
             "service pensé pour ses utilisateurs, avec un parcours, une proposition de valeur "
             "et des indicateurs de pilotage, plutôt qu'un centre de coûts ou un guichet de "
             "tickets.")
    chapo_h = _lignes(chapo, CONTENT_W, 8.5) * (8.5 * 1.2 / 72.0) + 0.05
    D.add_text(s, MARGIN, CONTENT_TOP, CONTENT_W, chapo_h, [
        (chapo, dict(size=8.5, color=NAVY, italic=True, line_spacing=1.2)),
    ])
    # (budget vertical serré — cf. contrainte de place du brief : chaque gap
    # ci-dessous a été resserré après un premier rendu réel qui montrait la
    # légende partiellement recouverte par le bandeau de citation, empiétement
    # que `verifier_geometrie` ne peut pas voir — seul le rendu le révèle.)

    citation = ("« Réussir une transformation Infra as a Product, ce n'est pas \"mettre des PO "
                "dans l'infra\". C'est concevoir, opérer et faire adopter une plateforme interne "
                "comme un produit, en équilibrant delivery, robustesse du RUN et valeur perçue "
                "par les utilisateurs internes. »")
    cit_pad = 0.10
    cit_usable = CONTENT_W - 2 * cit_pad
    cit_lines = _lignes(citation, cit_usable, 8.5)
    citation_h = cit_pad + (7 * 1.1 / 72.0) + 0.04 + cit_lines * (8.5 * 1.2 / 72.0) + cit_pad

    grid_n = 5
    row_h = 0.38
    schema_top = CONTENT_TOP + chapo_h + 0.06

    # --- Bande du haut : 2 mécanismes additifs (gauche/droite) + variantes (centre) ---
    x0n, w0n = col_x(0, grid_n)
    note_l_w = 2.55
    h_note_l = _note_mecanisme(s, x0n, schema_top, note_l_w, "EXTENSION POSSIBLE",
                                "Reconstitution d'incident avant Cadrage, si crise déclencheuse.")
    note_r_w = 2.85
    note_r_x = BORD_DROIT - note_r_w
    h_note_r = _note_mecanisme(s, note_r_x, schema_top, note_r_w, "EXTENSION POSSIBLE",
                                "Tri contraintes / habitudes entre Diagnostic et Segmentation, "
                                "si sites hétérogènes.")

    x1, w1 = col_x(1, grid_n)
    x2, w2 = col_x(2, grid_n)
    v_cx = (x1 + w1 + x2) / 2.0
    v_w = 1.35
    v_x = v_cx - v_w / 2.0
    D.add_text(s, v_x - 0.25, schema_top, v_w + 0.5, 0.14, [
        ("signal de contexte détecté ?", dict(size=6, italic=True, color=MUTED, align=PP_ALIGN.CENTER)),
    ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    pill_h = 0.17
    pill1_top = schema_top + 0.12
    pill2_top = pill1_top + pill_h + 0.02
    _pilule_variante(s, v_x, pill1_top, v_w, pill_h, "Contexte léger")
    _pilule_variante(s, v_x, pill2_top, v_w, pill_h, "Contexte politique")
    variantes_bottom = pill2_top + pill_h

    band_a_bottom = max(schema_top + h_note_l, schema_top + h_note_r, variantes_bottom)

    # --- Rangée 1 : mouvements du socle ---
    row1_top = band_a_bottom + 0.14
    D.add_text(s, x0n + w0n * 0.15, band_a_bottom + 0.04, 0.5, 0.14, [
        ("↓", dict(size=8, bold=True, color=D.PALETTE[3], align=PP_ALIGN.CENTER)),
    ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    x3, w3 = col_x(3, grid_n)
    D.add_text(s, x3 + w3 * 0.65, band_a_bottom + 0.04, 0.5, 0.14, [
        ("↓", dict(size=8, bold=True, color=D.PALETTE[3], align=PP_ALIGN.CENTER)),
    ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

    socle_row1 = [
        ("Premier contact", None, True),
        ("Cadrage", "note de cadrage", False),
        ("Diagnostic", "base factuelle partagée", False),
        ("Discovery gaspillages", "6 catégories", False),
        ("Segmentation / Product Discovery", None, False),
    ]
    for i, (titre, sous, oval) in enumerate(socle_row1):
        x, w = col_x(i, grid_n)
        _noeud_socle(s, x, row1_top, w, row_h, titre, sous, oval=oval)
        if i < grid_n - 1:
            _fleche_h(s, x + w, row1_top, GAP, row_h)
    row1_bottom = row1_top + row_h

    # --- + Contradictions structurelles (si contexte politique), sous Discovery gaspillages ---
    h_cs = _pilule_variante(s, x3, row1_bottom + 0.05, w3, None,
                             "+ Contradictions structurelles (si contexte politique)", size=6.0)
    cs_bottom = row1_bottom + 0.05 + h_cs

    # --- Connecteur en Z : fin rangée 1 (col 5) -> début rangée 2 (col 1) ---
    x4, w4 = col_x(grid_n - 1, grid_n)
    x0, w0 = col_x(0, grid_n)
    cx4, cx0 = x4 + w4 / 2.0, x0 + w0 / 2.0
    y_elbow = cs_bottom + 0.06
    row2_top = y_elbow + 0.10
    D.add_rect(s, cx4 - 0.01, row1_bottom, 0.02, y_elbow - row1_bottom, fill=LINE)
    D.add_rect(s, cx0 - 0.01, y_elbow, cx4 - cx0 + 0.02, 0.02, fill=LINE)
    D.add_rect(s, cx0 - 0.01, y_elbow, 0.02, row2_top - y_elbow, fill=LINE)
    D.add_text(s, cx0 - 0.11, y_elbow - 0.01, 0.22, row2_top - y_elbow + 0.02, [
        ("▾", dict(size=6, bold=True, color=MUTED, align=PP_ALIGN.CENTER)),
    ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

    # --- Rangée 2 : mouvements du socle (suite, « boustrophédon ») ---
    socle_row2 = [
        ("Product Definition", "fiche produit cible", False),
        ("Operating Model", "responsabilités + mesure", False),
        ("Adoption / Pilote", "bilan avant / après", False),
        ("Mission close", None, True),
        ("Retour d'expérience", "vers la bibliothèque partagée", False),
    ]
    for i, (titre, sous, oval) in enumerate(socle_row2):
        x, w = col_x(i, grid_n)
        _noeud_socle(s, x, row2_top, w, row_h, titre, sous, oval=oval)
        if i < grid_n - 1:
            _fleche_h(s, x + w, row2_top, GAP, row_h)
    row2_bottom = row2_top + row_h

    # --- + Dispositif de revue (si contexte politique), sous Adoption / Pilote ---
    x2b, w2b = col_x(2, grid_n)
    h_dr = _pilule_variante(s, x2b, row2_bottom + 0.05, w2b, None,
                             "+ Dispositif de revue (si contexte politique)", size=6.0)
    dr_bottom = row2_bottom + 0.05 + h_dr

    # --- Checklist gouvernance IA (mécanisme transverse, pleine largeur) ---
    checklist_top = dr_bottom + 0.05
    h_checklist = _note_mecanisme(s, MARGIN, checklist_top, CONTENT_W,
                                   "CHECKLIST GOUVERNANCE IA — TRANSVERSE",
                                   "Indépendante des mouvements, mobilisable dès qu'un usage IA "
                                   "est identifié — à tout moment de la mission.")
    checklist_bottom = checklist_top + h_checklist

    # --- Légende (3 registres) ---
    legend_top = checklist_bottom + 0.04
    legend = [
        ("#dce6f5", NAVY, False, "mouvement du socle (toujours présent)"),
        ("#f7ecd2", D.PALETTE[3], False, "variante ou section conditionnée au contexte"),
        ("#fdf8ec", D.PALETTE[3], True, "mécanisme additif (extension, checklist transverse)"),
    ]
    lx = MARGIN
    sw = 0.14
    for fill, line, dashed, label in legend:
        if dashed:
            _dashed_rect(s, lx, legend_top, sw, sw, fill=fill, line=line, line_w=0.9, radius=0.3)
        else:
            D.add_rect(s, lx, legend_top, sw, sw, fill=fill, line=line, line_w=1.0, rounded=True, radius=0.3)
        tw = 2.55
        D.add_text(s, lx + sw + 0.06, legend_top - 0.02, tw, sw + 0.05, [
            (label, dict(size=6.3, color=MUTED)),
        ], anchor=MSO_ANCHOR.MIDDLE)
        lx += sw + 0.06 + tw + 0.12

    # --- Citation-thèse (verbatim), bandeau bas ---
    cit_top = CONTENT_BOTTOM - citation_h
    D.add_rect(s, MARGIN, cit_top, CONTENT_W, citation_h, fill=NAVY, rounded=True, radius=0.08)
    D.add_rect(s, MARGIN, cit_top, 0.07, citation_h, fill=ACCENT, rounded=True, radius=0.5)
    D.add_text(s, MARGIN + 0.24, cit_top + cit_pad * 0.5, CONTENT_W - 0.44, citation_h - cit_pad, [
        ("LA THÈSE", dict(size=7, bold=True, color="#8fd6db")),
        (citation, dict(size=8.5, bold=True, color="#ffffff", space_before=3, line_spacing=1.2)),
    ])
    return s


# --- Pictogrammes de la v2.9 : vocabulaire de SILHOUETTES (pas de glyphes
# exotiques, qui rendent en tofu dans la police du template — cf.
# _GLYPHES_SANS_GRAS). Trois formes distinctes, réutilisées à l'identique par
# slide_pitch_iap et slide_demarche_avec_sans_agentic pour que les deux faces de
# l'agentic se reconnaissent d'une slide à l'autre :
#   "alerte"  (ovale + « ! », rouge)   = la douleur du client ;
#   "engrenage" (GEAR_6, teal)         = le module qui outille LE CONSULTANT ;
#   "deploiement" (PENTAGON, violet)   = l'agentic DÉPLOYÉ CHEZ LE CLIENT.
# Deux teintes et deux silhouettes : la confusion « accélérateur agentic » lu
# deux fois de suite (l'outil du consultant vs l'option sous gate IA) rendrait
# la slide FAUSSE — c'est la condition posée en table ronde.
def _picto(slide, kind, x, y, d, color):
    if kind == "alerte":
        _oval(slide, x, y, d, d, fill=color)
        D.add_text(slide, x, y, d, d, [
            ("!", dict(size=11, bold=True, color="#ffffff", align=PP_ALIGN.CENTER)),
        ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        return
    forme = {"engrenage": MSO_SHAPE.GEAR_6, "deploiement": MSO_SHAPE.PENTAGON}[kind]
    shp = slide.shapes.add_shape(forme, Inches(x), Inches(y), Inches(d), Inches(d))
    try:
        shp.shadow.inherit = False
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color)
    shp.line.fill.background()
    shp.text_frame.paragraphs[0].text = ""


# ---------------------------------------------------------------- chapitre 01 (v2.9)
# Nouveau (v2.9, 2026-09-02) : ouvre le chapitre Exec summary à la place de
# slide_offre_iap (partie au chapitre 07 · Démarche). Trois cartes ÉGALES en
# rangée — le client, le consultant, l'option chez le client — puis, HORS de la
# rangée, une bande pleine largeur collée sous les cartes (fond gris, texte plus
# petit) : un poids visuel moindre sans casser l'alignement d'une rangée censée
# être régulière.
# CONTRAINTE FACTUELLE (garde-fou, vérifiée en source) : aucun REX n'existe dans
# les sources. docs/bmad-iap-cadrage.md:115 dit que la source brute CONTENAIT des
# noms de clients réels (secteurs télécom/plateformes numériques, banque de
# détail, GIE informatique bancaire) — c'est une note de rédaction, pas un acte
# de mission ; `rex-library.md` est planifié, non peuplé. La bande parle donc du
# MATÉRIAU de cadrage, jamais d'une intervention : pas de « missions menées », pas
# de verbe d'action au passé, pas de chiffre de résultat, pas de durée, pas
# d'avant/après, pas de nom de client. Elle n'exhibe pas non plus l'incertitude
# interne (le manque de REX vit ailleurs dans le deck) — écrire notre propre
# objection à la place du prospect serait un contresens de slide d'ouverture.
def slide_pitch_iap(prs):
    s = content_slide(prs, "Exec summary",
                       "Trois faces à ne pas confondre : leurs douleurs, notre outillage, l'agentic chez eux",
                       color=NAVY)

    chapo = ("L'offre part des douleurs de ces organisations ; un module agentic accélère le "
             "consultant qui les traite ; déployer de l'agentic chez le client reste une option, "
             "sous gate IA.")
    chapo_h = _lignes(chapo, CONTENT_W, 8.5) * (8.5 * 1.25 / 72.0) + 0.06
    D.add_text(s, MARGIN, CONTENT_TOP, CONTENT_W, chapo_h, [
        (chapo, dict(size=8.5, color=NAVY, italic=True, line_spacing=1.25)),
    ])

    cartes = [
        ("CE QU'ILS VIVENT", D.PALETTE[2], "alerte",
         "Douleurs & besoins de ces organisations",
         "Quatre personas interrogés séparément, huit familles de gaspillage.",
         ["Infra & RUN — incidents subis en boucle, BUILD sacrifié à l'astreinte",
          "Utilisateur applicatif — sans self-service : guichet, contournements",
          "Management — expert devenu manager malgré lui, reporting miroir",
          "Sponsor — pression à « mettre de l'IA » sans cas d'usage démontré"],
         "Détaillé aux chapitres 03 et 04"),
        ("NOTRE OUTILLAGE · CÔTÉ CONSULTANT", D.PALETTE[5], "engrenage",
         "Notre démarche, outillée par un module agentic",
         "Le module BMAD IAP tourne sur le poste du consultant.",
         ["11 agents spécialisés : stratégie, produit, RUN, gaspillages, métriques, IA, changement",
          "Workflows, templates et checklists outillent la démarche — ils ne la remplacent pas",
          "Rien ne s'installe chez le client par défaut : les livrables sortent, c'est tout"],
         "L'outil du consultant — chapitres 07 et 08"),
        ("EN OPTION · CÔTÉ CLIENT", D.PALETTE[4], "deploiement",
         "Déployer de l'agentic chez le client",
         "Ouverte seulement si le contexte du client s'y prête.",
         ["Gate IA d'abord : la donnée est classée (D0–D4) avant tout usage d'IA",
          "Agentic Readiness [0]–[1] oriente vers documentation-first ; [2]–[3] ouvre agentic-implementation",
          "Le processus doit être explicite avant l'agent, jamais l'inverse"],
         "Une décision de mission — chapitre 06"),
    ]

    band_h = 0.56
    band_top = CONTENT_BOTTOM - band_h
    cards_top = CONTENT_TOP + chapo_h + 0.14
    card_h = band_top - cards_top          # la bande est COLLÉE sous les cartes
    n = len(cartes)
    pad = 0.15
    _, cw = col_x(0, n)
    usable = cw - 2 * pad
    picto_d = 0.30
    chip_h = 0.34
    ITEM_SIZE_DEFAUT = 7
    ITEM_SIZE_PLANCHER = 6.25   # jamais en dessous : déjà la taille du kicker/chip de cette slide

    # Étages communs aux 3 cartes (hauteurs dérivées du contenu le plus long) —
    # les titres, accroches et pieds s'alignent d'une carte à l'autre, sinon la
    # rangée « égale » ne l'est qu'en boîte, pas à la lecture.
    titre_h = max(_lignes(c[3], usable, 9.5) for c in cartes) * (9.5 * 1.2 / 72.0) + 0.04
    accr_h = max(_lignes(c[4], usable, 7.5) for c in cartes) * (7.5 * 1.2 / 72.0) + 0.04

    for i, (label, color, picto, titre, accroche, items, pied) in enumerate(cartes):
        x, _ = col_x(i, n)
        D.add_rect(s, x, cards_top, cw, card_h, fill="#ffffff", line=LINE, line_w=0.75,
                   rounded=True, radius=0.08)
        D.add_rect(s, x, cards_top, cw, 0.075, fill=color, rounded=True, radius=0.5)

        y = cards_top + 0.16
        _picto(s, picto, x + pad, y, picto_d, color)
        D.add_text(s, x + pad + picto_d + 0.10, y, usable - picto_d - 0.10, picto_d, [
            (label, dict(size=6.5, bold=True, color=color, line_spacing=1.1)),
        ], anchor=MSO_ANCHOR.MIDDLE)

        y += picto_d + 0.10
        D.add_text(s, x + pad, y, usable, titre_h, [
            (titre, dict(size=9.5, bold=True, color=NAVY, line_spacing=1.2)),
        ])
        y += titre_h + 0.04
        D.add_text(s, x + pad, y, usable, accr_h, [
            (accroche, dict(size=7.5, color=MUTED, italic=True, line_spacing=1.2)),
        ])
        y += accr_h + 0.10
        D.add_rect(s, x + pad, y, usable, 0.012, fill=LINE)
        y += 0.12

        # Le mou restant se répartit ENTRE les puces (jamais en vide sous le
        # pied de carte) : les 3 cartes n'ont pas le même nombre de puces, leurs
        # pieds doivent malgré tout s'aligner. Défaut « panneau flottant » évité
        # par construction, pas par relecture.
        pied_top = cards_top + card_h - 0.14 - chip_h
        dispo = pied_top - 0.12 - y
        # v2.11 : la taille de puce n'était pas adaptée au NOMBRE de puces —
        # une carte à 4 items (ex. "CE QU'ILS VIVENT", 4 personas) débordait
        # silencieusement sous le chip du pied, dont le fond plein masquait le
        # dernier item entier (seule sa puce ronde dépassait, visible comme un
        # petit disque coloré au-dessus du chip — repéré au rendu réel, pas au
        # self-check géométrique qui ne mesurait pas ce débordement). On
        # réduit taille puis interligne, PALIER PAR PALIER, jusqu'à tenir dans
        # `dispo` avec une marge de sécurité — jamais en dessous du plancher.
        marge = 0.03
        item_size = item_ls = items_h = None
        for size in (7, 6.75, 6.5, 6.25):
            for ls in (1.2, 1.1, 1.0):
                lh = size * ls / 72.0
                h = [_lignes(t, usable - 0.14, size) * lh + 0.02 for t in items]
                total = sum(h) + max(0, len(items) - 1) * 0.06
                if total <= dispo - marge:
                    item_size, item_ls, items_h = size, ls, h
                    break
            if item_size:
                break
        if item_size is None:
            # Rien ne tient même au plancher : dégrader plutôt que planter,
            # et le signaler comme une vraie anomalie de build (pas un print
            # perdu) — c'est la clause que "faire évoluer le check graphique"
            # demandait : un débordement de ce type n'est plus silencieux.
            item_size, item_ls = ITEM_SIZE_PLANCHER, 1.0
            lh = item_size * item_ls / 72.0
            items_h = [_lignes(t, usable - 0.14, item_size) * lh + 0.02 for t in items]
            _ANOMALIES_BUILD.append(
                f"slide_pitch_iap carte {i} ('{label}') : {len(items)} puces ne tiennent pas "
                f"dans la hauteur disponible ({dispo:.2f}in) même à la taille plancher "
                f"({ITEM_SIZE_PLANCHER}pt) — raccourcir le texte ou réduire le nombre d'items."
            )
        gap = max(0.06, (dispo - sum(items_h)) / max(1, len(items) - 1))
        for t, ih in zip(items, items_h):
            D.add_rect(s, x + pad + 0.01, y + 0.045, 0.06, 0.06, fill=color, rounded=True, radius=0.5)
            D.add_text(s, x + pad + 0.14, y, usable - 0.14, ih, [
                (t, dict(size=item_size, color=NAVY, line_spacing=item_ls)),
            ])
            y += ih + gap

        chip(s, x + pad, pied_top, usable, chip_h, pied, color, size=6.5)

    # --- Bande « matériau de cadrage » : PAS une 4e carte — pleine largeur,
    # collée sous la rangée, fond gris et texte plus petit (poids visuel moindre).
    D.add_rect(s, MARGIN, band_top, CONTENT_W, band_h, fill=TRACK, rounded=True, radius=0.06)
    D.add_text(s, MARGIN + 0.20, band_top, CONTENT_W - 0.40, band_h, [
        ("MATÉRIAU DE CADRAGE", dict(size=6.5, bold=True, color=MUTED)),
        ("Les personas, les douleurs et les familles de gaspillage de ce cadrage sont nourris "
         "de trois contextes sectoriels réels : télécom et plateformes numériques, banque de "
         "détail, GIE informatique bancaire.",
         dict(size=7.5, color=NAVY, space_before=2, line_spacing=1.2)),
    ], anchor=MSO_ANCHOR.MIDDLE)
    return s


# ---------------------------------------------------------------- chapitre 01 (v2.9)
# Nouveau (v2.9, 2026-09-02) : le niveau ZOOMÉ-ARRIÈRE de l'« avec ou sans
# agentic ». Trois slides du deck traitent déjà le sujet à des altitudes plus
# basses — slide_activites_humaines (phase par phase, chapitre 07),
# slide_iap_contexte_client (topologie poste consultant / poste client, chapitre
# 08), slide_export_markdown (bifurcation documentation-first ou
# agentic-implementation, chapitre 06). Celle-ci n'en redessine aucune : une
# LIGNE HORIZONTALE d'étapes, et sous chaque étape des pastilles empilées — la
# forme dit que c'est la MÊME démarche qui absorbe l'outillage, pas une démarche
# dédoublée. Deux colonnes en vis-à-vis ont été écartées : elles reprendraient
# l'axe poste-consultant / poste-client de slide_iap_contexte_client et
# grilleraient la surprise du chapitre 08, deux slides après le sommaire.
# Langage de couleur : gris-navy = socle (aucun octogone or, qui pillerait le
# vocabulaire des variantes conditionnées du parcours de mission) ; teal =
# module qui outille le consultant ; violet = agentic déployé chez le client
# (même violet que slide_iap_contexte_client et que les badges du chapitre IA).
# v2.13 (2026-09-03, DRAFT non câblé dans build() — généré en isolation pour
# validation avant d'intégrer au reste du deck). Reprend le PATTERN visuel de
# slide_trajectoire (chapitre 07, "vue bout-en-bout") — badge rond numéroté,
# titre, description, chip de renvoi en pied de colonne — appliqué à un
# contenu différent : POURQUOI/QUOI/COMMENT de slide_executive_summary
# (l'ex-slide 3 pour l'utilisateur, une fois son sommaire déplacé après
# l'intercalaire), avec le POURQUOI enrichi de douleurs concrètes reprises de
# la carte "Douleurs & besoins de ces organisations" de slide_pitch_iap
# (l'ex-slide 4) — 3 des 4 bullets condensés en une phrase, pas une 4e
# répétition intégrale. OFFRE et RÉSULTAT du sommaire ne sont PAS repris ici :
# le sommaire garde son rôle de table des matières complète, cette slide-ci
# est une synthèse resserrée sur le fil POURQUOI→QUOI→COMMENT.
def slide_synthese_pourquoi_quoi_comment(prs):
    s = content_slide(prs, "Exec summary",
                       "Le pourquoi, le quoi, le comment — et ce qu'ils vivent concrètement",
                       color=NAVY)

    chapo = ("Le fil du deck en trois temps, le premier ancré dans les douleurs réelles de "
             "ces organisations — pas une généralité, une mesure.")
    chapo_h = _lignes(chapo, CONTENT_W, 8.5) * (8.5 * 1.25 / 72.0) + 0.06
    D.add_text(s, MARGIN, CONTENT_TOP, CONTENT_W, chapo_h, [
        (chapo, dict(size=8.5, color=NAVY, italic=True, line_spacing=1.25)),
    ])

    etapes = [
        ("①", "POURQUOI", D.PALETTE[2],
         "L'infra subie coûte de plus en plus cher.",
         "Mêmes incidents en boucle côté RUN, guichet et contournements côté utilisateurs, "
         "reporting miroir côté management — quatre personas interrogés séparément, huit "
         "familles de gaspillage mesurées, pas des plaintes.",
         "Chapitres 02–04"),
        ("②", "QUOI", D.PALETTE[1],
         "Traiter l'infra comme un produit — et assainir.",
         "Double mission, méthode scorée (impact × faisabilité − prudence IA), IA sous gate : "
         "jamais la réponse à un problème d'abord organisationnel.",
         "Chapitres 05–06"),
        ("③", "COMMENT", D.PALETTE[3],
         "Trois temps et une boucle, personnes comprises.",
         "Démarche ①②③⟲ avec son fil humain de bout en bout ; l'outillage IAP au service de "
         "la démarche — jamais l'inverse.",
         "Chapitres 07–08"),
    ]
    n = len(etapes)
    badge_d = 0.6
    top0 = CONTENT_TOP + chapo_h + 0.3
    line_y = top0 + badge_d / 2 - 0.012
    D.add_rect(s, MARGIN + badge_d / 2, line_y, CONTENT_W - badge_d, 0.024, fill=LINE)
    _, wcol = col_x(0, n)
    label_h = 0.26
    titre_h = max(_lignes(e[3], wcol - 0.2, 10) for e in etapes) * (10 * 1.2 / 72.0) + 0.06
    desc_h = max(_lignes(e[4], wcol - 0.2, 8) for e in etapes) * (8 * 1.25 / 72.0) + 0.05
    for i, (sym, label, color, titre, desc, renvoi) in enumerate(etapes):
        x, w = col_x(i, n)
        cx = x + w / 2 - badge_d / 2
        D.add_rect(s, cx, top0, badge_d, badge_d, fill=color, rounded=True, radius=0.5)
        D.add_text(s, cx, top0, badge_d, badge_d, [
            (sym, dict(size=17, bold=True, color="#ffffff", align=PP_ALIGN.CENTER)),
        ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        y = top0 + badge_d + 0.12
        D.add_text(s, x, y, w, label_h, [
            (label, dict(size=8, bold=True, color=color, align=PP_ALIGN.CENTER)),
        ], align=PP_ALIGN.CENTER)
        y += label_h
        D.add_text(s, x + 0.05, y, w - 0.1, titre_h, [
            (titre, dict(size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER, line_spacing=1.15)),
        ], align=PP_ALIGN.CENTER)
        y += titre_h + 0.06
        D.add_text(s, x + 0.1, y, w - 0.2, desc_h, [
            (desc, dict(size=8, color=MUTED, align=PP_ALIGN.CENTER, line_spacing=1.25)),
        ], align=PP_ALIGN.CENTER)
        y += desc_h + 0.14
        chip(s, x + w / 2 - 0.75, y, 1.5, 0.28, renvoi, color, size=7.5)

    return s


def slide_demarche_avec_sans_agentic(prs):
    s = content_slide(prs, "Exec summary",
                       "La même démarche absorbe l'outillage : elle n'est jamais dédoublée",
                       color=NAVY)

    chapo = ("Trois temps identiques avec ou sans outillage — ce qui change, c'est ce que le "
             "consultant a dans les mains, et ce qui peut être déployé chez le client.")
    chapo_h = _lignes(chapo, CONTENT_W, 8.5) * (8.5 * 1.25 / 72.0) + 0.06
    D.add_text(s, MARGIN, CONTENT_TOP, CONTENT_W, chapo_h, [
        (chapo, dict(size=8.5, color=NAVY, italic=True, line_spacing=1.25)),
    ])

    SOCLE_FILL, SOCLE_LINE = "#dce6f5", MUTED           # gris-navy neutre
    MODULE_FILL, MODULE_LINE = "#e1fdfa", D.PALETTE[5]  # teal — outillage consultant
    CLIENT_FILL, CLIENT_LINE = "#efe7f7", D.PALETTE[4]  # violet — agentic chez le client

    etapes = [
        ("①", "COMPRENDRE",
         "Le matériau est structuré par le module — les questions posées, elles, ne changent pas.",
         ["Interviews et terrain", "Trames et synthèse outillées", "Rien de déployé à ce stade"]),
        ("②", "DÉFINIR",
         "Le scoring devient explicite et rejouable ; la décision, elle, reste humaine.",
         ["Ateliers de conception", "Fiche produit et scoring", "Candidats qualifiés, gate IA"]),
        ("③", "FAIRE ADOPTER & PROUVER",
         "Les livrables sont générés ; l'adoption, elle, reste de la présence de consultant.",
         ["Coaching, présence dégressive", "Decks et export markdown", "Supervisé, puis délégué"]),
    ]
    registres = [
        ("SANS OUTILLAGE", "présence du consultant", None, SOCLE_FILL, SOCLE_LINE, NAVY),
        ("AVEC LE MODULE", "côté consultant · décliné chapitre 07", "engrenage",
         MODULE_FILL, MODULE_LINE, D.PALETTE[5]),
        ("AGENTIC CHEZ LE CLIENT", "en option, sous gate IA · déployé chapitre 08", "deploiement",
         CLIENT_FILL, CLIENT_LINE, D.PALETTE[4]),
    ]

    label_w = 1.30
    grid_x0 = MARGIN + label_w + 0.14
    grid_w = BORD_DROIT - grid_x0
    n = len(etapes)
    col_gap = 0.16
    col_w = (grid_w - (n - 1) * col_gap) / n

    def _col_px(i):
        return grid_x0 + i * (col_w + col_gap)

    # --- En-tête : les 3 temps sur UN fil continu (la démarche est une, l'outillage
    # ne la coupe pas) — badges navy, pas de couleur de registre ici.
    head_top = CONTENT_TOP + chapo_h + 0.18
    badge_d = 0.42
    change_h = max(_lignes(e[2], col_w, 6.8) for e in etapes) * (6.8 * 1.25 / 72.0) + 0.04
    cx_first = _col_px(0) + col_w / 2
    cx_last = _col_px(n - 1) + col_w / 2
    D.add_rect(s, cx_first, head_top + badge_d / 2 - 0.011, cx_last - cx_first, 0.022, fill=LINE)
    D.add_text(s, MARGIN, head_top, label_w, badge_d, [
        ("LA MÊME DÉMARCHE", dict(size=8, bold=True, color=NAVY, line_spacing=1.1)),
        ("quel que soit l'outillage", dict(size=6.2, color=MUTED, italic=True, space_before=1)),
    ], anchor=MSO_ANCHOR.MIDDLE)
    for i, (sym, nom, change, _pastilles) in enumerate(etapes):
        x = _col_px(i)
        bx = x + col_w / 2 - badge_d / 2
        D.add_rect(s, bx, head_top, badge_d, badge_d, fill=NAVY, rounded=True, radius=0.5)
        D.add_text(s, bx, head_top, badge_d, badge_d, [
            (sym, dict(size=14, bold=True, color="#ffffff", align=PP_ALIGN.CENTER)),
        ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        nom_y = head_top + badge_d + 0.08
        D.add_text(s, x, nom_y, col_w, 0.22, [
            (nom, dict(size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)),
        ], align=PP_ALIGN.CENTER)
        D.add_text(s, x, nom_y + 0.24, col_w, change_h, [
            (change, dict(size=6.8, color=MUTED, italic=True, align=PP_ALIGN.CENTER,
                          line_spacing=1.25)),
        ], align=PP_ALIGN.CENTER)
        if i < n - 1:
            _fleche_h(s, x + col_w, head_top, col_gap, badge_d, color=MUTED, size=12)
    D.add_text(s, MARGIN, head_top + badge_d + 0.30, label_w, change_h, [
        ("CE QUE ÇA CHANGE", dict(size=6.5, bold=True, color=MUTED, line_spacing=1.1)),
    ])
    head_bottom = head_top + badge_d + 0.08 + 0.24 + change_h

    # --- Bandeau de renvoi (pied) : dimensionné à son texte, jamais étiré.
    renvoi = ("Une seule démarche, deux faces de l'agentic : le module outille le consultant "
              "(décliné chapitre 07 · Démarche) ; déployer des agents chez le client est une "
              "décision distincte, sous gate IA (déployé chapitre 08 · Outillage IAP).")
    r_pad = 0.12
    renvoi_h = r_pad + (6.8 * 1.1 / 72.0) + 0.03 \
        + _lignes(renvoi, CONTENT_W - 0.44, 8) * (8 * 1.25 / 72.0) + r_pad
    renvoi_top = CONTENT_BOTTOM - renvoi_h

    # --- Trois rangées de pastilles : la hauteur disponible est CONSOMMÉE (le
    # piège de cette slide est un schéma d'1,5 in laissant 3 in de vide dessous).
    rows_top = head_bottom + 0.16
    row_gap = 0.12
    pill_h = (renvoi_top - 0.16 - rows_top - (len(registres) - 1) * row_gap) / len(registres)
    for r, (label, sous, picto, fill, line, txt_color) in enumerate(registres):
        y = rows_top + r * (pill_h + row_gap)
        px = MARGIN
        if picto:
            _picto(s, picto, px, y + pill_h / 2 - 0.11, 0.22, line)
            px += 0.30
        D.add_text(s, px, y, MARGIN + label_w - px, pill_h, [
            (label, dict(size=7, bold=True, color=NAVY, line_spacing=1.1)),
            (sous, dict(size=6, color=MUTED, italic=True, space_before=2, line_spacing=1.15)),
        ], anchor=MSO_ANCHOR.MIDDLE)
        for i, _e in enumerate(etapes):
            x = _col_px(i)
            texte = etapes[i][3][r]
            D.add_rect(s, x, y, col_w, pill_h, fill=fill, line=line, line_w=1.0,
                       rounded=True, radius=0.22)
            D.add_text(s, x + 0.10, y, col_w - 0.20, pill_h, [
                (texte, dict(size=8, bold=True, color=txt_color, align=PP_ALIGN.CENTER,
                             line_spacing=1.15)),
            ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

    D.add_rect(s, MARGIN, renvoi_top, CONTENT_W, renvoi_h, fill=NAVY, rounded=True, radius=0.08)
    D.add_rect(s, MARGIN, renvoi_top, 0.07, renvoi_h, fill=ACCENT, rounded=True, radius=0.5)
    D.add_text(s, MARGIN + 0.24, renvoi_top + r_pad * 0.5, CONTENT_W - 0.44, renvoi_h - r_pad, [
        ("CE QUE LE DECK DÉCLINE ENSUITE", dict(size=6.8, bold=True, color="#8fd6db")),
        (renvoi, dict(size=8, color="#ffffff", space_before=3, line_spacing=1.25)),
    ])
    return s


# ---------------------------------------------------------------- slide 3
def slide_mission(prs):
    s = content_slide(prs, "Contexte", "Une double mission : transformer ET assainir", color=D.PALETTE[0])
    cards = [
        ("TRANSFORMER", D.PALETTE[0],
         "Cible produit/plateforme : utilisateurs identifiés, valeur, roadmap, "
         "engagements de qualité, gouvernance lisible.",
         "La vision à moyen terme — ce que le sponsor achète."),
        ("ASSAINIR", D.PALETTE[2],
         "Traitement mesurable des gaspillages : flux, RUN, humain, financier, "
         "cognitif, décisionnel, environnemental, IA.",
         # v2.3 : promesse instrumentée, pas un acquis — même honnêteté que le
         # statut interne du cadrage (fin du double discours, finding C1).
         # 2026-09-01 : la moitié CONDITIONNELLE de la formule du cadrage (l.23,
         # « Hypothèse porteuse à prouver, pas un invariant acquis […] suppose un
         # mécanisme de réallocation budgétaire côté client ») avait été effacée —
         # la slide affirmait au présent un KPI qui est un point ouvert MVP3.
         "La capacité récupérée finance la trajectoire produit — hypothèse à "
         "prouver, qui suppose une réallocation budgétaire côté client."),
    ]
    # v2.5 (chantier ③) : les cartes flottaient à CONTENT_TOP+0.5 sans rien
    # au-dessus (~0.5in de blanc sous le titre). La note « ni séquentiels ni
    # optionnels » devient le CHAPEAU (à CONTENT_TOP), les cartes suivent, et
    # la rangée de tensions gagne son étiquette — l'espace se redistribue dans
    # le contenu, pas en vide.
    chapeau_h = 0.5
    D.add_text(s, MARGIN, CONTENT_TOP, CONTENT_W, chapeau_h, [
        ("Deux piliers ni séquentiels ni optionnels : une cible produit sans traitement du "
         "gaspillage manque de capacité pour s'y déployer ; l'inverse reste une réduction "
         "de coûts sans vision.",
         dict(size=D.TYPE["small"], color=NAVY, italic=True, line_spacing=1.25)),
    ])
    card_h = 1.95
    top0 = CONTENT_TOP + chapeau_h + 0.12
    for i, (titre, color, vise, finance) in enumerate(cards):
        x, w = col_x(i, 2)
        D.add_card(s, x, top0, w, card_h, color)
        pad = 0.22
        D.add_text(s, x + pad, top0 + 0.18, w - 2 * pad, 0.3, [
            (titre, dict(size=D.TYPE["h3"], bold=True, color=color))
        ])
        D.add_text(s, x + pad, top0 + 0.58, w - 2 * pad, 0.62, [
            ("CE QU'IL VISE", dict(size=D.TYPE["tiny"], bold=True, color=MUTED)),
            (vise, dict(size=D.TYPE["tiny"], color=NAVY, space_before=2, line_spacing=1.25)),
        ])
        D.add_text(s, x + pad, top0 + 1.28, w - 2 * pad, 0.58, [
            ("CE QU'IL FINANCE", dict(size=D.TYPE["tiny"], bold=True, color=MUTED)),
            (finance, dict(size=D.TYPE["tiny"], color=NAVY, space_before=2, line_spacing=1.25)),
        ])

    label_top = top0 + card_h + 0.22
    D.add_text(s, MARGIN, label_top, CONTENT_W, 0.22, [
        ("L'ÉQUILIBRE QUE LA DOUBLE MISSION TIENT EN PERMANENCE",
         dict(size=7.5, bold=True, color=MUTED)),
    ])
    tens_top = label_top + 0.26
    tens_h = 0.62
    tensions = ["Efficacité du delivery", "Robustesse du RUN", "Valeur perçue (utilisateurs internes)"]
    for i, t in enumerate(tensions):
        x, w = col_x(i, 3)
        D.add_rect(s, x, tens_top, w, tens_h, fill=TRACK, rounded=True, radius=0.12)
        D.add_text(s, x + 0.1, tens_top, w - 0.2, tens_h, [
            (t, dict(size=D.TYPE["tiny"], bold=True, color=NAVY, align=PP_ALIGN.CENTER))
        ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    return s


def slide_pourquoi_contexte(prs):
    """Nouveau (point ②) : dans le chapitre Contexte, le POURQUOI — pourquoi
    proposer cette transformation à un client infra, et maintenant. Trois
    déclencheurs + un pont trait-pour-trait vers la double mission (slide_mission)."""
    s = content_slide(prs, "Contexte",
                       "Pourquoi cette transformation, pour un client infra — et maintenant",
                       color=D.PALETTE[0])
    D.add_text(s, MARGIN, CONTENT_TOP, CONTENT_W, 0.5, [
        ("Trois bascules rendent l'Infra-as-a-Product pertinente — et urgente — pour un client "
         "dont l'infrastructure est encore vécue comme un centre de coûts et un guichet.",
         dict(size=D.TYPE["small"], color=NAVY, italic=True, line_spacing=1.25)),
    ])
    triggers = [
        (D.PALETTE[2], "L'infra subie n'est plus tenable",
         "RUN subi, experts seniors drainés sur du répétitif, gaspillage cloud non maîtrisé, "
         "plateforme contournée : le coût du statu quo ne cesse de monter."),
        # Le « MAIS » du cadrage (l.43) est le motif d'achat du pilier Assainir :
        # sans lui, ce declencheur ne declenche rien. Restaure le 2026-09-01.
        (D.PALETTE[1], "Le modèle produit/plateforme est prouvé",
         "Devenu un standard — mais Gartner : 80 % de grandes organisations avec platform "
         "team en 2026, moins de 30 % de gains mesurables. C'est cet écart qu'Assainir adresse."),
        (D.PALETTE[4], "L'IA rebat les cartes — l'organisation d'abord",
         "L'IA amplifie une organisation mûre, jamais l'inverse. S'y préparer maintenant "
         "(doctrine confidentialité-first) évite de la subir plus tard."),
    ]
    lead_h, bridge_h = 0.55, 0.72
    top0 = CONTENT_TOP + lead_h + 0.1
    card_h = (CONTENT_BOTTOM - bridge_h - 0.18) - top0
    pad = 0.22
    for i, (color, titre, corps) in enumerate(triggers):
        x, w = col_x(i, 3)
        D.add_card(s, x, top0, w, card_h, color)
        tx = x + 0.08 + pad
        tw = w - 0.08 - 2 * pad
        D.add_text(s, tx, top0 + 0.22, tw, card_h - 0.44, [
            (f"DÉCLENCHEUR {i + 1}", dict(size=6.5, bold=True, color=MUTED)),
            (titre, dict(size=D.TYPE["small"], bold=True, color=color, space_before=3, line_spacing=1.05)),
            (corps, dict(size=9, color=NAVY, space_before=8, line_spacing=1.25)),
        ])
    bridge_top = CONTENT_BOTTOM - bridge_h
    D.add_rect(s, MARGIN, bridge_top, CONTENT_W, bridge_h, fill=TRACK, rounded=True, radius=0.1)
    D.add_rect(s, MARGIN, bridge_top, 0.08, bridge_h, fill=D.PALETTE[0], rounded=True, radius=0.5)
    D.add_text(s, MARGIN + 0.28, bridge_top, CONTENT_W - 0.46, bridge_h, [
        ("Et surtout — nos deux missions répondent trait pour trait aux deux douleurs du client.",
         dict(size=8.5, bold=True, color=NAVY, line_spacing=1.05)),
        ("Subir le RUN → TRANSFORMER (cible produit/plateforme) ; le gaspillage → ASSAINIR "
         "(capacité récupérée à réinvestir dans la trajectoire — sous réserve d'une "
         "réallocation côté client).",
         dict(size=8, color=MUTED, space_before=2, line_spacing=1.15)),
    ], anchor=MSO_ANCHOR.MIDDLE)
    return s


# --- Nouveau (2026-09-01) : « qui achète, contre quoi ». La section
# §Positionnement & achat du cadrage (l.36) fait foi POUR LE DECK depuis la
# v2.3, mais n'y avait jamais été redescendue : 40 slides disaient COMMENT on
# fait la mission, aucune contre quel achat alternatif elle se gagne.
# Forme retenue (deck-design-library) : grille-référentiel à libellé propre par
# ligne (pattern 15) — lignes = les 4 achats alternatifs (cadrage l.52-57),
# colonnes = ce qu'il apporte / ce qui lui manque ; la règle « un sur N en
# accent » est appliquée à la COLONNE de droite (teinte pâle), c'est là que
# l'offre gagne. En pied, bandeau transverse (pattern 18) : collision de nom
# (l.63) et réponse au sponsor « je ne veux que la baisse de coûts » (l.57, l.78).
def slide_qui_achete(prs):
    s = content_slide(prs, "Contexte",
                       "Le sponsor est la DSI — l'offre se gagne contre quatre achats partiels",
                       color=D.PALETTE[0])
    couleur = D.PALETTE[0]

    def lh(pt):
        return pt * 1.25 / 72.0

    # --- Chapeau : qui achète, et dans quel langage (cadrage l.48) -----------
    lead1 = ("Sponsor qualifié : DSI (direction des systèmes d'information) ou direction "
             "infrastructure, sur une ligne budgétaire transformation — jamais le budget "
             "RUN (exploitation courante).")
    lead2 = ("La modernisation infra recule face à la cyber dans les priorités budgétaires "
             "2026 (baromètre Abraxio) : le langage de vente est « récupérer de la capacité "
             "humaine rare », pas « moderniser l'infra ».")
    lead_h = (_lignes(lead1, CONTENT_W, 9.5) * lh(9.5)
              + _lignes(lead2, CONTENT_W, 7.5) * lh(7.5) + 0.09)
    D.add_text(s, MARGIN, CONTENT_TOP, CONTENT_W, lead_h, [
        (lead1, dict(size=9.5, color=NAVY, italic=True, line_spacing=1.25)),
        (lead2, dict(size=7.5, color=MUTED, space_before=3, line_spacing=1.25)),
    ])

    # --- Bandeau transverse : dimensionné AVANT la grille, qui prend le reste
    # (jamais de panneau étiré sur la hauteur restante).
    bandeaux = [
        (TRACK, MUTED, NAVY, "CE QUE LES QUATRE ALTERNATIVES N'ONT PAS",
         "L'étiquette « Infrastructure as a Product » est déjà prise — Thoughtworks (conseil), "
         "Itential (plateforme). Le différenciateur n'est pas le nom mais le couplage produit "
         "+ gaspillage + doctrine IA, angle mort commun des quatre."),
        (NAVY, "#8891b3", "#ffffff", "LA RÉPONSE AU « JE NE VEUX QUE LA BAISSE DE COÛTS »",
         "Une mission flash d'entrée — intake, gate confidentialité, pilote court — puis la "
         "trajectoire ; jamais l'assainissement seul. Sous pression IA, un cas d'usage sur "
         "données publiques est packagé dès l'intake : « celui-ci, tout de suite, sous gate » "
         "(chapitre IA)."),
    ]
    _, band_w = col_x(0, 2)
    band_pad = 0.12
    band_lignes = max(_lignes(b[4], band_w - 2 * band_pad - 0.04, 7) for b in bandeaux)
    band_h = 2 * band_pad + lh(7) + 0.04 + band_lignes * lh(7) + 0.03
    band_top = CONTENT_BOTTOM - band_h

    # --- Grille : 4 lignes (les alternatives) x 3 colonnes ------------------
    alternatives = [
        ("Ne rien faire",
         "Zéro coût apparent.",
         "Le coût du statu quo monte",
         "C'est lui que l'Assessment flash chiffre (déclencheur ①)."),
        ("FinOps outillé seul",
         "Mesure le gaspillage : marché mature, gaspillage cloud estimé à 29 % (Flexera).",
         "Ni cible produit, ni réallocation",
         "IAP se place en aval : le chiffre devient une capacité produit gouvernée, pas "
         "seulement une économie."),
        ("Platform engineering pur",
         "La cible produit/plateforme, un modèle devenu standard.",
         "La cible sans le financement",
         "Reproduit l'écart 80/30 : 80 % de platform teams en 2026, moins de 30 % de gains "
         "mesurables (Gartner)."),
        ("AIOps / agentic outillé",
         "Time-to-value court (ServiceNow, Datadog).",
         "Automatise le RUN sans transformer",
         "Plus de 40 % des projets agentic seront abandonnés d'ici 2027 (Gartner, juin 2025)."),
    ]
    c1_w, c2_w = 1.80, 2.45
    c3_w = CONTENT_W - c1_w - c2_w
    c1_x = MARGIN
    c2_x = c1_x + c1_w
    c3_x = c2_x + c2_w

    head_h = 0.18
    head_top = CONTENT_TOP + lead_h + 0.10
    grid_top = head_top + head_h + 0.05
    grid_h = band_top - 0.16 - grid_top
    row_h = grid_h / len(alternatives)

    _header_cell(s, c1_x, head_top, c1_w, head_h, "L'ACHAT ALTERNATIF")
    _header_cell(s, c2_x + 0.12, head_top, c2_w - 0.12, head_h, "CE QU'IL APPORTE")
    _header_cell(s, c3_x + 0.14, head_top, c3_w - 0.14, head_h,
                 "CE QUI LUI MANQUE — LA RÉPONSE IAP", color=couleur)

    # « Un sur N en accent » porté par la COLONNE, pas par une ligne : la
    # troisième colonne est la seule teintée — c'est là que l'offre répond.
    D.add_rect(s, c3_x, grid_top, c3_w, grid_h, fill="#D9ECFF", rounded=True, radius=0.06)

    chip_h = 0.40
    for i, (nom, apporte, manque, reponse) in enumerate(alternatives):
        y = grid_top + i * row_h
        if i:  # séparateurs fins : la grille sans le tableau (pattern 11)
            D.add_rect(s, MARGIN, y, CONTENT_W, 0.012, fill=LINE)
        chip_y = y + row_h / 2 - chip_h / 2
        D.add_rect(s, c1_x, chip_y, c1_w - 0.16, chip_h, fill=TRACK, rounded=True, radius=0.1)
        D.add_text(s, c1_x + 0.10, chip_y, c1_w - 0.36, chip_h, [
            (nom, dict(size=8, bold=True, color=NAVY, line_spacing=1.05)),
        ], anchor=MSO_ANCHOR.MIDDLE)
        D.add_text(s, c2_x + 0.12, y, c2_w - 0.24, row_h, [
            (apporte, dict(size=7.5, color=NAVY, line_spacing=1.25)),
        ], anchor=MSO_ANCHOR.MIDDLE)
        D.add_text(s, c3_x + 0.14, y, c3_w - 0.28, row_h, [
            (manque, dict(size=7.5, bold=True, color=couleur, line_spacing=1.1)),
            (reponse, dict(size=7.5, color=NAVY, space_before=2, line_spacing=1.25)),
        ], anchor=MSO_ANCHOR.MIDDLE)

    for i, (fill, label_c, texte_c, label, corps) in enumerate(bandeaux):
        x, w = col_x(i, 2)
        D.add_rect(s, x, band_top, w, band_h, fill=fill, rounded=True, radius=0.08)
        D.add_text(s, x + band_pad + 0.02, band_top + band_pad, w - 2 * band_pad - 0.04,
                   band_h - 2 * band_pad, [
            (label, dict(size=7, bold=True, color=label_c)),
            (corps, dict(size=7, color=texte_c, space_before=3, line_spacing=1.25)),
        ])
    return s


# ---------------------------------------------------------------- slide 4
def slide_gate_ia(prs):
    s = content_slide(prs, "IA", "Les données du client gouvernent le choix du modèle IA", color=D.PALETTE[4])
    D.add_text(s, MARGIN, CONTENT_TOP, CONTENT_W, 0.35, [
        ("Checkpoint toujours humain avant tout usage IA sur données client — "
         "iap-ai-data-confidentiality-gate, quel que soit le mode d'exécution retenu.",
         dict(size=D.TYPE["tiny"], color=MUTED, line_spacing=1.2)),
    ])
    rows = [
        ("D0", "Public", "Articles publics, docs méthodo", "IA externe possible"),
        ("D1", "Interne", "Organisation macro, catalogue anonymisé", "IA client recommandée"),
        ("D2", "Confidentiel", "Notes d'interview, reporting, portefeuille", "IA client ou LLM privé"),
        ("D3", "Restreint", "Tickets détaillés, logs, CMDB, IAM", "LLM local, contrôles forts"),
        ("D4", "Critique", "Secrets de production, données réglementées", "Local/on-prem, sans IA générative"),
    ]
    row_top = CONTENT_TOP + 0.45
    row_h = 0.62
    row_gap = 0.1
    label_w = 1.35
    desc_w = 4.2
    usage_w = CONTENT_W - label_w - desc_w - 2 * 0.2
    for i, (code, nom, exemples, usage) in enumerate(rows):
        y = row_top + i * (row_h + row_gap)
        chip(s, MARGIN, y, label_w, row_h, f"{code} · {nom}", SEVERITE[i], size=D.TYPE["tiny"])
        D.add_text(s, MARGIN + label_w + 0.2, y, desc_w, row_h, [
            (exemples, dict(size=D.TYPE["tiny"], color=NAVY, line_spacing=1.1)),
        ], anchor=MSO_ANCHOR.MIDDLE)
        D.add_text(s, MARGIN + label_w + 0.2 + desc_w + 0.2, y, usage_w, row_h, [
            (usage, dict(size=D.TYPE["tiny"], color=MUTED, italic=True, line_spacing=1.1)),
        ], anchor=MSO_ANCHOR.MIDDLE)
    return s


# ---------------------------------------------------------------- slide 5
def slide_why_iap(prs):
    """Nouveau (point ⑤) : OUVRE le chapitre Proposition (la thèse). Le POURQUOI de
    l'Infra-as-a-Product — trois bascules produit, chacune ancrée sur un persona/une
    douleur déjà posés. (2e passe : la maturité est partie au chapitre KPI, le
    sous-chapitre « Technique IAP » a donc disparu — why_iap ouvre la Proposition.)"""
    s = content_slide(prs, "Proposition",
                       "Pourquoi « Infrastructure as a Product » — le socle de la proposition",
                       color=D.PALETTE[1])
    claim_h = 0.95
    D.add_rect(s, MARGIN, CONTENT_TOP, CONTENT_W, claim_h, fill=TRACK, rounded=True, radius=0.1)
    D.add_rect(s, MARGIN, CONTENT_TOP, 0.08, claim_h, fill=D.PALETTE[1], rounded=True, radius=0.5)
    D.add_text(s, MARGIN + 0.3, CONTENT_TOP, CONTENT_W - 0.5, claim_h, [
        ("Traiter l'infrastructure comme un produit, pas comme un guichet de tickets.",
         dict(size=14, bold=True, color=NAVY, line_spacing=1.05)),
        ("Un produit a des utilisateurs, un cycle de vie et une valeur mesurée — trois bascules "
         "qui répondent directement aux personas et à leurs douleurs.",
         dict(size=9, color=MUTED, space_before=4, line_spacing=1.15)),
    ], anchor=MSO_ANCHOR.MIDDLE)

    piliers = [
        (D.PALETTE[5], "Des utilisateurs, pas des tickets",
         "On conçoit l'adoption — self-service, onboarding, parcours — au lieu de subir un "
         "guichet que le contournement rend inutile.",
         "l'Utilisateur applicatif"),
        (D.PALETTE[0], "Un cycle de vie, une équipe qui en répond",
         "Le produit a un propriétaire, une roadmap et une dette gérée : on sort du RUN subi "
         "et on récupère de la capacité.",
         "Infra & RUN"),
        (D.PALETTE[3], "Un pilotage par la valeur",
         "On mesure l'usage et la valeur produite, pas l'activité : un signal de flux fiable, "
         "des KPIs de mission — pas du reporting-miroir.",
         "Management & Sponsor"),
    ]
    top = CONTENT_TOP + claim_h + 0.25
    card_h = CONTENT_BOTTOM - top
    pad = 0.22
    for i, (color, titre, corps, ancre) in enumerate(piliers):
        cx, cw = col_x(i, 3)
        D.add_card(s, cx, top, cw, card_h, color)
        tx = cx + 0.08 + pad
        tw = cw - 0.08 - 2 * pad
        D.add_text(s, tx, top + 0.22, tw, card_h - 0.44, [
            (titre, dict(size=D.TYPE["small"], bold=True, color=color, line_spacing=1.05)),
            (corps, dict(size=9, color=NAVY, space_before=8, line_spacing=1.25)),
            ("RÉPOND À", dict(size=6.5, bold=True, color=MUTED, space_before=12)),
            (ancre, dict(size=9, bold=True, color=color, space_before=2)),
        ])
    return s


def slide_maturite(prs):
    s = content_slide(prs, "KPI",
                       "La grille de maturité : deux échelles distinctes, mesurées dans le temps",
                       color=D.PALETTE[0])
    # Placée en fin de chapitre KPI (juste avant le cas chiffré) et CLARIFIÉE
    # (point ①) : c'est la 3e famille de KPIs (grille de maturité). Message resserré :
    # deux échelles ne mesurant PAS la même chose, chacune gouvernant une décision
    # différente ; le KPI = le DELTA dans le temps, pas le niveau absolu. Ambiguïté
    # « Remplace le M0–M4 » toujours levée (badge « Où se lit l'axe IA »).
    D.add_text(s, MARGIN, CONTENT_TOP, CONTENT_W, 0.42, [
        ("La 3ᵉ famille de KPIs. Deux lectures qui ne mesurent pas la même chose et se lisent "
         "séparément ; le KPI de progression, c'est le DELTA par pilier entre T0 et chaque "
         "réévaluation — pas le niveau absolu.",
         dict(size=8, color=MUTED, italic=True, line_spacing=1.2)),
    ])
    x0, w0 = col_x(0, 2)
    x1, w1 = col_x(1, 2)

    head_top = CONTENT_TOP + 0.5
    D.add_text(s, x0, head_top, w0, 0.5, [
        ("CAPACITÉ IA DU CLIENT (M0–M4)", dict(size=D.TYPE["tiny"], bold=True, color=NAVY)),
        ("→ gouverne le choix du modèle IA et le gate",
         dict(size=7.5, color=D.PALETTE[0], space_before=2, line_spacing=1.05)),
    ])
    niveaux = [
        ("M0", "Pas d'IA interne utilisable", "Méthodo générique, données anonymisées"),
        ("M1", "IA interne basique", "Synthèses internes, pas d'analyse critique auto"),
        ("M2", "IA privée avec RAG", "Diagnostic documentaire, consolidation"),
        ("M3", "Plateforme IA gouvernée", "Workflows agentic contrôlés"),
        ("M4", "IA industrielle", "Agents spécialisés à fort volume, contrôle humain"),
    ]
    row_top = head_top + 0.52
    row_h = 0.58
    row_gap = 0.06
    for i, (code, titre, strat) in enumerate(niveaux):
        y = row_top + i * (row_h + row_gap)
        chip(s, x0, y, 0.62, row_h, code, D.PALETTE[0], size=D.TYPE["tiny"])
        D.add_text(s, x0 + 0.62 + 0.15, y, w0 - 0.77 - 0.15, row_h, [
            (titre, dict(size=D.TYPE["tiny"], bold=True, color=NAVY)),
            (strat, dict(size=8, color=MUTED, space_before=1, line_spacing=1.1)),
        ], anchor=MSO_ANCHOR.MIDDLE)

    D.add_text(s, x1, head_top, w1, 0.5, [
        ("MATURITÉ PRODUIT / PLATEFORME (grille V3.2)", dict(size=D.TYPE["tiny"], bold=True, color=NAVY)),
        ("→ gouverne la trajectoire de transformation",
         dict(size=7.5, color=D.PALETTE[1], space_before=2, line_spacing=1.05)),
    ])
    piliers = [
        ("Équipe Produit", "Adjacent", False),
        ("Excellence Technique", "Cœur du périmètre", True),
        ("Culture de l'Entreprise Agile", "Adjacent", False),
        ("Agilité à l'Échelle", "Cœur du périmètre", True),
        ("IA, Agentic et Organisation Augmentée", "Où se lit l'axe IA (M0–M4)", True),
    ]
    for i, (nom, badge, coeur) in enumerate(piliers):
        y = row_top + i * (row_h + row_gap)
        color = D.couleur_pilier(i)
        D.add_dot(s, x1, y + row_h / 2 - 0.07, 0.14, color)
        tx = x1 + 0.28
        tw = w1 - 0.28
        D.add_text(s, tx, y, tw, row_h, [
            (nom, dict(size=D.TYPE["tiny"], bold=coeur, color=NAVY, line_spacing=1.1)),
            (badge, dict(size=8, color=(color if coeur else MUTED), space_before=1)),
        ], anchor=MSO_ANCHOR.MIDDLE)
    return s


# --- Nouveau (réouverture de périmètre, arbitrage 2026-07-21) : la Product
# Discovery (personas/parcours/pain points), délibérément fusionnée dans
# iap-product-definition pour MVP1 (§Décision de cadrage, ligne 236), est
# rouverte ici en une slide dédiée. Quatre parties prenantes de la couverture
# d'interview (§Synthesis, "répartition par persona : infra/utilisateur/
# management/sponsor", ligne 457) — chacune sa voix, sa question directrice
# (reprise des questions des §Agents), son irritant, son attente. Ouvre le
# chapitre Personas (02) : on sait QUI l'on transforme avant d'exposer ses
# douleurs (chapitre Besoins & douleurs) puis notre réponse (Proposition).
def slide_personas(prs):
    s = content_slide(prs, "Personas",
                       "Quatre parties prenantes interrogées séparément — leur voix, leur posture",
                       color=D.PALETTE[5])
    D.add_text(s, MARGIN, CONTENT_TOP, CONTENT_W, 0.5, [
        ("Product Discovery fusionnée dans iap-product-definition en MVP1 — mais chaque partie "
         "prenante répond à la même trame, pour révéler convergences ET divergences plutôt qu'un "
         "diagnostic monolithique.",
         dict(size=8, color=MUTED, italic=True, line_spacing=1.2)),
    ])

    P = D.PALETTE
    # 2×2 cartes persona. Tuple : nom, accent, rôle (1 ligne), verbatim, ce qu'il
    # SUBIT, ce qu'il VISE, posture face à la transformation + sa couleur (feu
    # tricolore sémantique : allié=vert, sceptique=rouge, vigilant=or). La posture
    # est l'ajout du brainstorm — le comité lit le paysage politique, pas 4 listes.
    personas = [
        ("Infra & RUN", P[0], "Tient l'exploitation, subit les astreintes.",
         "« Opérable sans sacrifier le delivery ? »",
         "Experts seniors mobilisés sur du répétitif.",
         "Capacité récupérée, RUN maîtrisé.",
         "Vigilant", P[3]),
        ("Utilisateur applicatif", P[5], "Consomme la plateforme — ou la contourne.",
         "« Pourquoi adopterais-je la plateforme ? »",
         "Guichet unique, contournement plus rapide.",
         "Un self-service adopté par choix.",
         "Sceptique", P[2]),
        ("Management", P[3], "Expert devenu manager, pilote à vue.",
         "« Comment piloter avec un signal fiable ? »",
         "Reporting-miroir et micromanagement.",
         "Un signal de flux de confiance.",
         "Allié", P[1]),
        ("Sponsor", P[4], "Porte le budget et la promesse business.",
         "« Quel problème business règle-t-on ? »",
         "Craint une transformation cosmétique.",
         "Problème business réglé, KPIs de mission.",
         "Allié exigeant", P[1]),
    ]
    top0 = CONTENT_TOP + 0.62
    row_gap = 0.18
    card_h = (CONTENT_BOTTOM - top0 - row_gap) / 2
    pad = 0.2
    chip_w, chip_h = 1.2, 0.24
    for i, (nom, accent, role, verbatim, subit, vise, posture, cposture) in enumerate(personas):
        r, c = divmod(i, 2)
        cx, cw = col_x(c, 2)
        cy = top0 + r * (card_h + row_gap)
        D.add_card(s, cx, cy, cw, card_h, accent)
        tx = cx + 0.07 + pad
        tw = cw - 0.07 - 2 * pad
        D.add_text(s, tx, cy + 0.15, tw, card_h - 0.15 - chip_h - 0.18, [
            (nom, dict(size=D.TYPE["tiny"], bold=True, color=accent, line_spacing=1.0)),
            (role, dict(size=8, color=MUTED, space_before=1, line_spacing=1.1)),
            (verbatim, dict(size=8, italic=True, color=NAVY, space_before=5, line_spacing=1.1)),
            ("Subit — " + subit, dict(size=8, color=MUTED, space_before=6, line_spacing=1.1)),
            ("Vise — " + vise, dict(size=8, color=accent, space_before=3, line_spacing=1.1)),
        ])
        chip(s, cx + cw - pad - chip_w, cy + card_h - 0.14 - chip_h, chip_w, chip_h,
             posture.upper(), cposture, size=6.5)
    return s


# --- Nouveau (arbitrage cadrage validé) : corollaire direct de slide_personas.
# Interviewer chaque partie prenante SÉPARÉMENT (§Synthesis) n'a de sens que si
# l'on garde les divergences au lieu de les lisser en consensus — cette slide
# les rend explicites. Réutilise les couleurs d'accent persona de
# slide_personas (Infra=bleu, Utilisateur=teal, Management=or, Sponsor=violet)
# et introduit le RSSI (porteur du gate) en rouge = criticité/blocage. Le
# symbole de tension « ⟂ » du cadrage est rendu par le connecteur texte « en
# tension avec » plutôt que par le glyphe (non garanti dans la fonte du
# template, cf. _GLYPHES_SANS_GRAS) — même prudence que pour « ⟲ ». Rangées
# dimensionnées à leur contenu (pas de panneau sur-étiré).
def slide_personas_divergences(prs):
    s = content_slide(prs, "Personas",
                       "Interroger chaque persona séparément révèle des tensions "
                       "qu'un diagnostic fusionné lisserait",
                       color=D.PALETTE[5])
    # Note d'intro retirée (redondante avec le sous-titre) : les rangées démarrent
    # plus haut pour laisser place, en bas, à la synthèse « pont » vers la Proposition.
    # Passe de design 2026-07-23 — pattern 7 du catalogue deck-design-library
    # (« rangée de cartes, une en accent ») : la rangée ANGLE MORT (Sponsor ⟂ RSSI)
    # est la seule teintée (fond rouge très pâle + contour rouge) — rouge =
    # sémantique d'alerte, pas décoration ; les 3 tensions instruites restent
    # des cartes blanches identiques.
    c_infra = D.PALETTE[0]   # Infra & RUN — bleu (comme slide_personas)
    c_user = D.PALETTE[5]    # Utilisateur applicatif — teal
    c_mgmt = D.PALETTE[3]    # Management — or
    c_spon = D.PALETTE[4]    # Sponsor — violet
    c_rssi = D.PALETTE[2]    # RSSI — rouge = porteur du gate, criticité/blocage
    rows = [
        (("Management", c_mgmt), ("Infra & RUN", c_infra), None,
         "Le même métrique de flux, lu « signal de pilotage de confiance » d'un côté, "
         "« surveillance » de l'autre."),
        (("Sponsor", c_spon), ("Infra & RUN", c_infra), None,
         "Horizon : valeur business rapide et visible d'un côté, soulagement durable et "
         "structurel du RUN de l'autre."),
        (("Utilisateur applicatif", c_user), ("Infra & RUN", c_infra), None,
         "Self-service adopté par choix face à l'opérabilité sans sacrifier le delivery : "
         "qui absorbe le coût du self-service ?"),
        (("Sponsor", c_spon), ("RSSI", c_rssi), "ANGLE MORT",
         "Vitesse de démonstration face au gate confidentialité, bloquant sur donnée client."),
    ]
    top0 = CONTENT_TOP + 0.1
    n = len(rows)
    synth_h, note_h = 0.56, 0.34
    bottom_reserve = synth_h + 0.14 + note_h + 0.10
    row_gap = 0.12
    region_bot = CONTENT_BOTTOM - bottom_reserve
    row_h = (region_bot - top0 - (n - 1) * row_gap) / n
    name_w = 2.7
    x_name = MARGIN + 0.2
    x_fric = x_name + name_w + 0.25
    fric_w = (MARGIN + CONTENT_W) - x_fric - 0.15
    for i, ((nomA, colA), (nomB, colB), tag, friction) in enumerate(rows):
        y = top0 + i * (row_h + row_gap)
        accent = tag is not None   # « un sur N » : l'angle mort, seul non instruit
        D.add_rect(s, MARGIN, y, CONTENT_W, row_h,
                   fill="#fbeeed" if accent else "#ffffff",
                   line=c_rssi if accent else LINE,
                   line_w=1.0 if accent else 0.75,
                   rounded=True, radius=0.08)
        # liseré scindé : moitié haute = couleur A, moitié basse = couleur B
        D.add_rect(s, MARGIN, y, 0.06, row_h / 2, fill=colA, rounded=True, radius=0.5)
        D.add_rect(s, MARGIN, y + row_h / 2, 0.06, row_h / 2, fill=colB, rounded=True, radius=0.5)
        lignes = [
            (nomA, dict(size=D.TYPE["tiny"], bold=True, color=colA, line_spacing=1.0)),
            ("en tension avec", dict(size=7, italic=True, color=MUTED, space_before=3, space_after=3)),
            (nomB, dict(size=D.TYPE["tiny"], bold=True, color=colB, line_spacing=1.0)),
        ]
        if tag:
            lignes.append((tag, dict(size=6.5, bold=True, color=c_rssi, space_before=3)))
        D.add_text(s, x_name, y + 0.08, name_w, row_h - 0.16, lignes, anchor=MSO_ANCHOR.MIDDLE)
        D.add_text(s, x_fric, y + 0.08, fric_w, row_h - 0.16, [
            (friction, dict(size=8, color=NAVY, line_spacing=1.2)),
        ], anchor=MSO_ANCHOR.MIDDLE)

    # Ligne de synthèse « pont » (issue du brainstorm) : les divergences ne se
    # tranchent pas, la méthode (ch. Proposition) les tient des deux bouts —
    # transforme la slide de « voici des conflits » en « voici pourquoi on n'a pas
    # à choisir un camp », et donne l'élan vers la suite.
    synth_top = top0 + n * row_h + (n - 1) * row_gap + 0.14
    D.add_rect(s, MARGIN, synth_top, CONTENT_W, synth_h, fill=TRACK, rounded=True, radius=0.12)
    D.add_rect(s, MARGIN, synth_top, 0.07, synth_h, fill=c_user, rounded=True, radius=0.5)
    D.add_text(s, MARGIN + 0.26, synth_top, CONTENT_W - 0.42, synth_h, [
        ("Ces tensions ne se tranchent pas — on les tient des deux bouts.",
         dict(size=8.5, bold=True, color=c_user, line_spacing=1.05)),
        ("La méthode (ch. Proposition) : la métrique de flux = signal partagé, le gate "
         "confidentialité = non négociable.",
         dict(size=8, color=NAVY, space_before=2, line_spacing=1.1)),
    ], anchor=MSO_ANCHOR.MIDDLE)

    note_top = synth_top + synth_h + 0.10
    D.add_text(s, MARGIN, note_top, CONTENT_W, note_h, [
        ("Angles morts, non interrogés à ce stade : le client métier consommateur des services, "
         "le RSSI (porteur du gate), le junior / nouvel arrivant.",
         dict(size=7.5, color=MUTED, italic=True, line_spacing=1.15)),
    ])
    return s


# ---------------------------------------------------------------- Besoins & douleurs
# Nouveau (restructuration 2026-07-22) : la grille des 8 familles de gaspillage,
# jusqu'ici empaquetée dans slide_gaspillages avec la chaîne de traitement et le
# score, est isolée ici — elle appartient au chapitre « Besoins & douleurs » (le
# langage commun qui rend une douleur nommable, donc détectable et traitable),
# tandis que la MÉTHODE de traitement (chaîne + score) reste au chapitre
# Proposition. Accent unifié sur la couleur du chapitre Douleurs (PALETTE[2]) :
# les 8 familles se distinguent par leur libellé, pas par 8 teintes sans clé.
# Passe de design 2026-07-23 — règle « un sur N en accent » (principes
# transversaux + pattern 3 du catalogue deck-design-library) : la famille IA,
# seule famille que cette méthode NOMME comme gaspillage (cas gadget,
# automatisation sans garde-fous — la doctrine du deck), reçoit un fill navy
# plein ; les 7 autres restent des cartes blanches identiques.
def slide_familles(prs):
    s = content_slide(prs, "Besoins & douleurs",
                       "Les 8 familles de gaspillage — le langage commun qui rend les douleurs traitables",
                       color=D.PALETTE[2])
    D.add_text(s, MARGIN, CONTENT_TOP, CONTENT_W, 0.42, [
        ("Nommer la famille, c'est déjà pouvoir la détecter, la quantifier et la prioriser "
         "(méthode de traitement → chapitre Proposition).",
         dict(size=D.TYPE["small"], color=NAVY, italic=True, line_spacing=1.25)),
    ])
    familles = [
        ("Flux", "Attentes, validations multiples"),
        ("Humain", "Experts seniors sur tâches répétitives"),
        ("RUN", "Incidents récurrents, demandes répétées"),
        ("Financier", "Surdimensionnement, ressources non décommissionnées"),
        ("Cognitif", "Trop d'outils, procédures complexes"),
        ("Décisionnel", "Arbitrages subjectifs, priorisation opaque"),
        ("Environnemental", "Ressources inutilisées, environnements non éteints"),
        ("IA", "Cas d'usage gadget, automatisation sans garde-fous"),
    ]
    # 2 colonnes x 4 rangées : remplit la hauteur de la slide dédiée sans étirer
    # chaque carte (défaut « panneau sur-étiré »). Lecture gauche->droite par
    # paire (col = i % 2, row = i // 2).
    n_rows = 4
    region_top = CONTENT_TOP + 0.55
    row_gap = 0.14
    row_h = (CONTENT_BOTTOM - region_top - (n_rows - 1) * row_gap) / n_rows
    for i, (nom, ex) in enumerate(familles):
        col = i % 2
        row = i // 2
        x, w = col_x(col, 2)
        y = region_top + row * (row_h + row_gap)
        accent = (nom == "IA")   # « un sur N » : la famille portée par la doctrine
        if accent:
            D.add_rect(s, x, y, w, row_h, fill=NAVY, rounded=True, radius=0.1)
        else:
            D.add_rect(s, x, y, w, row_h, fill="#ffffff", line=LINE, line_w=0.75,
                       rounded=True, radius=0.1)
        D.add_rect(s, x, y, 0.06, row_h, fill=D.PALETTE[2], rounded=True, radius=0.5)
        D.add_text(s, x + 0.2, y + 0.06, w - 0.34, row_h - 0.12, [
            (f"{i + 1}. {nom}", dict(size=D.TYPE["small"], bold=True,
                                     color="#ffffff" if accent else NAVY)),
            (ex, dict(size=8, color="#c7cbe0" if accent else MUTED,
                      space_before=2, line_spacing=1.1)),
        ], anchor=MSO_ANCHOR.MIDDLE)
    return s


# ---------------------------------------------------------------- Proposition
# Reframe (restructuration 2026-07-22) : la grille des 8 familles est partie au
# chapitre « Besoins & douleurs » (slide_familles). Ne reste ici que la MÉTHODE
# de traitement — la chaîne Détecter->Prévenir et le score de priorisation — ré-
# ancrée vers le haut pour combler l'espace libéré par la grille retirée, avec
# une accroche en tête pour éviter un vide sous le titre.
# Passe de design 2026-07-23 : la chaîne, jusqu'ici 10 pilules grises identiques
# en grille 2×5 (effet « tableau de chips »), est redessinée selon le pattern 4
# du catalogue deck-design-library (« flux numéroté en quinconce, badges +
# connecteur, sans cadres ») : badges numérotés reliés par un fil, 2e rangée
# décalée d'un demi-slot, et « un sur N en accent » — seule l'étape 6 (Prioriser)
# est remplie en couleur pleine, car c'est elle qui produit le score détaillé
# dans le panneau navy juste en dessous.
def slide_gaspillages(prs):
    s = content_slide(prs, "Proposition",
                       "Du gaspillage au backlog priorisé : chaîne de traitement + score",
                       color=D.PALETTE[1])
    D.add_text(s, MARGIN, CONTENT_TOP, CONTENT_W, 0.5, [
        ("Chaque famille de gaspillage (chapitre précédent) passe par la même chaîne de "
         "traitement, puis reçoit un score explicite qui la classe dans un backlog priorisé "
         "— jamais un tri à l'intuition.",
         dict(size=D.TYPE["small"], color=NAVY, italic=True, line_spacing=1.25)),
    ])

    chain_top = CONTENT_TOP + 0.65
    D.add_text(s, MARGIN, chain_top, CONTENT_W, 0.24, [
        ("CHAÎNE DE TRAITEMENT — de la détection à la prévention",
         dict(size=8, bold=True, color=MUTED))
    ])
    etapes = ["Détecter", "Qualifier", "Quantifier", "Cause racine", "Pattern",
              "Prioriser", "Expérimenter", "Mesurer", "Industrialiser", "Prévenir"]
    step_top = chain_top + 0.34
    n = 5
    slot = CONTENT_W / (n + 0.5)   # 2e rangée décalée d'un demi-slot (quinconce)
    badge_d = 0.32
    unit_h = 0.62                  # badge + libellé
    row_gap2 = 0.10
    accent_idx = 5                 # « Prioriser » — l'étape qui produit le score
    for row in range(2):
        x0 = MARGIN + (slot / 2 if row == 1 else 0.0)
        y = step_top + row * (unit_h + row_gap2)
        cy = y + badge_d / 2
        # fil du flux : connecteur horizontal reliant les badges de la rangée
        D.add_rect(s, x0 + slot / 2, cy - 0.01, (n - 1) * slot, 0.02, fill=LINE)
        for col in range(n):
            i = row * n + col
            accent = (i == accent_idx)
            bx = x0 + col * slot + slot / 2 - badge_d / 2
            D.add_rect(s, bx, y, badge_d, badge_d,
                       fill=D.PALETTE[1] if accent else "#ffffff",
                       line=None if accent else D.PALETTE[1], line_w=1.0,
                       rounded=True, radius=0.5)
            D.add_text(s, bx, y, badge_d, badge_d, [
                (str(i + 1), dict(size=8, bold=True,
                                  color="#ffffff" if accent else D.PALETTE[1],
                                  align=PP_ALIGN.CENTER)),
            ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
            D.add_text(s, x0 + col * slot, y + badge_d + 0.04, slot, 0.22, [
                (etapes[i], dict(size=8, bold=True,
                                 color=D.PALETTE[1] if accent else NAVY,
                                 align=PP_ALIGN.CENTER)),
            ], align=PP_ALIGN.CENTER)

    score_top = step_top + 2 * unit_h + row_gap2 + 0.30
    score_h = CONTENT_BOTTOM - score_top - 0.12
    D.add_rect(s, MARGIN, score_top, CONTENT_W, score_h, fill=NAVY, rounded=True, radius=0.08)
    text_w = CONTENT_W * 0.5
    D.add_text(s, MARGIN + 0.22, score_top, text_w - 0.3, score_h, [
        ("Priorité = (impact × faisabilité) − prudence IA",
         dict(size=D.TYPE["small"], bold=True, color="#ffffff", line_spacing=1.15)),
        ("Support de discussion ORDINAL, pas une métrique calculée : à lire en paliers "
         "(fort / moyen / faible), jamais comme un nombre exact. Il rend la discussion "
         "explicite et classe les candidats — il ne remplace pas l'arbitrage humain.",
         dict(size=8, color="#c7cbe0", space_before=4, line_spacing=1.2)),
    ], anchor=MSO_ANCHOR.MIDDLE)

    # Jauge à points — pattern repris de l'autre template analysé
    # (analyse-template-alternatif.md §4) pour illustrer un score 1-5.
    gauge_x = MARGIN + text_w
    gauge_w = CONTENT_W - text_w
    D.add_text(s, gauge_x, score_top + 0.12, gauge_w - 0.15, 0.18, [
        ("SCORE ILLUSTRATIF · ORDINAL", dict(size=7, bold=True, color="#8891b3")),
    ])
    rows_top = score_top + 0.38
    row_h2 = (score_h - 0.38 - 0.1) / 3
    gauge_rows = [
        ("Impact", 4, "#ffffff"),
        ("Faisabilité", 3, ACCENT),
        ("Prudence IA", 1, SEVERITE[4]),
    ]
    for i, (label, score, color) in enumerate(gauge_rows):
        ry = rows_top + i * row_h2
        D.add_text(s, gauge_x, ry, 1.2, row_h2, [
            (label, dict(size=7, color="#c7cbe0")),
        ], anchor=MSO_ANCHOR.MIDDLE)
        dot_scale(s, gauge_x + 1.25, ry + row_h2 / 2 - 0.07, 5, score, color,
                  empty_color="#3a4568")
    return s


# ---------------------------------------------------------------- Besoins & douleurs
# Nouveau (restructuration 2026-07-22) : va PLUS LOIN que slide_personas (qui porte
# un irritant + une attente d'une ligne par persona). Ici chaque douleur est
# approfondie, dotée d'un signal/mesure qui la rend objectivable, et rattachée à
# une ou plusieurs familles de gaspillage — le pont direct vers slide_familles.
# Réutilise les couleurs d'accent persona (Infra=bleu, Utilisateur=teal,
# Management=or, Sponsor=violet). Rangées dimensionnées à leur contenu.
def slide_douleurs(prs):
    s = content_slide(prs, "Besoins & douleurs",
                       "Les douleurs des clients infra : mesurables, pas des plaintes",
                       color=D.PALETTE[2])
    D.add_text(s, MARGIN, CONTENT_TOP, CONTENT_W, 0.4, [
        ("Chaque douleur appartient à un persona et se range dans une famille de gaspillage "
         "— c'est ce qui la rend traitable plutôt que subie.",
         dict(size=8, color=MUTED, italic=True, line_spacing=1.2)),
    ])

    rows = [
        ("Infra & RUN", D.PALETTE[0],
         "RUN subi : les mêmes incidents reviennent et mobilisent les experts seniors, le BUILD "
         "est sacrifié à l'astreinte.",
         "Tickets récurrents/mois, part du temps en RUN non maîtrisé.",
         "RUN · Humain"),
        ("Utilisateur applicatif", D.PALETTE[5],
         "Aucun self-service ni parcours conçu : tout passe par un guichet, le contournement "
         "(shadow IT) va plus vite que la demande officielle.",
         "Taux de contournement, délai de mise à disposition.",
         "Flux · Cognitif"),
        ("Management", D.PALETTE[3],
         "« Expert devenu manager malgré lui » : reporting miroir et micromanagement "
         "compensatoire, faute de signal fiable sur le flux.",
         "Ratio temps reporting / temps résolution d'obstacles.",
         "Décisionnel · Humain"),
        ("Sponsor", D.PALETTE[4],
         "Pression à « mettre de l'IA » sans cas d'usage, peur d'une transformation cosmétique : "
         "beaucoup d'activité, peu de valeur démontrée.",
         "Valeur / capacité récupérée démontrée vs promise.",
         "Décisionnel · IA (gadget)"),
    ]

    # Colonnes (comme slide_architecture_si) : en-têtes une seule fois, puis 4
    # rangées à liseré = couleur persona, anchor MIDDLE. Hauteur de rangée
    # calée sur le contenu le plus long (colonne douleur), pas d'étirement.
    headers = ["PERSONA", "LA DOULEUR, APPROFONDIE", "SIGNAL / MESURE", "FAMILLE(S)"]
    col_widths = [1.3, 3.6, 1.85, 1.425]
    col_gap = 0.12
    xs = []
    cx = MARGIN
    for cw in col_widths:
        xs.append(cx)
        cx += cw + col_gap

    header_y = CONTENT_TOP + 0.45
    for x, w, label in zip(xs, col_widths, headers):
        D.add_text(s, x + (0.12 if x == xs[0] else 0), header_y, w, 0.2, [
            (label, dict(size=7, bold=True, color=MUTED)),
        ])

    # Hauteur de rangée calée sur le contenu, MAIS bornée pour toujours réserver
    # la ligne-pont du bas (« → slide suivante ») — sans ce plafond, 4 rangées de
    # 4 lignes remplissaient jusqu'au bas et escamotaient la ligne-pont.
    n = len(rows)
    row_gap = 0.12
    note_reserve = 0.5
    row_top = header_y + 0.26
    region_bot = CONTENT_BOTTOM - note_reserve
    size = 8
    lh = size * 1.2 / 72.0
    row_lines = max(
        max(_lignes(r[2], col_widths[1] - 0.2, size), _lignes(r[3], col_widths[2] - 0.15, size))
        for r in rows)
    row_h = min(row_lines * lh + 0.26, (region_bot - row_top - (n - 1) * row_gap) / n)
    for i, (nom, color, douleur, signal, famille) in enumerate(rows):
        y = row_top + i * (row_h + row_gap)
        D.add_rect(s, MARGIN, y, CONTENT_W, row_h, fill="#ffffff", line=LINE, line_w=0.75,
                   rounded=True, radius=0.08)
        D.add_rect(s, MARGIN, y, 0.06, row_h, fill=color, rounded=True, radius=0.5)
        D.add_text(s, xs[0] + 0.12, y + 0.08, col_widths[0] - 0.12, row_h - 0.16, [
            (nom, dict(size=D.TYPE["tiny"], bold=True, color=color, line_spacing=1.1)),
        ], anchor=MSO_ANCHOR.MIDDLE)
        D.add_text(s, xs[1], y + 0.08, col_widths[1] - 0.15, row_h - 0.16, [
            (douleur, dict(size=size, color=NAVY, line_spacing=1.2)),
        ], anchor=MSO_ANCHOR.MIDDLE)
        D.add_text(s, xs[2], y + 0.08, col_widths[2] - 0.12, row_h - 0.16, [
            (signal, dict(size=size, color=MUTED, italic=True, line_spacing=1.2)),
        ], anchor=MSO_ANCHOR.MIDDLE)
        D.add_text(s, xs[3], y + 0.08, col_widths[3], row_h - 0.16, [
            (famille, dict(size=size, bold=True, color=color, line_spacing=1.2)),
        ], anchor=MSO_ANCHOR.MIDDLE)

    note_top = row_top + len(rows) * row_h + (len(rows) - 1) * row_gap + 0.14
    note_h = min(0.4, CONTENT_BOTTOM - note_top)
    if note_h > 0.12:
        D.add_text(s, MARGIN, note_top, CONTENT_W, note_h, [
            ("Ces douleurs se rangent en 8 familles de gaspillage → slide suivante.",
             dict(size=8, bold=True, color=D.PALETTE[2], line_spacing=1.15)),
        ])
    return s


# ---------------------------------------------------------------- slide 9
def slide_team_topologies(prs):
    s = content_slide(prs, "Proposition", "La cible IAP est une Platform Team — agents IA compris", color=D.PALETTE[1])
    types = [
        ("Stream-aligned", D.PALETTE[0], "Flux de valeur métier continu",
         "Équipes applicatives clientes de la plateforme infra"),
        ("Platform", D.PALETTE[1], "Capacités en self-service (X-as-a-Service)",
         "La cible même de la transformation IAP"),
        ("Enabling", D.PALETTE[3], "Montée en compétence temporaire",
         "Posture du coach BMAD IAP — jamais permanente"),
        ("Complicated-subsystem", D.PALETTE[4], "Expertise pointue, compétences rares",
         "Un vrai sous-système complexe, pas un produit plateforme classique"),
    ]
    n = 4
    card_h = 1.55
    top0 = CONTENT_TOP + 0.05
    for i, (titre, color, role, lecture) in enumerate(types):
        x, w = col_x(i, n)
        D.add_card(s, x, top0, w, card_h, color)
        pad = 0.16
        D.add_text(s, x + pad, top0 + 0.14, w - 2 * pad, 0.45, [
            (titre, dict(size=8, bold=True, color=color, line_spacing=1.05)),
        ])
        D.add_text(s, x + pad, top0 + 0.58, w - 2 * pad, 0.4, [
            (role, dict(size=8, color=NAVY, line_spacing=1.15)),
        ])
        D.add_text(s, x + pad, top0 + 1.0, w - 2 * pad, card_h - 1.1, [
            (lecture, dict(size=8, color=MUTED, italic=True, line_spacing=1.15)),
        ])

    note_top = top0 + card_h + 0.18
    note_h = min(1.5, CONTENT_BOTTOM - note_top)
    D.add_rect(s, MARGIN, note_top, CONTENT_W, note_h, fill=TRACK, rounded=True, radius=0.08)
    D.add_text(s, MARGIN + 0.22, note_top, CONTENT_W - 0.44, note_h, [
        ("Extension — les agents IA comme coéquipiers, et leur mise en œuvre (v1.7)",
         dict(size=D.TYPE["tiny"], bold=True, color=NAVY)),
        ("Un agent peut être membre d'une Stream-aligned team ou capacité exposée par la "
         "Platform Team. Aux 3 modes d'interaction Team Topologies — Collaboration, "
         "X-as-a-Service, Facilitating — s'ajoute un 4ᵉ candidat : Supervision. "
         "L'adoption suit la trajectoire "
         "Coach → Délégué (assisté → supervisé → délégué) : mandat écrit (ce que l'agent "
         "décide seul / ce qui escalade / qui répond de ses erreurs) avant tout palier "
         "au-delà de l'assisté — jamais un usage qui dérive à l'implicite.",
         dict(size=8, color=NAVY, space_before=3, line_spacing=1.25)),
    ], anchor=MSO_ANCHOR.MIDDLE)
    return s


# ---------------------------------------------------------------- slide 8
# Formes inspirées des slides d'exemple du template lui-même (« Notre
# approche ») : badges circulaires connectés par une ligne, chip de durée,
# description centrée sous chaque étape.
# Nouveau (v2.3) : le schéma de fonctionnement du §Trajectoire de
# bmad-iap-cadrage.md n'avait jusqu'ici qu'un résumé en une ligne dans
# slide_trajectoire (phase ①, "= Schéma de fonctionnement déjà cadré") —
# jamais sa propre slide. Reprend la lecture verticale du schéma ASCII
# source : bandeau Gate IA transversal, 4 colonnes du pipeline, bandeau
# iap-risk-reviewer, bandeau boucle de réévaluation.
def slide_schema_fonctionnement(prs):
    # v2.5 (chantier ④) : déplacée de la Proposition vers la Démarche — c'est
    # « comment la mission tourne », pas ce qu'on propose.
    s = content_slide(prs, "Démarche",
                       "La Gate IA s'applique à chaque étape, de la collecte à la boucle de réévaluation",
                       color=D.PALETTE[3])
    D.add_text(s, MARGIN, CONTENT_TOP, CONTENT_W, 0.26, [
        ("Deux sources de collecte convergent vers un diagnostic structuré ; une boucle de "
         "réévaluation referme le cycle.", dict(size=8, color=MUTED, italic=True, line_spacing=1.1)),
    ])

    band_top = CONTENT_TOP + 0.3
    band_h = 0.36
    D.add_rect(s, MARGIN, band_top, CONTENT_W, band_h, fill=NAVY, rounded=True, radius=0.12)
    D.add_text(s, MARGIN + 0.2, band_top, CONTENT_W - 0.4, band_h, [
        ("GATE IA & CONFIDENTIALITÉ — checkpoint humain, transversal à chaque étape qui invoque un LLM",
         dict(size=8, bold=True, color="#ffffff", line_spacing=1.1)),
    ], anchor=MSO_ANCHOR.MIDDLE)

    etapes = [
        ("COLLECTE", D.PALETTE[0],
         ["Interviews par persona (trame / thème / question)",
          "Import outils : ServiceNow/Jira/CMDB si accès"]),
        ("DIAGNOSTIC", D.PALETTE[4],
         ["Synthèse par thème puis synthèse globale",
          "Registre de gaspillage (tags CONFIRMÉ/DÉDUIT/INCERTAIN)"]),
        ("CONCEPTION", D.PALETTE[3],
         ["Définition produit (+ cible MVP)",
          "Operating model + traitement du gaspillage (décisions actées)"]),
        ("RESTITUTION", D.PALETTE[1],
         ["Deck exécutif : axes valeur/complexité",
          "+ radar de maturité"]),
    ]
    n = len(etapes)
    col_top = band_top + band_h + 0.16
    card_h = 1.85
    for i, (titre, color, lignes) in enumerate(etapes):
        x, w = col_x(i, n)
        D.add_rect(s, x, col_top, w, 0.05, fill=color)
        D.add_text(s, x, col_top + 0.1, w, 0.26, [
            (titre, dict(size=8, bold=True, color=color, align=PP_ALIGN.CENTER)),
        ], align=PP_ALIGN.CENTER)
        card_y = col_top + 0.4
        D.add_rect(s, x, card_y, w, card_h - 0.4, fill=TRACK, rounded=True, radius=0.08)
        lignes_fmt = [(f"·  {l}", dict(size=7, color=NAVY, space_after=4, line_spacing=1.2)) for l in lignes]
        D.add_text(s, x + 0.1, card_y + 0.12, w - 0.2, card_h - 0.4 - 0.24, lignes_fmt,
                   anchor=MSO_ANCHOR.MIDDLE)
        if i < n - 1:
            D.add_text(s, x + w, col_top + 0.08, GAP, 0.26, [
                ("→", dict(size=10, bold=True, color=MUTED, align=PP_ALIGN.CENTER)),
            ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

    reviewer_top = col_top + card_h + 0.16
    reviewer_h = 0.4
    D.add_rect(s, MARGIN, reviewer_top, CONTENT_W, reviewer_h, fill="#ffffff", line=LINE,
               rounded=True, radius=0.1)
    D.add_text(s, MARGIN + 0.2, reviewer_top, CONTENT_W - 0.4, reviewer_h, [
        ("iap-risk-reviewer — lecture seule, challenge Product definition / Operating model → Deck exécutif",
         dict(size=7.5, italic=True, color=MUTED, line_spacing=1.1)),
    ], anchor=MSO_ANCHOR.MIDDLE)

    loop_top = reviewer_top + reviewer_h + 0.14
    loop_h = min(0.55, CONTENT_BOTTOM - loop_top)
    D.add_rect(s, MARGIN, loop_top, CONTENT_W, loop_h, fill=D.PALETTE[2], rounded=True, radius=0.1)
    D.add_text(s, MARGIN + 0.2, loop_top, CONTENT_W - 0.4, loop_h, [
        ("⟲ Boucle de réévaluation — iap-re-assessment, T+6-12 mois, alimente la bibliothèque de REX, "
         "reboucle vers la Collecte", dict(size=8, bold=False, color="#ffffff", line_spacing=1.15)),
    ], anchor=MSO_ANCHOR.MIDDLE)
    return s


# Fusion v2.5 (chantier ①) : slide_trajectoire et slide_schema_bout_en_bout
# déroulaient la même trame ①②③⟲ sur deux slides — fusionnées ici. La ligne de
# badges + durées + actions clés (de l'ancienne trajectoire) est enrichie du
# LIVRABLE-CLÉ par phase en une ligne (l'apport de la vue bout-en-bout — les
# NOMS seulement ; le détail des 4 profils de deck reste à slide_livrables_ppt).
def slide_trajectoire(prs):
    s = content_slide(prs, "Démarche",
                       "Trois temps et une boucle — chaque phase produit son livrable de décision",
                       color=D.PALETTE[3])
    phases = [
        ("①", "Assessment flash", "1–2 sem.", D.PALETTE[0],
         "= Schéma de fonctionnement déjà cadré (Collecte → Diagnostic → Conception → Restitution).",
         "Deck exécutif de restitution"),
        ("②", "Premier déploiement", "4–5 sem.", D.PALETTE[3],
         "1-2 équipes pilotes, mode Coach dominant. Piste agent IA (si retenue) : qualifier, cadrer, mandater.",
         "Deck de plan de déploiement · export markdown (1re version)"),
        ("③", "Implémentation itérative", "→ T+6-12 mois", D.PALETTE[1],
         "Généralisation équipe par équipe, bascule Coach → Délégué. Piste agent IA : supervisé puis délégué.",
         "Deck de comité de pilotage (périodique)"),
        ("⟲", "Boucle de réévaluation", "T+6-12 mois", D.PALETTE[2],
         "iap-re-assessment reboucle vers la Collecte — alimente la bibliothèque de REX.",
         "Deck de bilan / ré-évaluation · markdown amendé"),
    ]
    n = len(phases)
    badge_d = 0.55
    top0 = CONTENT_TOP + 0.1
    line_y = top0 + badge_d / 2 - 0.012
    D.add_rect(s, MARGIN + badge_d / 2, line_y, CONTENT_W - badge_d, 0.024, fill=LINE)
    _, wcol = col_x(0, n)
    desc_h = max(_lignes(p[4], wcol - 0.1, 7) for p in phases) * (7 * 1.2 / 72.0) + 0.05
    livr_h = max(_lignes(p[5], wcol - 0.2, 7.5) for p in phases) * (7.5 * 1.2 / 72.0) + 0.24
    for i, (sym, titre, duree, color, desc, livrable) in enumerate(phases):
        x, w = col_x(i, n)
        cx = x + w / 2 - badge_d / 2
        D.add_rect(s, cx, top0, badge_d, badge_d, fill=color, rounded=True, radius=0.5)
        D.add_text(s, cx, top0, badge_d, badge_d, [
            # bold=False pour "⟲" : sa variante grasse manque dans la police du
            # template (rendu LibreOffice = case vide) — ①②③ n'ont pas ce problème.
            (sym, dict(size=16, bold=(sym != "⟲"), color="#ffffff", align=PP_ALIGN.CENTER))
        ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        ty = top0 + badge_d + 0.12
        D.add_text(s, x, ty, w, 0.35, [
            (titre, dict(size=8, bold=True, color=NAVY, align=PP_ALIGN.CENTER, line_spacing=1.05)),
        ], align=PP_ALIGN.CENTER)
        chip_y = ty + 0.36
        chip(s, x + w / 2 - 0.55, chip_y, 1.1, 0.24, duree, color, size=7)
        desc_y = chip_y + 0.34
        D.add_text(s, x + 0.05, desc_y, w - 0.1, desc_h, [
            (desc, dict(size=7, color=MUTED, align=PP_ALIGN.CENTER, line_spacing=1.2)),
        ], align=PP_ALIGN.CENTER)
        # Livrable-clé (fusion bout-en-bout) : le NOM du livrable, encadré,
        # au pied de chaque colonne — le détail vit dans slide_livrables_ppt.
        livr_y = desc_y + desc_h + 0.10
        D.add_rect(s, x, livr_y, w, livr_h, fill=TRACK, rounded=True, radius=0.1)
        D.add_text(s, x + 0.1, livr_y, w - 0.2, livr_h, [
            ("LIVRABLE-CLÉ", dict(size=6.5, bold=True, color=MUTED, align=PP_ALIGN.CENTER)),
            (livrable, dict(size=7.5, bold=True, color=color, space_before=2,
                            align=PP_ALIGN.CENTER, line_spacing=1.15)),
        ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

    note_top = top0 + badge_d + 0.12 + 0.36 + 0.34 + desc_h + 0.10 + livr_h + 0.18
    note_h = min(1.05, CONTENT_BOTTOM - note_top)
    D.add_rect(s, MARGIN, note_top, CONTENT_W, note_h, fill=TRACK, rounded=True, radius=0.08)
    D.add_text(s, MARGIN + 0.2, note_top, CONTENT_W - 0.4, note_h, [
        ("Bifurcation avec/sans agents IA déployés", dict(size=8, bold=True, color=NAVY)),
        ("Le tronc commun ①→②→③→⟲ ne change pas de structure — la piste agent IA (si retenue) "
         "se greffe sur ②/③ via la démarche d'accompagnement en 5 phases déjà cadrée, plutôt "
         "que d'être un chemin séparé à maintenir. Owner proposé (non tranché) : "
         "iap-operating-model-architect + iap-change-coach sur le volet humain.",
         dict(size=7, color=NAVY, space_before=3, line_spacing=1.25)),
    ], anchor=MSO_ANCHOR.MIDDLE)
    return s


# Nouveau (v2.4) : le fil humain de la trajectoire (§Accompagnement de
# l'humain dans la trajectoire, transposé de l'offre SCALE). Décline la même
# trame ①②③⟲ que slide_trajectoire côté personnes — pattern 14 du catalogue
# deck-design-library (« processus en 4 étapes numérotées, colonnes de détail
# dans UNE carte », fil rouge transversal repérable en diagonale : ici
# iap-change-coach répété en pied de chaque colonne) + badge chevauchant le
# bord haut de la carte (pattern 10). L'accroche Kotter (70 %) est LE seul
# élément en aplat plein de la slide (« un sur N en accent »).
def slide_fil_humain(prs):
    s = content_slide(prs, "Démarche",
                       "La trajectoire accompagne aussi les personnes, de bout en bout",
                       color=D.PALETTE[3])

    # --- Accroche argumentaire : le chiffre Kotter en bloc accent + intro.
    stat_w, stat_h = 1.05, 0.62
    strip_top = CONTENT_TOP + 0.02
    D.add_rect(s, MARGIN, strip_top, stat_w, stat_h, fill=NAVY, rounded=True, radius=0.12)
    D.add_text(s, MARGIN, strip_top, stat_w, stat_h, [
        ("70 %", dict(size=18, bold=True, color="#ffffff", align=PP_ALIGN.CENTER)),
    ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    tx = MARGIN + stat_w + 0.18
    tw = CONTENT_W - stat_w - 0.18
    D.add_text(s, tx, strip_top, tw, stat_h, [
        ("des transformations sont inachevées ou échouent parce que les facteurs humains "
         "et culturels sont mal pris en compte (Kotter, Harvard Business Review).",
         dict(size=8, bold=True, color=NAVY, line_spacing=1.15)),
        # ⟲ en run NON gras (paragraphe italic non-bold) — cf. _GLYPHES_SANS_GRAS.
        ("Le fil humain suit la même trame ①②③⟲ que la trajectoire — un fil dans les "
         "phases, pas un stream séparé, porté par iap-change-coach de bout en bout.",
         dict(size=7.5, italic=True, color=MUTED, space_before=3, line_spacing=1.15)),
    ], anchor=MSO_ANCHOR.MIDDLE)

    phases = [
        ("①", "ASSESSMENT FLASH", "Engager", D.PALETTE[0],
         "L'engagement personnel du sponsor est testé dès l'intake, avant signature ; "
         "les interviews écoutent les tensions ; la restitution revient aux interviewés, "
         "pas au seul sponsor.",
         "iap-change-coach · iap-intake"),
        ("②", "PREMIER DÉPLOIEMENT", "Expérimenter", D.PALETTE[3],
         "Équipes pilotes volontaires, jamais désignées d'office ; formation sur les cas "
         "réels de l'équipe — « pas de formation sans coaching ».",
         "iap-change-coach · équipe pilote"),
        ("③", "IMPLÉMENTATION ITÉRATIVE", "Outiller & relayer", D.PALETTE[1],
         "Les résistances sont un signal à écouter ; communauté de managers (N+1/N+2 "
         "embarqués) ; relais internes formés — le consultant se rend dispensable.",
         "iap-change-coach · operating-model-architect"),
        ("⟲", "BOUCLE DE RÉÉVALUATION", "Mesurer", D.PALETTE[2],
         "Satisfaction et adhésion mesurées au même instrument à T0 et à la réévaluation — "
         "le delta humain se lit à côté du delta de maturité.",
         "iap-change-coach · metrics-sre-finops-lead"),
    ]
    n = len(phases)
    badge_d = 0.45
    pad = 0.16
    _, col_w = col_x(0, n, gap=0)
    usable = col_w - 2 * pad

    # Hauteurs dérivées du CONTENU (jamais « jusqu'à CONTENT_BOTTOM ») —
    # défaut récurrent « panneau étiré vide », cf. slide_livrables_ppt.
    body_lines = max(_lignes(p[4], usable, 7.5) for p in phases)
    body_h = body_lines * (7.5 * 1.25 / 72.0) + 0.04
    owner_lines = max(_lignes(p[5], usable, 6.5) for p in phases)
    owner_h = owner_lines * (6.5 * 1.2 / 72.0) + 0.03
    # card_h relatif au bord haut de la carte : demi-badge + en-têtes + corps
    # + séparateur + owner + respiration basse.
    card_h = badge_d / 2 + 0.10 + 0.44 + body_h + 0.14 + owner_h + 0.14
    band_h = 0.60
    group_h = badge_d / 2 + card_h + 0.18 + band_h
    region_top = strip_top + stat_h + 0.16
    top1 = region_top + min(0.30, max(0.0, (CONTENT_BOTTOM - region_top - group_h) / 2))
    card_top = top1 + badge_d / 2

    D.add_rect(s, MARGIN, card_top, CONTENT_W, card_h, fill="#ffffff",
               line=LINE, line_w=0.75, rounded=True, radius=0.06)
    for i, (sym, phase, verbe, color, desc, owner) in enumerate(phases):
        x, w = col_x(i, n, gap=0)
        if i > 0:  # séparateurs fins internes — UNE carte, pas 4 (pattern 14)
            D.add_rect(s, x, card_top + 0.14, 0.012, card_h - 0.28, fill=LINE)
        cx = x + w / 2 - badge_d / 2
        D.add_rect(s, cx, top1, badge_d, badge_d, fill=color, rounded=True, radius=0.5)
        D.add_text(s, cx, top1, badge_d, badge_d, [
            # bold=False pour "⟲" : variante grasse absente de la police du
            # template (tofu au rendu) — même correctif que slide_trajectoire.
            (sym, dict(size=13, bold=(sym != "⟲"), color="#ffffff", align=PP_ALIGN.CENTER)),
        ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        head_y = top1 + badge_d + 0.10
        D.add_text(s, x + pad, head_y, usable, 0.22, [
            (verbe, dict(size=9, bold=True, color=color, align=PP_ALIGN.CENTER)),
        ], align=PP_ALIGN.CENTER)
        D.add_text(s, x + pad, head_y + 0.22, usable, 0.16, [
            (phase, dict(size=6, bold=True, color=MUTED, align=PP_ALIGN.CENTER)),
        ], align=PP_ALIGN.CENTER)
        body_y = head_y + 0.44
        D.add_text(s, x + pad, body_y, usable, body_h, [
            (desc, dict(size=7.5, color=NAVY, line_spacing=1.25)),
        ])
        sep_y = body_y + body_h + 0.06
        D.add_rect(s, x + pad, sep_y, usable, 0.012, fill=LINE)
        D.add_text(s, x + pad, sep_y + 0.06, usable, owner_h, [
            (owner, dict(size=6.5, bold=True, color=color, line_spacing=1.2)),
        ])

    band_top = card_top + card_h + 0.18
    D.add_rect(s, MARGIN, band_top, CONTENT_W, band_h, fill=TRACK, rounded=True, radius=0.08)
    D.add_text(s, MARGIN + 0.2, band_top, CONTENT_W - 0.4, band_h, [
        ("Ce que ça ne crée pas", dict(size=8, bold=True, color=NAVY)),
        ("Pas de nouvel agent (iap-change-coach porte le fil de bout en bout), pas de phase "
         "en plus — et jamais d'évaluation individuelle des personnes (déontologie du "
         "consultant, transposée telle quelle).",
         dict(size=7, color=NAVY, space_before=3, line_spacing=1.25)),
    ], anchor=MSO_ANCHOR.MIDDLE)
    return s


# --- Nouveau (v2.6, point ②) : les activités humaines de la démarche, en DEUX
# registres — outillées par IAP vs purement humaines (sans l'outil). Contenu
# ancré dans docs/Import/notes-extraction-scale.md (micro-lancement sponsor,
# ateliers collaboratifs, communauté de managers N+1/N+2, coaching sous
# déontologie, relais internes, présence dégressive) et docs/bmad-iap-cadrage.md
# §Accompagnement de l'humain dans la trajectoire (v2.4) — rien d'inventé.
# Pattern 11 du catalogue deck-design-library (« trajectoire à 4 phases en
# colonnes × lignes de catégorie ») : colonnes = les temps ①②③⟲ (mêmes badges
# et couleurs que slide_trajectoire/slide_fil_humain), lignes = les deux
# registres — bande « avec IAP » teintée cyan pâle (la teinte distingue une
# famille sans lire l'étiquette, pattern 8), bande « sans IAP » blanche à
# contour. Cellules ancrées MIDDLE dans leur bande et hauteurs dérivées du
# contenu (défaut récurrent « panneau sur-étiré », évité d'office).
def slide_activites_humaines(prs):
    s = content_slide(prs, "Démarche",
                       "IAP outille une partie du fil humain — le reste est de la présence de consultant",
                       color=D.PALETTE[3])
    D.add_text(s, MARGIN, CONTENT_TOP, CONTENT_W, 0.28, [
        ("Les mêmes quatre temps que la trajectoire : en haut, ce que le consultant fait "
         "avec le module ; en bas, ce qu'il fait sans lui.",
         dict(size=8, color=MUTED, italic=True, line_spacing=1.15)),
    ])

    phases = [
        ("①", "Assessment flash", D.PALETTE[0]),
        ("②", "Premier déploiement", D.PALETTE[3]),
        ("③", "Implémentation itérative", D.PALETTE[1]),
        ("⟲", "Boucle de réévaluation", D.PALETTE[2]),
    ]
    avec_iap = [
        ["Trames d'interview par persona, versées à la Collecte",
         "La restitution s'appuie sur le deck exécutif généré",
         "Sondage humain de référence (T0)"],
        ["Deck de plan de déploiement · export markdown (1re version)"],
        ["Comité de pilotage : deck périodique, santé humaine incluse"],
        ["KPI humain rejoué au même instrument : delta T0 → réévaluation",
         "Export markdown amendé · deck de bilan"],
    ]
    sans_iap = [
        ["Le sponsor présente lui-même l'ambition (micro-lancement)",
         "Restitution-embarquement : feedback des interviewés"],
        ["Équipes pilotes volontaires, ateliers co-construits",
         "Formation sur les cas réels — pas de formation sans coaching"],
        ["Communauté de managers, N+1/N+2 embarqués",
         "Coaching individuel sous déontologie · relais internes formés"],
        ["Présence dégressive : le consultant se rend dispensable",
         "Lecture du delta humain avec les équipes"],
    ]

    label_w = 1.05
    grid_x0 = MARGIN + label_w + 0.12
    grid_w = BORD_DROIT - grid_x0
    n = 4
    col_gap = 0.08
    col_w = (grid_w - (n - 1) * col_gap) / n

    def _col_px(i):
        return grid_x0 + i * (col_w + col_gap)

    # En-têtes de phase : badge rond + nom, centrés par colonne (⟲ jamais en
    # gras — cf. _GLYPHES_SANS_GRAS, même correctif que slide_trajectoire).
    head_top = CONTENT_TOP + 0.34
    badge_d = 0.26
    for i, (sym, nom, color) in enumerate(phases):
        x = _col_px(i)
        cx = x + col_w / 2 - badge_d / 2
        D.add_rect(s, cx, head_top, badge_d, badge_d, fill=color, rounded=True, radius=0.5)
        D.add_text(s, cx, head_top, badge_d, badge_d, [
            (sym, dict(size=9, bold=(sym != "⟲"), color="#ffffff", align=PP_ALIGN.CENTER)),
        ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        D.add_text(s, x, head_top + badge_d + 0.02, col_w, 0.16, [
            (nom, dict(size=6.5, bold=True, color=color, align=PP_ALIGN.CENTER)),
        ], align=PP_ALIGN.CENTER)

    # Hauteur de bande dérivée du CONTENU (colonne la plus fournie).
    cell_usable = col_w - 0.16
    line_h = 7 * 1.2 / 72.0

    def _band_h(cells):
        return max(sum(_lignes(t, cell_usable, 7) for t in c) * line_h
                   + (len(c) - 1) * 0.06 for c in cells) + 0.24

    bands_top = head_top + badge_d + 0.24
    note_h = 0.34
    hA = _band_h(avec_iap)
    hB = _band_h(sans_iap)
    # Le mou vertical restant se répartit DANS les bandes (cellules centrées),
    # pas en vide sous la grille.
    slack = max(0.0, (CONTENT_BOTTOM - note_h - 0.12) - (bands_top + hA + 0.10 + hB))
    hA += slack / 2
    hB += slack / 2

    registres = [
        (bands_top, hA, "AVEC IAP", "outillé par le module", "#E1FDFA", None, avec_iap),
        (bands_top + hA + 0.10, hB, "SANS IAP", "présence du consultant", "#ffffff", LINE, sans_iap),
    ]
    for top, h, label, sous, fill, line, cells in registres:
        D.add_rect(s, MARGIN, top, CONTENT_W, h, fill=fill, line=line, line_w=0.75,
                   rounded=True, radius=0.06)
        D.add_text(s, MARGIN + 0.12, top, label_w - 0.12, h, [
            (label, dict(size=8, bold=True, color=NAVY, line_spacing=1.1)),
            (sous, dict(size=6.5, color=MUTED, italic=True, space_before=2, line_spacing=1.1)),
        ], anchor=MSO_ANCHOR.MIDDLE)
        D.add_rect(s, grid_x0 - 0.10, top + 0.10, 0.012, h - 0.20, fill=LINE)
        for i, items in enumerate(cells):
            x = _col_px(i)
            if i > 0:  # séparateurs fins (pattern 11 : la grille sans le tableau)
                D.add_rect(s, x - col_gap / 2, top + 0.10, 0.012, h - 0.20, fill=LINE)
            lignes_fmt = [(t, dict(size=7, color=NAVY, line_spacing=1.2,
                                   space_before=(4 if j else 0)))
                          for j, t in enumerate(items)]
            D.add_text(s, x + 0.08, top + 0.08, col_w - 0.16, h - 0.16, lignes_fmt,
                       anchor=MSO_ANCHOR.MIDDLE)

    note_top = CONTENT_BOTTOM - note_h
    D.add_text(s, MARGIN, note_top, CONTENT_W, note_h, [
        ("La rangée du bas ne s'outille pas : c'est la présence du consultant — dégressive, "
         "jusqu'aux relais internes qui portent le modèle après la mission (offre SCALE, transposée).",
         dict(size=8, color=MUTED, italic=True, line_spacing=1.2)),
    ], anchor=MSO_ANCHOR.BOTTOM)
    return s


# --- Nouveau (2026-09-01) : conditions de réussite et non-engagement. Le deck
# ne disait nulle part ce que la mission EXIGE du client ni ce qui la fait
# échouer — un sponsor destinataire ne pouvait pas savoir à quoi il s'engage.
# Prolonge (sans la répéter) la phrase du fil humain « testé dès l'intake ».
# Forme retenue (deck-design-library) : chaîne verticale de conditions reliées
# par un connecteur + encart de mise en exergue latéral (pattern 6) — la
# colonne porte le processus, le panneau porte le message, ici l'issue
# NÉGATIVE. Bandeau transverse en pied pour le critère de sortie (l.949).
# Matière : cadrage l.934 (non-engagement), l.951 (test à l'intake, RH),
# l.172 (anti-patterns), l.455 (deskilling-risk), l.667 (management-posture-risk).
def slide_conditions_reussite(prs):
    s = content_slide(prs, "Démarche",
                       "Sans ces quatre conditions, la mission ne s'engage pas",
                       color=D.PALETTE[3])
    couleur = D.PALETTE[3]

    def lh(pt):
        return pt * 1.25 / 72.0

    lead = ("Le test d'engagement porté par le fil humain n'a de valeur que s'il peut conclure "
            "« non » : voici ce qu'il vérifie avant signature — et ce que son échec déclenche.")
    lead_h = _lignes(lead, CONTENT_W, 9.5) * lh(9.5) + 0.06
    D.add_text(s, MARGIN, CONTENT_TOP, CONTENT_W, lead_h, [
        (lead, dict(size=9.5, color=NAVY, italic=True, line_spacing=1.25)),
    ])

    # --- Bandeau transverse (critère de sortie, l.949) : dimensionné d'abord,
    # la région centrale prend le reste — jamais l'inverse.
    sortie = ("L'équipe et ses relais tiennent-ils le modèle une période sans le consultant ? "
              "Le consultant se rend dispensable : il n'évalue jamais les personnes, ne fait "
              "pas le reporting à leur place et ne s'installe pas en intermédiaire permanent "
              "entre l'équipe et le sponsor.")
    band_pad = 0.12
    band_h = 2 * band_pad + lh(7) + 0.04 + _lignes(sortie, CONTENT_W - 0.28, 7) * lh(7) + 0.03
    band_top = CONTENT_BOTTOM - band_h

    region_top = CONTENT_TOP + lead_h + 0.10
    region_h = band_top - 0.16 - region_top

    # --- Colonne gauche : les 4 conditions, en chaîne -----------------------
    conditions = [
        ("Un sponsor qui porte la cible",
         "La transformation ne va pas plus loin que ce que le sponsor peut porter : c'est son "
         "engagement personnel qui se teste, pas son budget."),
        ("Des équipes réellement disponibles",
         "Interviews et ateliers supposent du temps réservé — la disponibilité réelle se "
         "vérifie à l'intake, jamais en cours de mission."),
        ("Une RH embarquée sur rôles et évaluation",
         "« Frein ou principal accélérateur » : coacher une posture que les grilles "
         "d'évaluation punissent revient à ramer contre le système."),
        ("Un processus documenté avant tout agent",
         "Sinon on fige une pratique mal définie dans du code — préalable non négociable de "
         "la doctrine d'automatisation."),
    ]
    gauche_w = 3.65
    fil_x = MARGIN + 0.13
    badge_d = 0.26
    card_x = MARGIN + 0.26
    card_w = MARGIN + gauche_w - card_x
    card_pad = 0.14
    card_usable = card_w - 0.08 - 2 * card_pad
    card_gap = 0.09
    card_h = (region_h - (len(conditions) - 1) * card_gap) / len(conditions)

    # Connecteur continu : une seule ligne, pas de flèches (pattern 6).
    D.add_rect(s, fil_x - 0.01, region_top + card_h / 2, 0.02,
               (len(conditions) - 1) * (card_h + card_gap), fill=LINE)
    for i, (titre, corps) in enumerate(conditions):
        y = region_top + i * (card_h + card_gap)
        D.add_card(s, card_x, y, card_w, card_h, couleur)
        D.add_rect(s, fil_x - badge_d / 2, y + card_h / 2 - badge_d / 2, badge_d, badge_d,
                   fill=couleur if i == 0 else "#ffffff",
                   line=None if i == 0 else couleur, line_w=1.0, rounded=True, radius=0.5)
        D.add_text(s, fil_x - badge_d / 2, y + card_h / 2 - badge_d / 2, badge_d, badge_d, [
            (str(i + 1), dict(size=8, bold=True,
                              color="#ffffff" if i == 0 else couleur,
                              align=PP_ALIGN.CENTER)),
        ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        D.add_text(s, card_x + 0.08 + card_pad, y, card_usable, card_h, [
            (titre, dict(size=8, bold=True, color=couleur, line_spacing=1.05)),
            (corps, dict(size=7, color=NAVY, space_before=3, line_spacing=1.25)),
        ], anchor=MSO_ANCHOR.MIDDLE)

    # --- Panneau droit : l'issue négative (le seul aplat plein de la slide) --
    pan_x = MARGIN + gauche_w + 0.24
    pan_w = BORD_DROIT - pan_x
    pad = 0.16
    tw = pan_w - 2 * pad

    chip_l, chip_h = 1.42, 0.22
    accroche = ("Un sponsor qui achète un audit mais ne portera pas la cible est un critère "
                "de non-engagement, pas un aléa.")
    nuance = ("Le refus de mission reste un point ouvert du cadrage ; le test d'engagement, "
              "lui, est acté — il conditionne la signature.")
    refus = ["Automatiser un processus mal conçu",
             "Livrer une plateforme techniquement bonne mais peu adoptée",
             "Séparer transformation organisationnelle et technique"]
    risques = [
        ("deskilling-risk — perte de la capacité tacite",
         "L'équipe saurait-elle reprendre la main une semaine sans l'agent ?"),
        ("management-posture-risk — incitations RH contraires",
         "Qu'est-ce qui, dans vos grilles d'évaluation, récompense encore le comportement "
         "qu'on vient de décourager ?"),
    ]
    b1_h = (chip_h + 0.06 + _lignes(accroche, tw, 9) * lh(9) + 0.04
            + _lignes(nuance, tw, 7.5) * lh(7.5))
    b2_h = (lh(7) + 0.03 + sum(_lignes(r, tw - 0.14, 7) for r in refus) * lh(7)
            + (len(refus) - 1) * 0.03)
    b3_h = (lh(7) + 0.03
            + sum(_lignes(a, tw, 7) + _lignes(b, tw, 7) for a, b in risques) * lh(7)
            + (len(risques) - 1) * 0.05)

    # Panneau dimensionné à SON contenu (le mou part dans les interlignes de
    # blocs, jamais en vide au pied du panneau).
    gap_int = 0.14
    contenu_h = 2 * pad + b1_h + b2_h + b3_h + 2 * gap_int
    slack = region_h - contenu_h
    if slack > 0:
        gap_int += min(slack / 2.0, 0.18)
        contenu_h = 2 * pad + b1_h + b2_h + b3_h + 2 * gap_int
    pan_h = min(region_h, contenu_h)
    pan_top = region_top + max(0.0, (region_h - pan_h) / 2.0)

    D.add_rect(s, pan_x, pan_top, pan_w, pan_h, fill=NAVY, rounded=True, radius=0.08)
    y = pan_top + pad
    D.add_rect(s, pan_x + pad, y, chip_l, chip_h, fill=D.PALETTE[2], rounded=True, radius=0.5)
    D.add_text(s, pan_x + pad, y, chip_l, chip_h, [
        ("NON-ENGAGEMENT", dict(size=7, bold=True, color="#ffffff", align=PP_ALIGN.CENTER)),
    ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    D.add_text(s, pan_x + pad, y + chip_h + 0.06, tw, b1_h - chip_h - 0.06, [
        (accroche, dict(size=9, bold=True, color="#ffffff", line_spacing=1.15)),
        (nuance, dict(size=7.5, color="#c7cbe0", space_before=3, line_spacing=1.25)),
    ])

    y += b1_h + gap_int
    D.add_rect(s, pan_x + pad, y - gap_int / 2, tw, 0.012, fill="#3a4568")
    D.add_text(s, pan_x + pad, y, tw, lh(7), [
        ("CE QU'ON REFUSE DE FAIRE — LES ANTI-PATTERNS DU CADRAGE",
         dict(size=7, bold=True, color="#8891b3")),
    ])
    ry = y + lh(7) + 0.03
    for r in refus:
        n = _lignes(r, tw - 0.14, 7)
        D.add_dot(s, pan_x + pad + 0.02, ry + 0.04, 0.05, ACCENT)
        D.add_text(s, pan_x + pad + 0.14, ry, tw - 0.14, n * lh(7), [
            (r, dict(size=7, color="#ffffff", line_spacing=1.25)),
        ])
        ry += n * lh(7) + 0.03

    y += b2_h + gap_int
    D.add_rect(s, pan_x + pad, y - gap_int / 2, tw, 0.012, fill="#3a4568")
    D.add_text(s, pan_x + pad, y, tw, lh(7), [
        ("CE QU'ON NE RÉSOUT PAS — MAIS QU'ON CONSIGNE ET QU'ON POSE",
         dict(size=7, bold=True, color="#8891b3")),
    ])
    ry = y + lh(7) + 0.03
    for nom, question in risques:
        n = _lignes(nom, tw, 7) + _lignes(question, tw, 7)
        D.add_text(s, pan_x + pad, ry, tw, n * lh(7), [
            (nom, dict(size=7, bold=True, color=D.PALETTE[3], line_spacing=1.25)),
            (question, dict(size=7, color="#c7cbe0", line_spacing=1.25)),
        ])
        ry += n * lh(7) + 0.05

    D.add_rect(s, MARGIN, band_top, CONTENT_W, band_h, fill=TRACK, rounded=True, radius=0.08)
    D.add_text(s, MARGIN + band_pad + 0.02, band_top + band_pad, CONTENT_W - 0.28,
               band_h - 2 * band_pad, [
        ("LE CRITÈRE DE SORTIE EST LE MIROIR DES CONDITIONS D'ENTRÉE",
         dict(size=7, bold=True, color=MUTED)),
        (sortie, dict(size=7, color=NAVY, space_before=3, line_spacing=1.25)),
    ])
    return s


# ---------------------------------------------------------------- slide 9
# Formes inspirées de la slide d'exemple « Une approche contextualisée » du
# template : colonne par étape avec badge + bandeau titre + ligne de
# séparation + bloc LIVRABLES, plutôt qu'un tableau plat.
def slide_livrables_ppt(prs):
    # v2.5 (chantier ④) : déplacée de la Proposition vers la Démarche, avec la
    # trajectoire fusionnée qui ne porte que les NOMS de livrables — ici le
    # détail des 4 profils (audience, contenu).
    s = content_slide(prs, "Démarche", "Livrables PPT par étape — piste à valider", color=D.PALETTE[3])
    D.add_text(s, MARGIN, CONTENT_TOP, CONTENT_W, 0.4, [
        ("iap-deck-builder est cadré comme un seul deck modulaire 16 sections, produit une fois "
         "à la Restitution — la trajectoire ci-avant implique plusieurs publics et moments de "
         "décision distincts. Piste : un profil de sections par étape, pas 4 générateurs séparés.",
         dict(size=8, color=MUTED, italic=True, line_spacing=1.2)),
    ])
    cols = [
        ("①", "Assessment flash", D.PALETTE[0], "Sponsor, comité de lancement",
         "Deck exécutif de restitution", "(déjà cadré) synthèse globale, axes valeur/complexité, radar T0"),
        ("②", "Premier déploiement", D.PALETTE[3], "Équipes pilotes + management",
         "Deck de plan de déploiement", "(nouveau) Cible TOM détaillée, backlog Coach/Délégué, mandat agent IA"),
        ("③", "Implémentation itérative", D.PALETTE[1], "Instance de comitologie",
         "Deck de comité de pilotage", "(nouveau, périodique) Avancement backlog, delta KPIs, risques actifs"),
        ("⟲", "Boucle de réévaluation", D.PALETTE[2], "Sponsor",
         "Deck de bilan / ré-évaluation", "(nouveau) Delta maturité T0→T+6-12, REX consolidé"),
    ]
    n = len(cols)
    pad = 0.16
    _, wcol = col_x(0, n)
    usable = wcol - 2 * pad
    # Carte plafonnée au contenu (bloc badge/titre → séparateur → audience →
    # LIVRABLES → contenu) au lieu de CONTENT_BOTTOM - top0 : sinon la colonne
    # s'étirait sur toute la hauteur et laissait ~60 % de vide sous le texte
    # (défaut « colonne timeline sur-étirée », slide 24).
    contenu_lines = max(_lignes(c[5], usable, 7) for c in cols)
    card_h = 1.56 + contenu_lines * (7 * 1.28 / 72.0) + 0.18
    region_top = CONTENT_TOP + 0.5
    top0 = region_top + min(0.45, max(0.0, (CONTENT_BOTTOM - region_top - card_h) / 2))
    badge_d = 0.34
    for i, (sym, titre, color, audience, deck, contenu) in enumerate(cols):
        x, w = col_x(i, n)
        D.add_card(s, x, top0, w, card_h, color)
        pad = 0.16
        D.add_rect(s, x + pad, top0 + 0.14, badge_d, badge_d, fill=color, rounded=True, radius=0.5)
        D.add_text(s, x + pad, top0 + 0.14, badge_d, badge_d, [
            (sym, dict(size=11, bold=(sym != "⟲"), color="#ffffff", align=PP_ALIGN.CENTER)),
        ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        D.add_text(s, x + pad + badge_d + 0.08, top0 + 0.14, w - 2 * pad - badge_d - 0.08, badge_d, [
            (titre, dict(size=8, bold=True, color=color, line_spacing=1.0)),
        ], anchor=MSO_ANCHOR.MIDDLE)
        line_y = top0 + 0.14 + badge_d + 0.1
        D.add_rect(s, x + pad, line_y, w - 2 * pad, 0.012, fill=LINE)
        D.add_text(s, x + pad, line_y + 0.08, w - 2 * pad, 0.22, [
            (audience, dict(size=7, italic=True, color=MUTED)),
        ])
        D.add_text(s, x + pad, line_y + 0.34, w - 2 * pad, 0.55, [
            ("LIVRABLES", dict(size=7, bold=True, color=NAVY)),
            (deck, dict(size=8, bold=True, color=NAVY, space_before=2, line_spacing=1.1)),
        ])
        contenu_top = line_y + 0.98
        D.add_text(s, x + pad, contenu_top, w - 2 * pad, top0 + card_h - contenu_top - 0.12, [
            (contenu, dict(size=7, color=MUTED, line_spacing=1.25)),
        ])
    return s


# Nouveau — brainstorm de design (v2.2) : reprend le pattern « cadre blanc »
# du template (déjà utilisé dans le REX "⛱️ L'Été de l'IA", VSCode1, en
# alternance gauche/droite pour chaque slide de contenu) — claim + puces à
# gauche, illustration encadrée à droite.
# Révisé (restructuration 2026-07-22) : slide 3 est l'énoncé de problème qui ouvre
# le deck — 3 puces d'une seule colonne vertébrale ancrées dans l'utilisateur
# (le constat → ce que ça coûte → ce qu'un bon diagnostic exige). La puce
# « l'IA amplifie l'organisation » a été RETIRÉE : elle est implicite dans la
# puce 2 et son point de doctrine est développé au chapitre IA. Le
# titre et l'image encadrée (sunset, cadre round2DiagRect) sont conservés.
def slide_vision(prs):
    """Slide 3 — la thèse qui lance le deck. Passe design 2026-07-23 : le pavé
    de 3 puces longues (dernier mur de texte du deck) devient un enchaînement
    vertical constat → coût → exigence (deck-design-library, pattern 6 « flux
    vertical connecté » transposé) : badges numérotés reliés par une ligne,
    un bloc par idée, kicker + claim + détail. « Un sur N en accent »
    (pattern 7) : le bloc 3 — LA thèse (partir des utilisateurs réels, pas
    d'une réponse toute faite) — est le seul en fond navy, mêmes teintes que
    la carte RÉSULTAT de l'executive summary (#8fd6db / #c7cbe0) pour rester
    dans le langage visuel déjà en place. Photo + cadre inchangés."""
    layout = prs.slide_masters[0].slide_layouts[LAYOUT_VISUEL_DROITE]
    s = prs.slides.add_slide(layout)
    phs = {ph.placeholder_format.idx: ph for ph in s.placeholders}

    phs[0].text_frame.text = ("Le risque n'est pas de manquer d'outils : "
                               "c'est de traiter le mauvais problème")
    for p in phs[0].text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = _rgb(NAVY)
    # Le corps est dessiné en blocs absolus — le placeholder BODY du layout
    # resterait un textframe vide par-dessus les cartes, on le retire.
    ph_body = phs[1]._element
    ph_body.getparent().remove(ph_body)

    # Couleurs = sémantique déjà en place : bleu Contexte (constat), rouge
    # ASSAINIR/gaspillage (coût), navy accent (l'exigence — la thèse).
    blocs = [
        ("LE CONSTAT", D.PALETTE[0],
         "Une infrastructure guichet ou centre de coûts subit la demande au lieu "
         "de la piloter.",
         "Ni utilisateurs identifiés, ni feuille de route, ni levier d'adoption."),
        ("CE QUE ÇA COÛTE", D.PALETTE[2],
         "La capacité disponible part en gaspillage.",
         "RUN subi (l'exploitation quotidienne), ressources orphelines, seniors "
         "sur du répétitif — et le réflexe « plus d'outils » ou « mettons de "
         "l'IA » aggrave le mal."),
        ("CE QU'UN BON DIAGNOSTIC EXIGE", NAVY,
         "Partir des utilisateurs réels et de leurs douleurs — pas d'une réponse "
         "toute faite.",
         "Le fil que déroulent tous les chapitres qui suivent."),
    ]

    # Colonne de gauche : s'arrête net avant le cadre photo (groupe du layout
    # à x=6.857in) — badges à la marge, cartes décalées à droite de la chaîne.
    badge_d = 0.34
    chain_cx = MARGIN + badge_d / 2
    card_x = MARGIN + 0.50
    card_w = 6.63 - card_x
    text_x = card_x + 0.22
    usable = card_w - 0.22 - 0.18
    kicker_h, gap_blocs = 0.20, 0.22

    # Hauteur de chaque bloc = son contenu (pas de panneau sur-étiré). Le
    # cpi_ref par défaut de l'estimateur (11.0) sous-estime nettement la
    # police du template sur cette largeur (mesuré au rendu réel : ~15-16
    # équivalent 10.5pt) — 14.0 garde une marge de sécurité sans laisser
    # 0.3in de vide sous la ligne de détail des cartes (défaut vu au zoom).
    def _est(texte, taille):
        return max(1, D.estimer_lignes(texte, usable, taille, cpi_ref=14.0))

    dims = []
    for _, _, claim, detail in blocs:
        claim_h = _est(claim, 10.5) * (10.5 * 1.2 / 72.0) + 0.04
        detail_h = _est(detail, 9) * (9 * 1.3 / 72.0) + 0.04
        h = 0.13 + kicker_h + claim_h + 0.07 + detail_h + 0.13
        dims.append((h, claim_h, detail_h))

    top0 = 1.40
    tops, y = [], top0
    for h, _, _ in dims:
        tops.append(y)
        y += h + gap_blocs

    # La chaîne : UNE ligne verticale continue sous les badges (pattern 6 —
    # pas de flèches), du centre du bloc 1 au centre du bloc 3.
    c1 = tops[0] + dims[0][0] / 2
    c3 = tops[2] + dims[2][0] / 2
    D.add_rect(s, chain_cx - 0.01, c1, 0.02, c3 - c1, fill=LINE)

    for i, ((label, color, claim, detail), (h, claim_h, detail_h), top) in \
            enumerate(zip(blocs, dims, tops)):
        accent = (i == len(blocs) - 1)   # « un sur N » : l'exigence, la thèse
        if accent:
            D.add_rect(s, card_x, top, card_w, h, fill=NAVY, rounded=True, radius=0.06)
            D.add_rect(s, card_x, top, 0.07, h, fill=ACCENT, rounded=True, radius=0.5)
        else:
            D.add_card(s, card_x, top, card_w, h, color)
        D.add_text(s, text_x, top + 0.13, usable, kicker_h, [
            (label, dict(size=8, bold=True, color="#8fd6db" if accent else color)),
        ])
        D.add_text(s, text_x, top + 0.13 + kicker_h, usable, claim_h, [
            (claim, dict(size=10.5, bold=True, color="#ffffff" if accent else NAVY,
                         line_spacing=1.2)),
        ])
        D.add_text(s, text_x, top + 0.13 + kicker_h + claim_h + 0.07, usable, detail_h, [
            (detail, dict(size=9, color="#c7cbe0" if accent else MUTED,
                          line_spacing=1.3)),
        ])
        # Badge numéroté sur la chaîne, centré sur son bloc.
        by = top + h / 2 - badge_d / 2
        D.add_rect(s, chain_cx - badge_d / 2, by, badge_d, badge_d,
                   fill=color, rounded=True, radius=0.5)
        D.add_text(s, chain_cx - badge_d / 2, by, badge_d, badge_d, [
            (str(i + 1), dict(size=11, bold=True, color="#ffffff",
                              align=PP_ALIGN.CENTER)),
        ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

    cadre = _find_frame_in_group(s.slide_layout.shapes, "Google Shape;212;p17", "Google Shape;213;p17")
    for pb in frame_obstructions(s, *cadre[:4]) if cadre else []:
        print("  [obstruction] vision:", pb["source"], pb["name"], pb["reason"])
    _remplir_cadre(s, cadre, "sunset", seed=1)
    return s


# (v2.5 — slide_schema_bout_en_bout supprimée, chantier ① : sa trame ①②③⟲
# doublonnait slide_trajectoire ; son apport — le livrable-clé par phase — est
# fusionné dans slide_trajectoire ci-dessus.)


# ---------------------------------------------------------------- slide 10
# Nouveau (v2.0) : la bifurcation "avec/sans agent IA" de slide_trajectoire
# gagne un livrable concret, distinct des 4 decks PPT — un markdown pour
# l'équipe qui exécute, pas pour le sponsor. Deux cartes (mêmes proportions
# que slide_mission) + un bandeau de routage + une note "pas un aller simple".
def slide_export_markdown(prs):
    s = content_slide(prs, "IA",
                       "Export markdown — agentic ou documentation, selon le contexte client (piste à valider)",
                       color=D.PALETTE[4])
    # v2.6 (point ④) : badge de série (comme les 3 slides d'agent candidat) —
    # l'intro cède la largeur du badge. Le renvoi aux 4 decks vise le chapitre
    # Démarche (slide_livrables_ppt y a déménagé en v2.5 — la mention
    # « Proposition » était restée, corrigée ici).
    badge_deploiement_agentic(s)
    D.add_text(s, MARGIN, CONTENT_TOP, CONTENT_W - BADGE_AGENTIC_W - 0.2, 0.5, [
        ("Pas un 5e deck PPT : un livrable markdown pour l'équipe qui exécute (versionnable, "
         "committable) — les 4 decks du chapitre Démarche restent pour sponsor et comité de pilotage.",
         dict(size=8, color=MUTED, italic=True, line_spacing=1.2)),
    ])

    cards = [
        ("DOCUMENTATION-FIRST", D.PALETTE[2],
         "Agentic Readiness [0]-[1], données D3-D4 sans LLM local, ou score de gaspillage faible.",
         "Runbook du processus", "iap-adoption-plan"),
        ("AGENTIC-IMPLEMENTATION", D.PALETTE[0],
         "Agentic Readiness [2]-[3], données D0-D2 (ou D3-D4 avec LLM local), score positif.",
         "Plan d'implémentation agentic", "iap-agentic-opportunities"),
    ]
    # v2.6 : +0.13 — l'intro, rétrécie par le badge de série, passe à 3 lignes.
    top0 = CONTENT_TOP + 0.55
    card_h = 1.55
    for i, (titre, color, quand, fichier, owner) in enumerate(cards):
        x, w = col_x(i, 2)
        D.add_card(s, x, top0, w, card_h, color)
        pad = 0.2
        D.add_text(s, x + pad, top0 + 0.14, w - 2 * pad, 0.3, [
            (titre, dict(size=D.TYPE["h3"], bold=True, color=color)),
        ])
        D.add_text(s, x + pad, top0 + 0.5, w - 2 * pad, 0.55, [
            ("QUAND", dict(size=D.TYPE["tiny"], bold=True, color=MUTED)),
            (quand, dict(size=8, color=NAVY, space_before=2, line_spacing=1.2)),
        ])
        D.add_text(s, x + pad, top0 + 1.08, w - 2 * pad, 0.4, [
            ("LIVRABLE · OWNER", dict(size=7, bold=True, color=MUTED)),
            (f"{fichier} — {owner}", dict(size=8, bold=True, color=NAVY, space_before=2)),
        ])

    signals_top = top0 + card_h + 0.15
    signals_h = 0.55
    signals = [
        ("PILIER AGENTIC READINESS", "[0-1] → documentation · [2-3] → agentic"),
        ("DONNÉES (GATE IA)", "D3-D4 sans LLM local → doc · D0-D2 → agentic"),
        ("SCORE DE GASPILLAGE", "faible/négatif → doc · positif → agentic"),
    ]
    for i, (label, mapping) in enumerate(signals):
        x, w = col_x(i, 3)
        D.add_rect(s, x, signals_top, w, signals_h, fill=TRACK, rounded=True, radius=0.1)
        D.add_text(s, x + 0.12, signals_top + 0.06, w - 0.24, signals_h - 0.12, [
            (label, dict(size=7, bold=True, color=NAVY)),
            (mapping, dict(size=7, color=MUTED, space_before=2, line_spacing=1.15)),
        ])

    note_top = signals_top + signals_h + 0.15
    note_h = min(1.05, CONTENT_BOTTOM - note_top)
    D.add_rect(s, MARGIN, note_top, CONTENT_W, note_h, fill=NAVY, rounded=True, radius=0.08)
    D.add_text(s, MARGIN + 0.22, note_top, CONTENT_W - 0.44, note_h, [
        ("Pas un aller simple", dict(size=D.TYPE["tiny"], bold=True, color="#ffffff")),
        ("À chaque boucle de réévaluation, le même fichier est amendé — jamais dupliqué — sur le "
         "modèle du registre de risques IA : un client documentation-first peut basculer en agentic si "
         "son Agentic Readiness progresse, et inversement en cas de perte de compétence avérée.",
         dict(size=8, color="#c7cbe0", space_before=4, line_spacing=1.25)),
    ], anchor=MSO_ANCHOR.MIDDLE)
    return s


# --- Nouveau (v2.6, point ③) : le schéma d'architecture de fonctionnement
# d'IAP en contexte client — joue aussi le rôle de la slide « ce que le module
# met dans les mains du consultant » annoncée par le plan v2.5 et jamais
# écrite. Trois zones (pattern 8 du catalogue deck-design-library : des zones
# teintées qui se saisissent sans lire les étiquettes) reliées par un
# vocabulaire de flèches uniforme (pattern 9) : le POSTE DU CONSULTANT (le
# module — 4 étapes aux couleurs de slide_schema_fonctionnement, 11 agents,
# gate confidentialité bloquant en bandeau navy comme slide_architecture_agents),
# le CONTEXTE CLIENT (sponsor/équipes + SI selon l'ambition A/B/C des slides
# SUIVANTES), et entre les deux les FLUX (collecte entrante, livrables
# sortants, agents retenus). La zone « déploiement agentic chez le client »
# est le seul bloc en aplat plein (« un sur N en accent »), violet
# D.PALETTE[4] = couleur du chapitre 06 · IA — la même que le badge de série
# posé sur les 4 slides de proposition agentic (point ④).
def slide_iap_contexte_client(prs):
    s = content_slide(prs, "Outillage IAP",
                       "IAP tourne sur le poste du consultant — seuls les agents retenus passent chez le client",
                       color=D.PALETTE[5])
    D.add_text(s, MARGIN, CONTENT_TOP, CONTENT_W, 0.30, [
        ("Rien ne s'installe côté client par défaut : les interviews et les exports entrent, "
         "les livrables sortent — le déploiement d'agents est une décision de ②/③, pas un prérequis.",
         dict(size=8, color=MUTED, italic=True, line_spacing=1.15)),
    ])

    z_top = CONTENT_TOP + 0.40
    z_h = 2.90
    cons_x, cons_w = MARGIN, 3.30
    flux_x, flux_w = cons_x + cons_w + 0.12, 1.30
    cli_x = flux_x + flux_w + 0.12
    cli_w = BORD_DROIT - cli_x
    pad = 0.16

    # --- Zone 1 : poste du consultant (le module IAP).
    D.add_rect(s, cons_x, z_top, cons_w, z_h, fill=TRACK, rounded=True, radius=0.06)
    D.add_text(s, cons_x + pad, z_top + 0.12, cons_w - 2 * pad, 0.42, [
        ("POSTE DU CONSULTANT", dict(size=8, bold=True, color=NAVY)),
        ("Le module IAP · 11 workflows outillés", dict(size=7, color=MUTED, space_before=2)),
    ])
    etapes = [("COLLECTE", D.PALETTE[0]), ("DIAGNOSTIC", D.PALETTE[4]),
              ("CONCEPTION", D.PALETTE[3]), ("RESTITUTION", D.PALETTE[1])]
    pill_w = (cons_w - 2 * pad - 0.10) / 2
    pill_h = 0.30
    pills_top = z_top + 0.60
    for i, (nom, color) in enumerate(etapes):
        px = cons_x + pad + (i % 2) * (pill_w + 0.10)
        py = pills_top + (i // 2) * (pill_h + 0.08)
        chip(s, px, py, pill_w, pill_h, nom, color, size=6.5)
    caption_top = pills_top + 2 * pill_h + 0.08 + 0.10
    D.add_text(s, cons_x + pad, caption_top, cons_w - 2 * pad, 0.40, [
        ("Mêmes étapes que le schéma de fonctionnement (chapitre 07 · Démarche) — "
         "fonctionne aussi sans IA externe si le gate l'impose (mode M0).",
         dict(size=6.5, color=MUTED, italic=True, line_spacing=1.2)),
    ])
    gate_h = 0.62
    gate_top = z_top + z_h - gate_h - 0.12
    D.add_rect(s, cons_x + pad - 0.04, gate_top, cons_w - 2 * pad + 0.08, gate_h,
               fill=NAVY, rounded=True, radius=0.10)
    D.add_text(s, cons_x + pad + 0.08, gate_top, cons_w - 2 * pad - 0.16, gate_h, [
        ("GATE CONFIDENTIALITÉ — BLOQUANT", dict(size=7, bold=True, color="#ffffff")),
        ("Classe la donnée (D0-D4) avant tout usage IA sur donnée client",
         dict(size=6.5, color="#c7cbe0", space_before=2, line_spacing=1.15)),
    ], anchor=MSO_ANCHOR.MIDDLE)

    # --- Zone 2 (milieu) : les flux — vocabulaire de flèches uniforme, une
    # flèche par ligne, cyan = données de la mission, violet = agents déployés.
    flux = [
        ("←", ACCENT, "COLLECTE ENTRANTE", "interviews · exports d'outils"),
        ("→", ACCENT, "LIVRABLES SORTANTS", "deck exécutif · export markdown"),
        ("→", D.PALETTE[4], "AGENTS RETENUS (②/③)", "supervisés puis délégués"),
    ]
    f_h = 0.72
    f_gap = (z_h - 3 * f_h) / 2
    for i, (fleche, color, label, detail) in enumerate(flux):
        fy = z_top + i * (f_h + f_gap)
        D.add_text(s, flux_x, fy, flux_w, 0.30, [
            (fleche, dict(size=16, bold=True, color=color, align=PP_ALIGN.CENTER)),
        ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        D.add_text(s, flux_x - 0.06, fy + 0.30, flux_w + 0.12, f_h - 0.30, [
            (label, dict(size=6.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER, line_spacing=1.1)),
            (detail, dict(size=6, color=MUTED, space_before=1, align=PP_ALIGN.CENTER, line_spacing=1.1)),
        ], align=PP_ALIGN.CENTER)

    # --- Zone 3 : contexte client.
    D.add_rect(s, cli_x, z_top, cli_w, z_h, fill="#ffffff", line=LINE, line_w=1.0,
               rounded=True, radius=0.06)
    D.add_text(s, cli_x + pad, z_top + 0.12, cli_w - 2 * pad, 0.24, [
        ("CONTEXTE CLIENT", dict(size=8, bold=True, color=NAVY)),
    ])
    bloc_x = cli_x + pad
    bloc_w = cli_w - 2 * pad
    b1_top = z_top + 0.42
    D.add_rect(s, bloc_x, b1_top, bloc_w, 0.55, fill=TRACK, rounded=True, radius=0.10)
    D.add_text(s, bloc_x + 0.12, b1_top, bloc_w - 0.24, 0.55, [
        ("SPONSOR & ÉQUIPES INTERVIEWÉES", dict(size=7, bold=True, color=NAVY)),
        ("les voix du diagnostic — interviews par persona",
         dict(size=6.5, color=MUTED, space_before=1, line_spacing=1.15)),
    ], anchor=MSO_ANCHOR.MIDDLE)
    b2_top = b1_top + 0.55 + 0.10
    b2_h = 0.72
    D.add_rect(s, bloc_x, b2_top, bloc_w, b2_h, fill=TRACK, rounded=True, radius=0.10)
    D.add_text(s, bloc_x + 0.12, b2_top, bloc_w - 0.24, b2_h, [
        ("SI CLIENT", dict(size=7, bold=True, color=NAVY)),
        ("ServiceNow · Jira · Confluence · Datadog · CMDB · FinOps — accès selon "
         "l'ambition A/B/C (slides suivantes)",
         dict(size=6.5, color=MUTED, space_before=1, line_spacing=1.15)),
    ], anchor=MSO_ANCHOR.MIDDLE)
    # Zone mise en évidence — le seul aplat plein du schéma (« un sur N »).
    b3_top = b2_top + b2_h + 0.12
    b3_h = z_top + z_h - 0.14 - b3_top
    D.add_rect(s, bloc_x, b3_top, bloc_w, b3_h, fill=D.PALETTE[4], rounded=True, radius=0.10)
    D.add_text(s, bloc_x + 0.12, b3_top, bloc_w - 0.24, b3_h, [
        ("DÉPLOIEMENT AGENTIC CHEZ LE CLIENT", dict(size=7, bold=True, color="#ffffff")),
        ("Les agents candidats retenus se déploient ici en ②/③ — supervisés puis délégués",
         dict(size=6.5, color="#e8def5", space_before=2, line_spacing=1.2)),
    ], anchor=MSO_ANCHOR.MIDDLE)

    # --- Renvoi de focus (point ④a) : même langage visuel que le badge de
    # série posé sur les 4 slides visées — liseré violet, renvoi par CHAPITRE.
    band_top = z_top + z_h + 0.16
    band_h = min(0.66, CONTENT_BOTTOM - band_top)
    D.add_rect(s, MARGIN, band_top, CONTENT_W, band_h, fill="#ffffff",
               line=D.PALETTE[4], line_w=1.0, rounded=True, radius=0.10)
    D.add_text(s, MARGIN + 0.2, band_top, CONTENT_W - 0.4, band_h, [
        ("QUATRE PROPOSITIONS DE DÉPLOIEMENT AGENTIC — DÉTAILLÉES AU CHAPITRE 06 · IA",
         dict(size=7, bold=True, color=D.PALETTE[4])),
        ("Agent de triage RUN · veille FinOps · agent documentaire (RAG) · export markdown "
         "(qui porte la décision agentic/documentation) — chacune porte le badge "
         "« déploiement agentic chez le client ».",
         dict(size=7, color=NAVY, space_before=2, line_spacing=1.2)),
    ], anchor=MSO_ANCHOR.MIDDLE)
    return s


# ---------------------------------------------------------------- slide 11
def slide_ambition(prs):
    # v2.5 (chantier ④) : déplacée de la Proposition vers l'Outillage IAP —
    # le niveau d'ambition qualifie l'outil, pas la proposition de transformation.
    s = content_slide(prs, "Outillage IAP", "Trois niveaux d'ambition, pas un spectre linéaire", color=D.PALETTE[5])
    niveaux = [
        ("A", "Aide au coach", D.PALETTE[0],
         "Génère un livrable à la demande — aucune initiative propre. Le consultant pilote à 100 %.",
         "État actuel du cadrage (MVP0–MVP5)"),
        ("B", "Assistant interactif", D.PALETTE[3],
         "Guide pas à pas, pose des questions de clarification, signale les incohérences.",
         "Palier intermédiaire, entre MVP5 et MVP6"),
        ("C", "Companion connecté", D.PALETTE[2],
         # « quasi autonome » seul survendait C comme une autonomie décisionnelle,
         # ce que le cadrage (l.733) demande explicitement d'éviter.
         "Connecté en direct à ServiceNow/Jira/Confluence/Datadog/CMDB/FinOps : quasi "
         "autonome sur la collecte et la préparation — jamais sur l'arbitrage.",
         "= MVP6, non engagé"),
    ]
    n = 3
    pad = 0.2
    _, wcol = col_x(0, n)
    usable = wcol - 2 * pad
    # Légende roadmap collée SOUS le texte de rôle (au lieu d'un y fixe en bas
    # de carte) et carte plafonnée à son contenu : supprime le « trou mort »
    # entre le rôle et la légende, et le sur-étirement de la carte (slide 26).
    # Passe de design 2026-07-23 — pattern 7 du catalogue deck-design-library
    # (« un sur N en accent ») : le niveau A, ÉTAT ACTUEL assumé du cadrage
    # (cf. executive summary, posture de gouvernance), est la seule carte en
    # fill navy plein — B et C restent des cartes blanches identiques.
    role_lines = max(_lignes(niv[3], usable, 8) for niv in niveaux)
    role_h = role_lines * (8 * 1.25 / 72.0) + 0.06
    roadmap_y = 0.55 + role_h + 0.12
    card_h = roadmap_y + 0.28 + 0.14
    top0 = CONTENT_TOP + 0.45
    accent_idx = 0   # A · Aide au coach — l'état actuel
    for i, (code, titre, color, role, roadmap) in enumerate(niveaux):
        x, w = col_x(i, n)
        accent = (i == accent_idx)
        if accent:
            D.add_rect(s, x, top0, w, card_h, fill=NAVY, rounded=True, radius=0.06)
            D.add_rect(s, x, top0, 0.07, card_h, fill=color, rounded=True, radius=0.5)
        else:
            D.add_card(s, x, top0, w, card_h, color)
        D.add_text(s, x + pad, top0 + 0.15, w - 2 * pad, 0.35, [
            (f"{code} · {titre}", dict(size=D.TYPE["small"], bold=True,
                                       color="#ffffff" if accent else color)),
        ])
        D.add_text(s, x + pad, top0 + 0.55, w - 2 * pad, role_h, [
            (role, dict(size=8, color="#e8ebf5" if accent else NAVY, line_spacing=1.25)),
        ])
        D.add_text(s, x + pad, top0 + roadmap_y, w - 2 * pad, 0.3, [
            (roadmap, dict(size=8, bold=True, color="#ffffff" if accent else MUTED)),
        ])

    note_top = top0 + card_h + 0.18
    note_h = CONTENT_BOTTOM - note_top
    D.add_text(s, MARGIN, note_top, CONTENT_W, note_h, [
        ("Monter de A à C n'est pas qu'une question de fonctionnalités : le niveau C suppose "
         "un accès direct aux données de production du client — risque sécurité/confidentialité "
         "d'un tout autre ordre. Un cabinet peut durablement rester au niveau A ou B par choix "
         "de gouvernance, pas seulement par contrainte technique transitoire. "
         "MVP0–6 = jalons de la roadmap de mise en œuvre ; MVP6, le « companion connecté », n'est pas engagé.",
         dict(size=8, color=MUTED, italic=True, line_spacing=1.3)),
    ])
    return s


# ---------------------------------------------------------------- slide 12
def slide_kpis(prs):
    s = content_slide(prs, "KPI", "Trois familles de KPIs, à ne jamais confondre", color=D.PALETTE[0])
    familles = [
        ("KPIs de mission", D.PALETTE[0], "Côté client",
         ["Gaspillage traité (capacité RUN récupérée)", "Adoption produit (self-service)",
          "Fiabilité & coût (MTTR : délai moyen de résolution, coût/capacité)", "Gouvernance IA (supervision, incidents)",
          "Maturité (delta par pilier, T0→réévaluation)"]),
        ("KPIs d'usage du module", D.PALETTE[1], "Côté cabinet",
         ["Accélération (temps pour un cadrage flash)", "Réutilisation (templates vs ad hoc)",
          "Cohérence (écarts détectés par risk-reviewer)", "Capitalisation (REX ajoutés)",
          "Adoption interne (consultants, missions)"]),
        ("Grille de maturité", D.PALETTE[3], "Progression dans le temps",
         ["Delta par pilier (Excellence Tech., Agilité, IA/Agentic)",
          "Re-assessment T+6–12 mois",
          "Score de priorisation ≠ KPI de résultat"]),
    ]
    n = 3
    pad = 0.18
    _, wcol = col_x(0, n)
    usable = wcol - 2 * pad
    # Carte plafonnée au contenu (titre + sous-titre + puces) puis bande de
    # cartes CENTRÉE verticalement — au lieu de card_h = CONTENT_H qui étirait
    # chaque colonne sur toute la hauteur et laissait un grand vide sous les
    # puces (défaut « panneau sur-étiré », slide 28).
    # Passe de design 2026-07-23 — pattern 3 du catalogue deck-design-library
    # (« grille de cartes stat, une en accent ») : la Grille de maturité est la
    # seule carte en fill navy plein — c'est la famille que le chapitre détaille
    # ensuite (slide_maturite + message « le KPI = le delta T0→réévaluation ») ;
    # corps monté à 8.5pt (la densité s'absorbe par la police, pas par le vide).
    accent_idx = 2   # « Grille de maturité »
    def _bloc_puces(items):
        lignes = sum(_lignes("·  " + it, usable, 8.5) for it in items)
        return lignes * (8.5 * 1.15 / 72.0) + len(items) * (4 / 72.0)
    bullets_h = max(_bloc_puces(items) for *_, items in familles)
    card_h = 0.8 + bullets_h + 0.22
    top0 = CONTENT_TOP + max(0.0, (CONTENT_H - card_h) / 2)
    for i, (titre, color, sous, items) in enumerate(familles):
        x, w = col_x(i, n)
        accent = (i == accent_idx)
        if accent:
            D.add_rect(s, x, top0, w, card_h, fill=NAVY, rounded=True, radius=0.06)
            D.add_rect(s, x, top0, 0.07, card_h, fill=color, rounded=True, radius=0.5)
        else:
            D.add_card(s, x, top0, w, card_h, color)
        D.add_text(s, x + pad, top0 + 0.16, w - 2 * pad, 0.55, [
            (titre, dict(size=D.TYPE["small"], bold=True,
                         color="#ffffff" if accent else color, line_spacing=1.05)),
            (sous, dict(size=8.5, color="#aeb6d4" if accent else MUTED,
                        italic=True, space_before=2)),
        ])
        lignes = [(f"·  {it}", dict(size=8.5, color="#e8ebf5" if accent else NAVY,
                                    space_after=4, line_spacing=1.15))
                  for it in items]
        D.add_text(s, x + pad, top0 + 0.8, w - 2 * pad, card_h - 0.95, lignes)
    return s

# --- Brainstorm KPIs relancé (v2.1, docs/bmad-iap-cadrage.md §KPIs) — pourquoi
# chaque famille, quoi mesurer précisément, comment la mettre en place, et un
# exemple chiffré sur le cas nominal déjà posé pour l'export markdown.
def slide_kpis_pourquoi_quoi(prs):
    s = content_slide(prs, "KPI", "KPIs : pourquoi chaque famille, et quoi mesurer précisément", color=D.PALETTE[0])
    familles = [
        ("KPIs de mission", D.PALETTE[0],
         "Sans eux, un deck peut être livré dans les règles sans jamais savoir si le client va "
         "réellement mieux — la « transformation cosmétique » appliquée cette fois au résultat.",
         "Capacité RUN récupérée en heures/mois ; delta de MTTR en minutes ; taux de self-service "
         "sur la capacité livrée — pas un pourcentage vague."),
        ("KPIs d'usage du module", D.PALETTE[1],
         "Le module est réutilisé mission après mission — sans mesure d'usage, impossible de "
         "distinguer une méthode qui s'améliore d'une méthode qui stagne.",
         "Temps en heures consultant pour un cadrage flash ; part des livrables issus d'un template "
         "sans réécriture substantielle — pas juste « utilisé un template »."),
        ("Grille de maturité", D.PALETTE[3],
         "Sans mesure répétée dans le temps, la maturité reste une opinion de consultant, pas un "
         "delta objectivable — ce qui rend la boucle ⟲ vérifiable plutôt que déclarative.",
         "Delta par pilier (pas un score agrégé qui masquerait un recul) ; même instrument (grille "
         "V3.2) à T0 et à chaque re-assessment."),
    ]
    top0 = CONTENT_TOP + 0.05
    row_h = (CONTENT_H - 0.05 - 2 * 0.12) / 3
    for i, (nom, color, pourquoi, quoi) in enumerate(familles):
        y = top0 + i * (row_h + 0.12)
        D.add_card(s, MARGIN, y, CONTENT_W, row_h, color)
        pad = 0.18
        D.add_text(s, MARGIN + pad, y + 0.12, 1.7, row_h - 0.24, [
            (nom, dict(size=8, bold=True, color=color, line_spacing=1.15)),
        ], anchor=MSO_ANCHOR.MIDDLE)
        colw = (CONTENT_W - 2 * pad - 1.7 - 0.2) / 2
        x1 = MARGIN + pad + 1.7 + 0.2
        D.add_text(s, x1, y + 0.12, colw, row_h - 0.24, [
            ("POURQUOI", dict(size=7, bold=True, color=MUTED)),
            (pourquoi, dict(size=8, color=NAVY, space_before=3, line_spacing=1.2)),
        ])
        x2 = x1 + colw + 0.15
        D.add_text(s, x2, y + 0.12, colw, row_h - 0.24, [
            ("QUOI, PRÉCISÉMENT", dict(size=7, bold=True, color=MUTED)),
            (quoi, dict(size=8, color=NAVY, space_before=3, line_spacing=1.2)),
        ])
    return s


def slide_kpis_mise_en_place(prs):
    s = content_slide(prs, "KPI", "KPIs : comment on les met en place, concrètement", color=D.PALETTE[0])
    familles = [
        ("KPIs de mission", D.PALETTE[0], "iap-metrics-sre-finops-lead",
         "ServiceNow/Jira/CMDB si accès (preuves externes), sinon déclaratif — tagué DÉDUIT",
         "Continu, lu à chaque étape ②③⟲"),
        ("KPIs d'usage du module", D.PALETTE[1], "Le consultant, au fil des missions",
         "Journal de mission + bibliothèque de REX",
         "Par mission, consolidé à MVP5"),
        ("Grille de maturité", D.PALETTE[3], "iap-strategy-lead",
         "Grille V3.2 repassée en atelier ou en interview",
         "T0 (① Assessment flash) puis chaque boucle ⟲"),
    ]
    top0 = CONTENT_TOP + 0.05
    row_h = 1.0
    name_w = 2.0
    for i, (nom, color, owner, source, cadence) in enumerate(familles):
        y = top0 + i * (row_h + 0.12)
        D.add_rect(s, MARGIN, y, CONTENT_W, row_h, fill="#ffffff", line=LINE, line_w=0.75, rounded=True, radius=0.08)
        D.add_rect(s, MARGIN, y, 0.06, row_h, fill=color, rounded=True, radius=0.5)
        D.add_text(s, MARGIN + 0.2, y + 0.1, name_w, row_h - 0.2, [
            (nom, dict(size=8, bold=True, color=color, line_spacing=1.15)),
        ], anchor=MSO_ANCHOR.MIDDLE)
        colw = (CONTENT_W - 0.2 - name_w - 3 * 0.15) / 3
        specs = [("OWNER", owner), ("SOURCE DES DONNÉES", source), ("CADENCE", cadence)]
        for j, (label, val) in enumerate(specs):
            xj = MARGIN + 0.2 + name_w + 0.15 + j * (colw + 0.15)
            D.add_text(s, xj, y + 0.1, colw, row_h - 0.2, [
                (label, dict(size=7, bold=True, color=MUTED)),
                (val, dict(size=8, color=NAVY, space_before=3, line_spacing=1.15)),
            ], anchor=MSO_ANCHOR.MIDDLE)

    note_top = top0 + 3 * row_h + 2 * 0.12 + 0.15
    note_h = min(0.85, CONTENT_BOTTOM - note_top)
    D.add_rect(s, MARGIN, note_top, CONTENT_W, note_h, fill=TRACK, rounded=True, radius=0.08)
    D.add_text(s, MARGIN + 0.2, note_top, CONTENT_W - 0.4, note_h, [
        ("Pas d'instrumentation automatique en MVP1", dict(size=8, bold=True, color=NAVY)),
        ("Cohérent avec le Niveau A/B assumé (§Ambition de l'outil) : recueil et rapport à la main "
         "tant qu'aucun tableau de bord temps réel n'est promis avant le Niveau C.",
         dict(size=8, color=MUTED, space_before=3, line_spacing=1.2)),
    ], anchor=MSO_ANCHOR.MIDDLE)
    return s


def slide_kpis_exemple(prs):
    s = content_slide(prs, "KPI", "KPIs en pratique : le cas nominal RUN massif, avant/après", color=D.PALETTE[0])
    D.add_text(s, MARGIN, CONTENT_TOP, CONTENT_W, 0.28, [
        ("Même fixture illustrative que le cas nominal de l'export markdown (chapitre IA) — pas un client réel.",
         dict(size=8, color=MUTED, italic=True)),
    ])
    col_widths = [2.85, 1.75, 1.85, 1.725]
    headers = ["KPI", "T0 · ① ASSESSMENT FLASH", "T+6-12 MOIS · ⟲ RÉÉVALUATION", "TAG"]
    xs = []
    cx = MARGIN
    for cw in col_widths:
        xs.append(cx)
        cx += cw + 0.12

    header_y = CONTENT_TOP + 0.4
    for x, w, label in zip(xs, col_widths, headers):
        # _header_cell : le "⟲" de « T+6-12 MOIS · ⟲ RÉÉVALUATION » serait un
        # tofu en gras (cf. _GLYPHES_SANS_GRAS) — posé bold=False pour ce seul
        # caractère, le reste du libellé reste en gras.
        _header_cell(s, x, header_y, w, 0.24, label, size=7, color=MUTED, bold=True)

    tagcolor = {"CONFIRMÉ": SEVERITE[0], "DÉDUIT": SEVERITE[2], "—": MUTED}
    rows = [
        ("Pilier Agentic Readiness", "[1] — process pas assez explicite",
         "[2] — process explicite, rôles définis", "CONFIRMÉ"),
        ("Tickets récurrents évités / mois", "0", "≈ 15", "DÉDUIT"),
        ("Temps de triage moyen / ticket", "25 min", "12 min", "DÉDUIT"),
        ("Recommandation associée", "Documentation-first (runbook)",
         "Agentic-implementation (même fichier amendé)", "—"),
    ]
    row_top = header_y + 0.32
    row_h = 0.62
    row_gap = 0.08
    for i, (kpi, t0, t1, tag) in enumerate(rows):
        y = row_top + i * (row_h + row_gap)
        D.add_rect(s, MARGIN, y + row_h, CONTENT_W, 0.012, fill=LINE)
        D.add_text(s, xs[0], y, col_widths[0], row_h, [
            (kpi, dict(size=8, bold=True, color=NAVY, line_spacing=1.15)),
        ], anchor=MSO_ANCHOR.MIDDLE)
        D.add_text(s, xs[1], y, col_widths[1], row_h, [
            (t0, dict(size=8, color=NAVY, line_spacing=1.15)),
        ], anchor=MSO_ANCHOR.MIDDLE)
        D.add_text(s, xs[2], y, col_widths[2], row_h, [
            (t1, dict(size=8, color=NAVY, line_spacing=1.15)),
        ], anchor=MSO_ANCHOR.MIDDLE)
        if tag == "—":
            D.add_text(s, xs[3], y, col_widths[3], row_h, [
                (tag, dict(size=8, color=MUTED, align=PP_ALIGN.CENTER)),
            ], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        else:
            chip(s, xs[3], y + row_h / 2 - 0.13, min(1.1, col_widths[3]), 0.26, tag, tagcolor[tag], size=7)

    note_top = row_top + 4 * row_h + 3 * row_gap + 0.15
    note_h = min(0.6, CONTENT_BOTTOM - note_top)
    D.add_text(s, MARGIN, note_top, CONTENT_W, note_h, [
        ("Le point à retenir n'est pas l'ampleur des chiffres (fixture, pas une preuve) mais la "
         "discipline de mesure : même instrument aux deux instants, tag de confiance explicite, et "
         "un KPI de maturité qui déclenche directement le changement de recommandation.",
         dict(size=8, color=MUTED, italic=True, line_spacing=1.25)),
    ])
    return s


# --- Nouveau (brainstorm) : rendre tangible, dans le chapitre IA (chapitre 05
# en v2.5, APRÈS la proposition), ce que "piste agentique" veut dire concrètement — 3
# candidats illustratifs ancrés sur des familles de gaspillage déjà cadrées
# (§Traitement des gaspillages), pas des exemples inventés hors cadre. Chaque
# carte reprend la couleur de sa famille de gaspillage (RUN=rouge, Financier=or,
# Cognitif=violet) via l'argument `color` — cohérence intentionnelle avec
# slide_familles, pas un hasard de palette.
def slide_agent_ia(prs, titre, nom_agent, famille, why, what, gain, color, note=None):
    s = content_slide(prs, "IA", titre, color=D.PALETTE[4])
    # v2.6 (point ④) : badge de série en haut à droite — l'en-tête cède la
    # largeur du badge pour ne pas passer dessous.
    D.add_text(s, MARGIN, CONTENT_TOP, CONTENT_W - BADGE_AGENTIC_W - 0.2, 0.45, [
        (nom_agent, dict(size=D.TYPE["h3"], bold=True, color=color)),
        (f"Gaspillage {famille}", dict(size=8, color=MUTED, italic=True, space_before=2)),
    ])
    badge_deploiement_agentic(s)
    top0 = CONTENT_TOP + 0.55
    bands = [
        ("POURQUOI", why),
        ("CE QUE FAIT L'AGENT", what),
        ("GAIN", gain),
    ]
    txt_size = 9
    line_h = txt_size * 1.25 / 72.0
    usable = CONTENT_W - 0.44
    # Chaque bandeau plafonné à SON contenu — le bandeau GAIN (souvent 1 ligne)
    # ne garde plus la hauteur fixe d'un bandeau à 2 lignes (défaut « panneau
    # sur-étiré » constaté slides 10/11) — puis les 3 bandeaux sont répartis
    # pour remplir la zone, donc pas de vide résiduel en bas non plus.
    heights = [0.42 + _lignes(t, usable, txt_size) * line_h for _, t in bands]
    n = len(bands)
    region_bot = CONTENT_BOTTOM - (0.5 if note else 0.0)
    total = sum(heights)
    gap = max(0.10, min(0.5, (region_bot - top0 - total) / (n - 1)))
    y = top0
    last_bottom = top0
    for i, (label, texte) in enumerate(bands):
        h = heights[i]
        D.add_rect(s, MARGIN, y, CONTENT_W, h, fill="#ffffff", line=LINE, line_w=0.75, rounded=True, radius=0.08)
        D.add_rect(s, MARGIN, y, 0.06, h, fill=color, rounded=True, radius=0.5)
        D.add_text(s, MARGIN + 0.22, y + 0.12, CONTENT_W - 0.44, h - 0.24, [
            (label, dict(size=7, bold=True, color=color)),
            (texte, dict(size=txt_size, color=NAVY, space_before=4, line_spacing=1.25)),
        ])
        last_bottom = y + h
        y = last_bottom + gap

    if note:
        note_top = last_bottom + 0.16
        note_h = min(0.45, CONTENT_BOTTOM - note_top)
        D.add_text(s, MARGIN, note_top, CONTENT_W, note_h, [
            (note, dict(size=8, color=MUTED, italic=True, line_spacing=1.2)),
        ])
    return s


# --- Nouveau (brainstorm) : la formule de priorisation (chapitre Proposition)
# cite "prudence IA" sans jamais l'expliquer — cette slide la décompose. Dans le
# chapitre IA, juste après le gate IA (qui l'a rejoint) et avant les 3 candidats
# d'agent : on pose d'abord le frein, ensuite seulement les cas d'usage.
def slide_prudence_ia(prs):
    s = content_slide(prs, "IA", "La prudence IA est un frein chiffré, pas un veto", color=D.PALETTE[4])
    D.add_text(s, MARGIN, CONTENT_TOP, CONTENT_W, 0.4, [
        ("Prudence IA = confidentialité + besoin de supervision + criticité de la décision",
         dict(size=D.TYPE["small"], bold=True, color=NAVY, line_spacing=1.2)),
    ])

    facteurs = [
        ("CONFIDENTIALITÉ", D.PALETTE[0],
         "Reprend directement la classification du gate IA (D0-D4, slide précédente) — "
         "plus la donnée est sensible, plus le score monte."),
        ("BESOIN DE SUPERVISION", D.PALETTE[3],
         "Le palier d'adoption visé (assisté / supervisé / délégué) — un agent encore "
         "au stade assisté pèse plus lourd qu'un agent déjà éprouvé."),
        ("CRITICITÉ DE LA DÉCISION", D.PALETTE[2],
         "L'impact d'une erreur si l'agent se trompe seul — une recommandation "
         "réversible pèse moins qu'une décision irréversible."),
    ]
    top0 = CONTENT_TOP + 0.55
    n = 3
    card_h = 1.55
    for i, (label, color, texte) in enumerate(facteurs):
        x, w = col_x(i, n)
        D.add_card(s, x, top0, w, card_h, color)
        pad = 0.16
        D.add_text(s, x + pad, top0 + 0.14, w - 2 * pad, 0.4, [
            (label, dict(size=8, bold=True, color=color, line_spacing=1.1)),
        ])
        D.add_text(s, x + pad, top0 + 0.55, w - 2 * pad, card_h - 0.65, [
            (texte, dict(size=8, color=NAVY, line_spacing=1.25)),
        ])

    note_top = top0 + card_h + 0.18
    note_h = min(1.15, CONTENT_BOTTOM - note_top)
    D.add_rect(s, MARGIN, note_top, CONTENT_W, note_h, fill=NAVY, rounded=True, radius=0.08)
    D.add_text(s, MARGIN + 0.22, note_top, CONTENT_W - 0.44, note_h, [
        ("Un frein, pas un veto automatique", dict(size=D.TYPE["tiny"], bold=True, color="#ffffff")),
        ("Le score est SOUSTRAIT de impact × faisabilité — un candidat facile et à fort "
         "impact peut quand même être écarté si sa prudence IA est trop "
         "haute. Le score ne remplace pas l'arbitrage humain : il le rend "
         "explicite. Avancer malgré un score élevé reste possible, mais se documente comme une "
         "décision à part entière (même discipline que la dérogation du gate DevOps).",
         dict(size=8, color="#c7cbe0", space_before=3, line_spacing=1.25)),
    ], anchor=MSO_ANCHOR.MIDDLE)
    return s


# --- Nouveau (brainstorm) : comment le tronc commun se branche concrètement
# sur le SI du client, et ce que ça change selon le niveau d'ambition déjà
# cadré (slide précédente) — synthèse d'éléments déjà posés (§Ambition de
# l'outil, §Solution technique envisagée), pas une nouvelle doctrine.
def slide_architecture_si(prs):
    # v2.5 (chantier ④) : déplacée de la Proposition vers l'Outillage IAP, avec
    # slide_ambition (qui reste la slide précédente).
    s = content_slide(prs, "Outillage IAP",
                       "Le lien avec le SI du client change avec le niveau d'ambition, pas la méthode",
                       color=D.PALETTE[5])
    headers = ["NIVEAU", "SOURCES", "MODE DE CONNEXION", "LIVRABLES"]
    col_widths = [1.1, 2.55, 2.75, 1.95]
    xs = []
    cx = MARGIN
    for cw in col_widths:
        xs.append(cx)
        cx += cw + 0.1

    header_y = CONTENT_TOP + 0.05
    for x, w, label in zip(xs, col_widths, headers):
        D.add_text(s, x, header_y, w, 0.22, [
            (label, dict(size=7, bold=True, color=MUTED)),
        ])

    rows = [
        ("A", D.PALETTE[0], "Aide au coach",
         "Exports ponctuels (ServiceNow/Jira), interviews",
         "Aucune — tout est apporté par le consultant",
         "Markdown + deck, à la demande"),
        ("B", D.PALETTE[3], "Assistant interactif",
         "Exports + App companion (capture terrain)",
         "Site web centralisé, orchestration assistée",
         "+ tableau de bord multi-engagements"),
        ("C", D.PALETTE[2], "Companion connecté (non engagé)",
         "ServiceNow/Jira/Confluence/Datadog/CMDB/FinOps",
         "Connecteurs API directs, en continu",
         "Livrables mis à jour en continu"),
    ]
    row_top = header_y + 0.3
    row_h = 1.05
    row_gap = 0.1
    for i, (code, color, niveau, sources, connexion, livrables) in enumerate(rows):
        y = row_top + i * (row_h + row_gap)
        D.add_rect(s, MARGIN, y, CONTENT_W, row_h, fill="#ffffff", line=LINE, line_w=0.75, rounded=True, radius=0.08)
        D.add_rect(s, MARGIN, y, 0.06, row_h, fill=color, rounded=True, radius=0.5)
        D.add_text(s, xs[0] + 0.12, y, col_widths[0] - 0.12, row_h, [
            (code, dict(size=D.TYPE["h3"], bold=True, color=color)),
            (niveau, dict(size=7, color=MUTED, space_before=2, line_spacing=1.1)),
        ], anchor=MSO_ANCHOR.MIDDLE)
        D.add_text(s, xs[1], y, col_widths[1] - 0.1, row_h, [
            (sources, dict(size=8, color=NAVY, line_spacing=1.2)),
        ], anchor=MSO_ANCHOR.MIDDLE)
        D.add_text(s, xs[2], y, col_widths[2] - 0.1, row_h, [
            (connexion, dict(size=8, color=NAVY, line_spacing=1.2)),
        ], anchor=MSO_ANCHOR.MIDDLE)
        D.add_text(s, xs[3], y, col_widths[3] - 0.1, row_h, [
            (livrables, dict(size=8, color=NAVY, line_spacing=1.2)),
        ], anchor=MSO_ANCHOR.MIDDLE)

    note_top = row_top + 3 * row_h + 2 * row_gap + 0.15
    note_h = min(0.5, CONTENT_BOTTOM - note_top)
    D.add_text(s, MARGIN, note_top, CONTENT_W, note_h, [
        ("Le niveau C suppose un accès direct aux données de production du client — un cabinet "
         "peut durablement rester au niveau A ou B par choix de gouvernance (slide précédente).",
         dict(size=8, color=MUTED, italic=True, line_spacing=1.25)),
    ])
    return s


# --- Nouveau (réouverture de périmètre, arbitrage 2026-07-21) : l'architecture
# des 11 agents-workflows (§Workflows, ligne 587+), délibérément retirée du deck
# au commit 4f0c9b7, est rouverte sur ce seul point. COMPLÉMENTAIRE de
# slide_schema_fonctionnement (le FLUX de données Collecte→Diagnostic→Conception
# →Restitution, avec flèches) : ici c'est l'INVENTAIRE des composants — les 11
# agents nommés, regroupés par étape en cartes (pas de flèches), le gate
# confidentialité posé comme un socle transversal et bloquant. Couleurs des
# familles reprises de slide_schema_fonctionnement (Diagnostic=violet,
# Conception=or, etc.) — cohérence inter-slides voulue, pas un hasard.
# v2.5 (chantier ④) : déplacée de l'IA vers la Démarche, juste après
# slide_schema_fonctionnement (le flux) — l'inventaire des composants.
def slide_architecture_agents(prs):
    s = content_slide(prs, "Démarche",
                       "Onze workflows outillés, un seul bloquant : le gate confidentialité les traverse tous",
                       color=D.PALETTE[3])
    D.add_text(s, MARGIN, CONTENT_TOP, CONTENT_W, 0.5, [
        ("Un mandat unique par workflow, regroupés par étape. Le gate confidentialité est le seul "
         "à pouvoir arrêter la chaîne — transversal, il précède tout usage d'un modèle IA sur "
         "donnée client.", dict(size=8, color=MUTED, italic=True, line_spacing=1.2)),
    ])

    familles = [
        ("INTAKE", D.PALETTE[0], [
            ("iap-intake",
             "Qualifie le contexte client, le positionne sur les deux échelles de maturité, "
             "puis choisit le chemin de mission : diagnostic, pilote, adoption ou gate d'abord."),
        ]),
        ("DIAGNOSTIC", D.PALETTE[4], [
            ("iap-diagnostic-systemique", "Structure, flux, RUN, posture management"),
            ("iap-discovery-gaspillage", "Preuves, causes racines, options de traitement"),
        ]),
        ("CONCEPTION", D.PALETTE[3], [
            ("iap-waste-treatment", "Backlog priorisé et scoré des gaspillages"),
            ("iap-product-definition", "Personas, capacités, valeur, roadmap"),
            ("iap-operating-model", "Rôles, gouvernance, financement (décisions actées)"),
            ("iap-agentic-opportunities", "Le gaspillage d'abord, l'IA ensuite"),
        ]),
        ("ADOPTION & RESTITUTION", D.PALETTE[1], [
            ("iap-adoption-plan", "Onboarding, documentation, communautés"),
            ("iap-scenario-playbook", "Adapte la démarche au scénario client"),
            ("iap-deck-builder", "Deck modulaire, restitution exécutive"),
        ]),
    ]
    n = len(familles)
    # Cartes de MÊME hauteur (cadence sur la colonne la plus fournie, CONCEPTION
    # à 4 agents) — mais au lieu d'un slot fixe aligné en haut qui laissait un
    # grand vide sous l'agent unique d'INTAKE, chaque colonne RÉPARTIT ses agents
    # sur toute la zone : chaque bloc-agent est dimensionné à SON texte, puis les
    # blocs sont espacés (space-between) pour couvrir la hauteur — colonne à 4 =
    # remplie ; colonne à 2-3 = espacée régulièrement ; colonne à 1 = centrée. La
    # rangée se lit ainsi « équilibrée » (cf. brief ppt-designer, défaut INTAKE).
    top0 = CONTENT_TOP + 0.6
    gate_h = 0.6
    note_h = 0.36
    card_h = CONTENT_BOTTOM - top0 - 0.15 - gate_h - 0.12 - note_h
    _, wcol = col_x(0, n)
    pad = 0.14
    usable_col = wcol - 2 * pad
    region_top = top0 + 0.52
    region_h = card_h - 0.52 - 0.08
    for i, (nom, color, agents) in enumerate(familles):
        x, w = col_x(i, n)
        D.add_card(s, x, top0, w, card_h, color)
        D.add_text(s, x + pad, top0 + 0.12, w - 2 * pad, 0.36, [
            (nom, dict(size=8, bold=True, color=color, line_spacing=1.0)),
            (f"{len(agents)} workflow" + ("s" if len(agents) > 1 else ""),
             dict(size=6.5, color=MUTED, space_before=1)),
        ])
        blocs = [0.16 + _lignes(role, usable_col, 6.5) * (6.5 * 1.15 / 72.0)
                 for _, role in agents]
        na = len(agents)
        total = sum(blocs)
        if na == 1:
            gap_a = 0.0
            start = region_top  # top-align le bloc unique (INTAKE) sous l'en-tête : le centrer le faisait flotter (défaut « panneau flottant », cf. revue 2026-07-21)
        else:
            gap_a = min(0.5, (region_h - total) / (na - 1))
            span = total + gap_a * (na - 1)
            start = region_top + max(0.0, (region_h - span) / 2)
        ay = start
        for j, (agent, role) in enumerate(agents):
            D.add_text(s, x + pad, ay, w - 2 * pad, blocs[j], [
                (agent, dict(size=7, bold=True, color=NAVY, line_spacing=1.0)),
                (role, dict(size=6.5, color=MUTED, space_before=2, line_spacing=1.1)),
            ], anchor=MSO_ANCHOR.TOP)
            ay += blocs[j] + gap_a

    gate_top = top0 + card_h + 0.15
    D.add_rect(s, MARGIN, gate_top, CONTENT_W, gate_h, fill=NAVY, rounded=True, radius=0.1)
    chip_w = 1.15
    chip(s, MARGIN + 0.16, gate_top + gate_h / 2 - 0.14, chip_w, 0.28, "BLOQUANT", SEVERITE[4], size=7)
    D.add_text(s, MARGIN + 0.16 + chip_w + 0.2, gate_top + 0.1, CONTENT_W - chip_w - 0.55, gate_h - 0.2, [
        ("iap-ai-data-confidentiality-gate", dict(size=8, bold=True, color="#ffffff")),
        ("Classe les données (D0-D4), décide le mode d'exécution IA et pose les garde-fous — "
         "transversal, avant tout usage d'un modèle IA sur donnée client.",
         dict(size=7.5, color="#c7cbe0", space_before=2, line_spacing=1.15)),
    ], anchor=MSO_ANCHOR.MIDDLE)

    note_top = gate_top + gate_h + 0.12
    note_h_real = min(note_h, CONTENT_BOTTOM - note_top)
    D.add_text(s, MARGIN, note_top, CONTENT_W, note_h_real, [
        ("Onze mandats distincts, un seul peut arrêter la chaîne — tous les autres proposent et "
         "produisent, la décision finale reste humaine.",
         dict(size=8, color=MUTED, italic=True, line_spacing=1.2)),
    ])
    return s


def build():
    prs = new_prs()
    slide_cover(prs)

    # === Chapitre 01 — EXEC SUMMARY : le pitch de l'offre (v2.8, refondu v2.9) ===
    # v2.9 (arbitrage utilisateur) : le grand schéma du parcours de mission
    # (slide_offre_iap) part au chapitre 07 · Démarche et la synthèse en une page
    # (slide_offre_synthese) est supprimée — slide_executive_summary reste LE
    # sommaire du deck. À la place, deux slides qui parlent au prospect : les
    # trois faces de l'offre, puis la même démarche avec ou sans agentic.
    # v2.13 (2026-09-03, arbitrage utilisateur) : slide_executive_summary
    # DÉMÉNAGE d'avant l'intercalaire à juste après — modifié manuellement par
    # l'utilisateur sur l'export, reporté ici pour que toute régénération le
    # conserve (sinon un rebuild écraserait l'édition manuelle sans bruit).
    slide_chapitre(prs, "01", "Exec summary",
                   "Les trois faces de l'offre — leurs douleurs, notre outillage de "
                   "consultant, l'agentic déployé chez eux en option — et la démarche qui les "
                   "absorbe.",
                   NAVY, "wheatfield", seed=0)
    slide_executive_summary(prs)
    slide_pitch_iap(prs)
    slide_demarche_avec_sans_agentic(prs)

    slide_vision(prs)

    # === Chapitre 02 — CONTEXTE : le problème ===
    slide_chapitre(prs, "02", "Contexte",
                   "La double mission, et pourquoi cette transformation a du sens pour un client infra maintenant.",
                   D.PALETTE[0], "mountains", seed=0)
    slide_mission(prs)
    slide_pourquoi_contexte(prs)
    # 2026-09-01 : « qui achète, contre quoi » — la section §Positionnement &
    # achat du cadrage (l.36) fait foi pour le deck depuis la v2.3 sans y avoir
    # jamais été redescendue. Vient APRÈS les 3 déclencheurs, qu'elle prolonge
    # (le déclencheur ① et l'écart 80/30 y sont repris comme réponse d'achat).
    slide_qui_achete(prs)

    # === Chapitre 03 — PERSONAS : qui l'on transforme ===
    slide_chapitre(prs, "03", "Personas",
                   "Quatre parties prenantes interrogées séparément — leurs voix, leurs postures, les tensions.",
                   D.PALETTE[5], "forest", seed=0)
    slide_personas(prs)
    slide_personas_divergences(prs)

    # === Chapitre 04 — BESOINS & DOULEURS : ce qui fait mal ===
    slide_chapitre(prs, "04", "Besoins & douleurs",
                   "Les douleurs approfondies et mesurables, et les 8 familles de gaspillage qui les rangent.",
                   D.PALETTE[2], "ocean", seed=0)
    slide_douleurs(prs)
    slide_familles(prs)

    # === Chapitre 05 — PROPOSITION : le QUOI ===
    # Fil rouge : la THÈSE (why_iap) ouvre, puis la MÉTHODE scorée (gaspillages),
    # puis la cible d'organisation. v2.5 (chantier ④) : schéma de fonctionnement
    # + livrables → Démarche ; ambition + lien SI → Outillage IAP. v2.6 : le
    # sous-chapitre « Exemples » (séparateur slide_sous_chapitre + 3 slides
    # illustratives) est supprimé à la demande — git garde l'historique (v2.5).
    slide_chapitre(prs, "05", "Proposition",
                   "Traiter l'infra comme un produit : la thèse, la méthode scorée et l'organisation cible.",
                   D.PALETTE[1], "dunes", seed=0)
    slide_why_iap(prs)
    slide_gaspillages(prs)
    slide_team_topologies(prs)

    # === Chapitre 06 — IA : tirée APRÈS la proposition (l'IA amplifie, n'est jamais la réponse) ===
    slide_chapitre(prs, "06", "IA",
                   "L'IA au service de la réponse : le gate confidentialité, la prudence, les agents candidats, l'export.",
                   D.PALETTE[4], "nightsky", seed=0)
    slide_gate_ia(prs)
    slide_prudence_ia(prs)
    slide_agent_ia(
        prs, "Un agent de triage peut absorber le gaspillage RUN le plus répétitif",
        "Agent de triage de tickets", "RUN",
        "Les mêmes types de tickets reviennent depuis des années et mobilisent des experts "
        "seniors sur un travail répétitif à faible valeur — le gaspillage RUN le plus classique.",
        "Lit chaque ticket entrant, le classe selon un runbook déjà documenté et le route vers "
        "la bonne équipe. Le processus doit être explicite avant l'agent (préalable non "
        "négociable) — jamais l'inverse.",
        "Capacité RUN récupérée : dans le cas nominal du cadrage, jusqu'à 15 tickets/mois "
        "traités sans intervention humaine, temps de triage divisé par deux.",
        D.PALETTE[2])
    slide_agent_ia(
        prs, "Un agent de veille FinOps rend le décommissionnement continu, pas ponctuel",
        "Agent de veille FinOps", "Financier",
        "Les ressources cloud surdimensionnées ou orphelines ne sont détectées qu'à l'occasion "
        "d'audits ponctuels — le gaspillage s'accumule entre deux revues manuelles.",
        "Scanne en continu la CMDB et la facturation cloud, détecte les ressources inactives ou "
        "surdimensionnées, et propose une liste de décommissionnement à valider — ne décommissionne "
        "jamais seul.",
        "Coût récupéré et directement mesurable (ressources décommissionnées par mois) — un KPI "
        "de mission déjà cadré, pas à inventer.",
        D.PALETTE[3])
    slide_agent_ia(
        prs, "Un agent documentaire réduit la charge cognitive sans remplacer l'expert",
        "Agent documentaire (RAG)", "Cognitif",
        "Trop d'outils, procédures dispersées : la charge cognitive pour retrouver l'information "
        "ralentit les équipes et sollicite en permanence les mêmes experts.",
        "Indexe la documentation existante (runbooks, wikis, tickets résolus) et répond aux "
        "questions fréquentes avec la source citée — jamais une réponse sans preuve.",
        "Charge cognitive réduite, onboarding plus rapide, moins d'interruptions des experts "
        "seniors pour des questions déjà documentées.",
        D.PALETTE[4],
        note=("Ces 3 candidats restent soumis au scoring (chapitre Proposition) et au gate IA (ouverture de ce chapitre) "
              "avant toute décision — des exemples illustratifs, pas une liste actée."))
    slide_export_markdown(prs)

    # === Chapitre 07 — DÉMARCHE : le COMMENT (après l'IA, pour enchaîner sur
    # l'outillage puis la preuve — cf. docstring v2.5) ===
    slide_chapitre(prs, "07", "Démarche",
                   "La trajectoire et ses livrables par phase, le fil humain, le schéma de "
                   "fonctionnement et l'inventaire des agents.",
                   D.PALETTE[3], "canyon", seed=0)
    # v2.5 (chantier ①) : trajectoire fusionnée avec la vue bout-en-bout.
    slide_trajectoire(prs)
    # v2.4 : le fil humain décline la trame ①②③⟲ de slide_trajectoire côté
    # personnes — placé juste après elle.
    slide_fil_humain(prs)
    # v2.6 (point ②) : les activités humaines de la démarche, avec/sans l'outil
    # — juste après le fil humain, qu'elle décline en registres d'activités.
    slide_activites_humaines(prs)
    # 2026-09-01 : conditions de réussite et non-engagement — juste après le fil
    # humain et ses activités, dont elle prolonge le « testé dès l'intake » : ce
    # que la mission exige du client, et ce que son absence déclenche.
    slide_conditions_reussite(prs)
    # v2.9 (arbitrage utilisateur) : le parcours de mission détaillé arrive du
    # chapitre 01 · Exec summary. Placé ICI, en tête du bloc des schémas (parcours
    # de mission → schéma de fonctionnement → inventaire des agents → livrables),
    # plutôt qu'accolé à slide_trajectoire : le fil humain ①②③⟲ (trajectoire →
    # fil humain → activités → conditions) reste d'un seul tenant.
    slide_offre_iap(prs)
    # v2.5 (chantier ④) : déplacées de la Proposition (schéma, livrables) et de
    # l'IA (inventaire des agents) vers la Démarche.
    slide_schema_fonctionnement(prs)
    slide_architecture_agents(prs)
    slide_livrables_ppt(prs)

    # === Chapitre 08 — OUTILLAGE IAP : l'AVEC QUOI (nouveau, v2.5) ===
    slide_chapitre(prs, "08", "Outillage IAP",
                   "Ce que le module met dans les mains du consultant : l'architecture en "
                   "contexte client, trois niveaux d'ambition, le lien avec le SI.",
                   D.PALETTE[5], "tropical", seed=0)
    # v2.6 (point ③) : le chapitre OUVRE sur le schéma d'architecture en
    # contexte client — ambition et lien SI le déclinent ensuite.
    slide_iap_contexte_client(prs)
    slide_ambition(prs)
    slide_architecture_si(prs)

    # === Chapitre 09 — KPI : la preuve (clôture du deck) ===
    # Les 3 familles ouvrent, puis leur pourquoi/quoi et leur mise en place ; la
    # grille de maturité (slide_maturite, la 3e famille détaillée) vient ensuite
    # (déplacée après kpis_mise_en_place, point ①), et le cas nominal chiffré
    # ferme le deck.
    slide_chapitre(prs, "09", "KPI",
                   "Trois familles de KPIs à ne jamais confondre, leur mise en place, la grille de maturité, et le cas chiffré.",
                   D.PALETTE[0], "meadow", seed=1)
    slide_kpis(prs)
    slide_kpis_pourquoi_quoi(prs)
    slide_kpis_mise_en_place(prs)
    slide_maturite(prs)
    slide_kpis_exemple(prs)

    problemes = D.verifier_geometrie(prs) + _ANOMALIES_BUILD
    if problemes:
        print(f"GEOMETRIE: {len(problemes)} probleme(s)")
        for p in problemes:
            print(" -", p)
    else:
        print("GEOMETRIE: OK — aucune forme hors cadre")

    # Un deck dont le controle signale un defaut n'ecrase PLUS le livrable.
    # `prs.save(out)` etait inconditionnel : seul le code de sortie signalait
    # l'echec, donc tout appelant qui l'ignore (IDE, double-clic, script sans
    # `set -e`) livrait un deck casse en silence (mesure du 2026-09-01). En cas
    # de defaut on ecrit a cote, sous un nom qui ne trompe personne, et l'export
    # precedent — valide — reste en place.
    base = os.path.join(os.path.dirname(__file__), "bmad-iap-cadrage-synthese")
    if problemes:
        out = base + ".INVALIDE.pptx"
        prs.save(out)
        print(f"Ecrit (NON LIVRABLE, {len(problemes)} probleme(s)):", out)
        print("Export livrable inchange:", base + ".pptx")
    else:
        out = base + ".pptx"
        prs.save(out)
        print("Ecrit:", out)
    return problemes


if __name__ == "__main__":
    problemes = build()
    sys.exit(1 if problemes else 0)
